from __future__ import annotations

# D-012 — streamlit run src/ui/app.py 형태로 실행되면 sibling import
# (rag_client / chatbot_flow / feedback_chat / enrich / doc_packaging)가
# 동작하도록 sys.path를 보강. 동시에 프로젝트 루트도 path에 넣어
# `from src.db.retrieve import ...` 같은 절대 import도 가능.
import sys as _sys
from pathlib import Path as _Path
_HERE = _Path(__file__).resolve().parent
_ROOT = _HERE.parent.parent
for _p in (_HERE, _ROOT):
    if str(_p) not in _sys.path:
        _sys.path.insert(0, str(_p))

import re
import streamlit as st
from rag_client import retrieve_docs
import pandas as pd

def _norm_merge_text(x: Any) -> str:
    return re.sub(r"\s+", " ", str(x or "").strip()).upper()

def _norm_merge_pno(x: Any) -> str:
    return re.sub(r"[^A-Z0-9]", "", str(x or "").strip().upper())

def _action_bucket(v: Any) -> str:
    a = _norm_merge_text(v or "")
    if a in ("ADD", "MODIFY", "CHECK", "추가", "변경", "⚠️확인필요"):
        return "CHANGE"
    if a in ("REMOVE", "DELETE", "삭제"):
        return "DELETE"
    if a in ("KEEP", "유지"):
        return "KEEP"
    return a

def _real_part_pno(pt: dict) -> str:
    """Return normalized real part number, excluding placeholders like TBD/채번."""
    for k in ("part_no", "display_pno", "db_new_pno", "db_base_pno", "품번", "Part No", "P/NO", "P/NO."):
        v = str(pt.get(k) or "").strip().upper()
        if not v:
            continue
        if any(x in v for x in ("채번", "TBD", "NEED", "미정")):
            continue
        n = re.sub(r"[^A-Z0-9]", "", v)
        if n:
            return n
    return ""

def _part_identity_key(pt: dict, *, l1_desc: str = "", row_uid: Any = "") -> tuple:
    """
    Stable dedup key used across pipeline:
    - scope by L1 + action bucket first
    - prefer real part_no
    - fallback to name/type/level when part_no is absent
    """
    l1_key = _norm_merge_text(l1_desc or pt.get("l1_desc") or "")
    act_bucket = _action_bucket(pt.get("action") or pt.get("변경유형") or "")
    pno = _real_part_pno(pt)
    if pno:
        return ("PNO", l1_key, act_bucket, pno)

    name_key = _norm_merge_text(pt.get("part_name") or pt.get("desc") or pt.get("부품명") or "")
    type_key = _norm_merge_text(pt.get("base_type") or pt.get("type") or pt.get("유형") or "")
    lvl_key = _norm_merge_text(pt.get("lvl") or pt.get("db_lvl") or pt.get("레벨") or "")
    if name_key:
        return ("NAME", l1_key, act_bucket, name_key, type_key, lvl_key)

    # Last-resort key keeps no-name/no-pno rows from collapsing accidentally.
    src_key = _norm_merge_text(pt.get("source_doc") or pt.get("출처") or "")
    rsn_key = _norm_merge_text(pt.get("rsn") or pt.get("chg") or pt.get("변경사유") or "")
    return ("ROW", l1_key, act_bucket, src_key, rsn_key, str(row_uid or ""))

def _normalize_rsn(rsn: str) -> str:
    """
    RSN(변경사유)에서 핵심 기술 키워드만 추출해 정규화
    예: "카메라 모듈 적용" → "카메라_모듈"
         "카메라 아셈블리 모듈 개발" → "카메라_모듈"
    같은 의미인데 다르게 표현되면 같은 key로 인식되도록 함
    """
    if not rsn:
        return ""
    
    # 주요 기술/부품 키워드
    tech_keywords = {
        "카메라", "센서", "모듈", "어셈블리", "아셈블리",
        "구조", "성능", "방열", "전력", "신호", "통신",
        "LED", "PCB", "메모리", "프로세서", "인터페이스",
        "전압", "전류", "저항", "커패시터", "인덕터",
        "브래킷", "렌즈", "이미지", "해상도", "곡률"
    }
    
    rsn_norm = str(rsn).upper().strip()
    normalized = []
    
    for kw in sorted(tech_keywords):  # 일관성 유지
        if kw in rsn_norm:
            normalized.append(kw)
    
    if normalized:
        return "_".join(normalized)
    
    # 키워드 없으면 원본 RSN 반환 (fallback)
    return _norm_merge_text(rsn)

def _part_merge_key(pt: dict) -> tuple:
    return _part_identity_key(pt, l1_desc=pt.get("l1_desc") or "")

def _merge_part_lists(a: list[dict] | None, b: list[dict] | None) -> list[dict]:
    tier_rank = {"CORE": 3, "CORE_GROUP": 2, "CASCADE": 1}
    lvl_source_rank = {"DB": 3, "BASE": 2, "ESTIMATED": 1}
    merged = {}
    for src in (a or []):
        merged[_part_merge_key(src)] = dict(src)
    for src in (b or []):
        k = _part_merge_key(src)
        if k not in merged:
            merged[k] = dict(src)
        else:
            dst = merged[k]
            # Keep stronger classification when same part appears with different tier.
            dst_tier = _norm_merge_text(dst.get("tier") or "")
            src_tier = _norm_merge_text(src.get("tier") or "")
            if tier_rank.get(src_tier, 0) > tier_rank.get(dst_tier, 0):
                dst["tier"] = src.get("tier")

            # Keep more reliable level source: DB > BASE > ESTIMATED.
            dst_src = _norm_merge_text(dst.get("lvl_source") or "")
            src_src = _norm_merge_text(src.get("lvl_source") or "")
            if lvl_source_rank.get(src_src, 0) > lvl_source_rank.get(dst_src, 0):
                if str(src.get("lvl") or "").strip():
                    dst["lvl"] = src.get("lvl")
                dst["lvl_source"] = src.get("lvl_source")

            # 빈 값만 보강해서 입력 순서에 따른 결과 차이를 줄인다.
            for fld in ["rsn", "chg", "source_doc", "sourcing", "sourcing_reason", "base_type", "qty", "lvl", "lvl_source"]:
                if not str(dst.get(fld, "") or "").strip() and str(src.get(fld, "") or "").strip():
                    dst[fld] = src.get(fld)
            merged[k] = dst

    # Preserve first-seen order to keep BOM toggle sequence intact.
    return list(merged.values())

def merge_proposals_order_independent(proposals: list[dict]) -> list[dict]:
    """
    변경점별로 생성된 proposals를 내용 기준으로 병합해
    입력 순서(변경점1+2 vs 2+1)에 관계없이 동일한 결과를 만든다.
    """
    grouped: dict[tuple, dict] = {}

    for p in proposals or []:
        lvl1 = p.get("lvl1") or {}
        gkey = (
            _norm_merge_text(lvl1.get("desc") or ""),
            _norm_merge_pno(lvl1.get("part_no") or ""),
        )

        if gkey not in grouped:
            cp = dict(p)
            cp["changed_parts"] = _merge_part_lists(p.get("changed_parts") or [], [])
            cp["indirect_parts"] = _merge_part_lists(p.get("indirect_parts") or [], [])
            cp["existing_parts"] = _merge_part_lists(p.get("existing_parts") or [], [])
            cp["source_docs"] = sorted(set(p.get("source_docs") or []))
            cp["ref_models"] = sorted(set(p.get("ref_models") or []))
            cp["_change_summaries"] = set([str(p.get("change_summary") or "").strip()])
            grouped[gkey] = cp
            continue

        dst = grouped[gkey]
        dst["changed_parts"] = _merge_part_lists(dst.get("changed_parts") or [], p.get("changed_parts") or [])
        dst["indirect_parts"] = _merge_part_lists(dst.get("indirect_parts") or [], p.get("indirect_parts") or [])
        dst["existing_parts"] = _merge_part_lists(dst.get("existing_parts") or [], p.get("existing_parts") or [])
        dst["source_docs"] = sorted(set((dst.get("source_docs") or []) + (p.get("source_docs") or [])))
        dst["ref_models"] = sorted(set((dst.get("ref_models") or []) + (p.get("ref_models") or [])))
        dst["confidence"] = max(float(dst.get("confidence") or 0.0), float(p.get("confidence") or 0.0))

        dst_s = dst.get("_change_summaries") or set()
        dst_s.add(str(p.get("change_summary") or "").strip())
        dst["_change_summaries"] = dst_s
        grouped[gkey] = dst

    out = list(grouped.values())
    out.sort(key=lambda x: (
        _norm_merge_text((x.get("lvl1") or {}).get("desc") or ""),
        _norm_merge_pno((x.get("lvl1") or {}).get("part_no") or ""),
    ))

    for i, p in enumerate(out, 1):
        p["proposal_id"] = f"P-{i:03d}"
        sums = [s for s in sorted(p.pop("_change_summaries", set())) if s]
        if sums:
            p["change_summary"] = " | ".join(sums)

    return out
from dataclasses import dataclass, asdict
from typing import Any
from collections import defaultdict


# ✅ set_page_config는 반드시 "첫 st 호출"이고 1번만!
st.set_page_config(page_title="BOM Agent - Chat Input", layout="wide")

# 🆕 플로팅 챗봇 초기화
from streamlit_float import float_init
float_init()

import feedback_chat
import importlib
importlib.reload(feedback_chat)
from feedback_chat import init_feedback_chat, render_floating_chat
from feedback_chat import prepare_active_review, auto_detect_model_diff

# ✅ chatbot_flow 모듈 임포트
import chatbot_flow
importlib.reload(chatbot_flow)
from chatbot_flow import (
    init_chat, chat_flow, chat_add, render_chat,
    render_grade_selector, render_status_panel, render_debug_sidebar,
)

# =========================
# FAST MODE (입력부 생략용)
# =========================
# ── 🔄 초기화 버튼 (등급 재선택 포함 전체 리셋) ──
FAST_MODE = st.sidebar.toggle("⚡ FAST MODE (입력부 생략)", value=True)

# ── 🔄 초기화 버튼 (사이드바) ──
if st.session_state.get("grade_locked"):
    if st.sidebar.button("🔄 재시작", use_container_width=True):
        _reset_keys = [
            "dev_grade", "grade_locked", "chat_step", "messages",
            "target_model", "change_items", "ref_model",
            "proposals", "decisions", "evidence_pool", "cards_for_review",
            "base_bom_df", "base_df", "base_df_raw", "base_snapshot",
            "bom_model", "bom_uploaded", "_was_fast", "key_diff",
            "region", "detected_diffs",
        ]
        for _k in _reset_keys:
            if _k in st.session_state:
                del st.session_state[_k]
        st.rerun()

# D-012 fix: 원본 260508은 chatbot_flow의 사이드바 함수들을 import만 하고 호출 누락.
# 사이드바에 개발등급 선택 + 디버그 토글이 안 보였던 원인. 여기서 명시적으로 호출.
init_chat()                  # session_state 기본값 setup (멱등)
render_grade_selector()      # 사이드바: 개발등급 selectbox
render_debug_sidebar()       # 사이드바: 🧪 디버그 확인 토글 + search_debug 패널

# =========================
# Base Excel: header auto detect
# =========================
import pandas as pd

def _norm_hdr(x) -> str:
    """헤더/셀 텍스트 정규화: 대문자 + 특수문자 제거 (P/no. -> PNO, Desc. -> DESC)"""
    s = "" if x is None else str(x)
    s = s.strip().upper()
    s = re.sub(r"[^A-Z0-9가-힣]+", "", s)
    return s

# 우리가 찾고 싶은 헤더 키(정규화된 형태)
_HDR_TARGETS = {"MODULE", "LVL", "DESC", "PNO"}  # base 템플릿 키

def detect_header_row(uploaded_file, sheet_name=0, scan_rows=40) -> int:
    uploaded_file.seek(0)
    probe = pd.read_excel(uploaded_file, sheet_name=sheet_name, header=None, nrows=scan_rows, dtype=str)

    best_i, best_hit = 0, -1
    for i in range(len(probe)):
        row = probe.iloc[i].tolist()
        normed = {_norm_hdr(v) for v in row if v is not None and str(v).strip() not in ("", "nan", "NaN")}
        hit = len(normed.intersection(_HDR_TARGETS))
        if hit > best_hit:
            best_hit, best_i = hit, i
        if hit >= 3:  # 충분히 확실
            return i
    return best_i

def read_excel_auto_header(uploaded_file, sheet_name=0, scan_rows=40) -> tuple[pd.DataFrame, int]:
    h = detect_header_row(uploaded_file, sheet_name=sheet_name, scan_rows=scan_rows)
    uploaded_file.seek(0)
    df = pd.read_excel(uploaded_file, sheet_name=sheet_name, header=h, dtype=str)
    return df, h

# -------------------------------
# 1) 등급 입력 범용화 (B/b/B급/등급:B 등)
# -------------------------------
GRADE_PATTERNS = [
    re.compile(r"(?i)(?:개발\s*등급|등급)\s*[:/ ]*\s*([ABCD])\b"),
    re.compile(r"(?i)(?<![A-Z0-9])([ABCD])\s*급\b"),
    re.compile(r"(?i)\(([ABCD])\)"),
    re.compile(r"(?i)^\s*([ABCD])\s*$"),
]

def extract_grade(text: str) -> str:
    t = (text or "").strip()
    for pat in GRADE_PATTERNS:
        m = pat.search(t)
        if m:
            return m.group(1).upper()
    return ""


# -------------------------------
# 2) 모델코드 파서 (최종 요구사항)
# - 결과는 "반드시" 9자리 + W/L 시작(core9)
# - 소문자 OK
# - K/C 채널 prefix가 앞에 오면 제거하고 2번째부터 1번째로 간주
# - 10자리 이상(악세사리 1글자 삽입 등): 앞8 + 마지막1 => core9
#   예) WS9D7652WM -> WS9D7652M
# -------------------------------
MODEL_TOKEN_RE = re.compile(r"(?i)\b[KC]?[A-Z0-9]{9,15}\b")  # 토큰을 넉넉히 잡고 정규화에서 걸러냄

def normalize_model_candidate(raw: str) -> dict:
    s = (raw or "").strip()
    if not s:
        return {"raw": raw, "channel_prefix": "", "core9": ""}

    cleaned = re.sub(r"[^A-Z0-9]", "", s.upper())  # 혹시 섞인 구분자 제거
    channel = ""

    # 채널 prefix(K/C) 제거 (단, 다음 글자가 W/L일 때만 "채널"로 인정)
    if len(cleaned) >= 2 and cleaned[0] in ("K", "C") and cleaned[1] in ("W", "L"):
        channel = cleaned[0]
        cleaned = cleaned[1:]  # 이제 W/L로 시작하는 본체

    # 본체는 반드시 W/L 시작이어야 유효
    if not cleaned or cleaned[0] not in ("W", "L"):
        return {"raw": raw, "channel_prefix": channel, "core9": ""}

    # 9자리면 그대로
    if len(cleaned) == 9:
        return {"raw": raw, "channel_prefix": channel, "core9": cleaned}

    # 길면(10자리 이상): 앞8 + 마지막1 (중간 삽입/추가 구분 무시)
    if len(cleaned) > 9:
        core9 = cleaned[:8] + cleaned[-1]
        if len(core9) == 9 and core9[0] in ("W", "L"):
            return {"raw": raw, "channel_prefix": channel, "core9": core9}

    return {"raw": raw, "channel_prefix": channel, "core9": ""}

def extract_model_code(text: str) -> dict:
    t = text or ""
    for m in MODEL_TOKEN_RE.finditer(t):
        info = normalize_model_candidate(m.group(0))
        if info["core9"]:
            return info
    return {"raw": "", "channel_prefix": "", "core9": ""}


def model_parts(model_core9: str, channel_prefix: str = "") -> dict:
    m = (model_core9 or "").upper().strip()
    ch = (channel_prefix or "").upper().strip()
    if len(m) != 9:
        return {"raw": m, "channel_prefix": ch}

    return {
        "raw": m,
        "channel_prefix": ch,  # K/C or ""
        "p1_product": m[0],    # W/L
        "p2_type": m[1],
        "p3_series_or_fuel": m[2],
        "p4_platform": m[3],
        "p56_capacity": m[4:6],
        "p7_design": m[6],
        "p8_grade": m[7],
        "p9_color": m[8],
    }

def match_prefix_by_dev_grade(model_core9: str, dev_grade: str) -> str:
    """
    등급별 Primary 검색 prefix
    - D: 앞7
    - C: 앞6
    - B 이상/미정: 앞4
    """
    m = (model_core9 or "").upper().strip()
    g = (dev_grade or "").upper().strip()
    if len(m) != 9:
        return m
    if g == "D":
        return m[:7]
    if g == "C":
        return m[:6]
    return m[:4]

def select_policy(target_model_text: str, dev_grade: str, change_items: list[str]) -> dict:
    model_info = extract_model_code(target_model_text)
    core9 = model_info.get("core9", "")
    ch = model_info.get("channel_prefix", "")

    parsed = model_parts(core9, ch)
    prefix = match_prefix_by_dev_grade(core9, dev_grade)

    return {
        "dev_grade": (dev_grade or "").upper().strip(),
        "target_model_raw": target_model_text,
        "target_model_core9": core9,
        "channel_prefix": ch,
        "target_parsed": parsed,
        "primary_prefix": prefix,
        # ✅ W/L 고정 제거: core9 1번째 자리로 게이트
        "primary_product_gate": parsed.get("p1_product"),
        "secondary_enabled": is_reference_trigger(change_items),
    }

def is_reference_trigger(change_items: list[str]) -> bool:
    """Secondary(참고용) 풀 ON 트리거"""
    text = " ".join(change_items or []).lower()
    triggers = ["최초", "first", "최근", "트렌드", "camera", "카메라", "led", "하네스", "화각", "옵셋"]
    return any(k in text for k in triggers)

# =========================
# BASE MASTER 업로드 + Snapshot
# =========================

# ---- 컬럼 동의어(엑셀 헤더가 달라도 자동 매핑)
COL_SYNONYMS = {
    "part_name": ["부품명", "품명", "DESC", "DESC.", "Description", "DESCRIPTION", "Desc", "Desc."],
    "part_no":   ["품번", "부품번호", "P/NO", "P/NO.", "P/no.", "Pno", "P/N", "PARTNO", "Part No"],
    "module":    ["Module", "MODULE", "모듈"],
    "lvl":       ["Lvl", "LVL", "Level", "LEVEL", "레벨"],
    "cmdt":      ["CMDT", "공정", "분류"],
    "grade":     ["Grade", "GRADE", "등급", "Part Grade", "PARTGRADE"],
    "supplier":  ["Supplier", "Supplier Code", "SUPPLIERCODE", "협력사", "업체"],
    "qty":       ["Qty", "QTY", "수량", "개수", "수량(EA)"],
}

import json
from pathlib import Path
import chromadb

RAG_JSON_PREFIX = "__RAG_JSON__:"  # 문서 문자열에 JSON이 들어있다는 표시

import hashlib

@st.cache_data(ttl=3600)
def _load_bom_cached(file_path: str) -> pd.DataFrame:
    try:
        with open(file_path, "rb") as _f:
            df = pd.read_excel(_f, header=0, dtype=str, engine="openpyxl")
        df = df.dropna(how="all").fillna("")
        df.columns = [str(c).strip() for c in df.columns]
        return df
    except Exception as e:
        st.warning(f"⚠️ BOM 로드 실패: {e}")
        return pd.DataFrame()

def _file_digest(uploaded_file) -> str:
    try:
        b = uploaded_file.getvalue()
        return hashlib.md5(b).hexdigest()
    except Exception:
        return ""

def wrap_rag_json(doc_obj: dict) -> str:
    """Chroma documents는 문자열만 받으니 JSON을 문자열로 래핑"""
    return RAG_JSON_PREFIX + json.dumps(doc_obj, ensure_ascii=False)

def unwrap_rag_json(doc_text: str) -> dict | None:
    """retrieve_docs 결과에서 JSON 문서를 다시 dict로 복원"""
    if not doc_text or not isinstance(doc_text, str):
        return None
    if not doc_text.startswith(RAG_JSON_PREFIX):
        return None
    try:
        return json.loads(doc_text[len(RAG_JSON_PREFIX):])
    except Exception:
        return None
    
# =========================================================
# Structured RAG 전용 Chroma 컬렉션
# - 기존 get_collection() 컬렉션과 완전히 분리
# =========================================================
STRUCTURED_DB_DIR = str(Path(__file__).resolve().parent / "chroma_structured")
# D-012: v3 → v4 → v5.
#   v3: default(onnx) embedder → Windows [Errno 22]
#   v4: bge-m3 OK, but doc text 형식이 generate_proposals_from_docs의
#       '[L1]' 마커 expectation과 mismatch → proposals 0건
#   v5: build_structured_docs_from_base가 [L1] / '- ' 형식으로 출력
STRUCTURED_COLLECTION_NAME = "bom_structured_v5"

# ── Chroma 클라이언트 캐싱 ──
_STRUCT_CLIENT = None
_STRUCT_COLLECTION = None

def get_structured_client():
    global _STRUCT_CLIENT
    if _STRUCT_CLIENT is None:
        _STRUCT_CLIENT = chromadb.PersistentClient(path=STRUCTURED_DB_DIR)
    return _STRUCT_CLIENT

def get_structured_collection():
    global _STRUCT_COLLECTION
    if _STRUCT_COLLECTION is None:
        client = get_structured_client()
        # D-012: embedding_function을 명시적으로 None — chromadb의 기본 onnxruntime
        # 모델 자동 다운로드/실행을 막아 Windows [Errno 22] 회피. 우리 bge-m3로
        # 직접 임베딩해 upsert(embeddings=...) 형태로 전달.
        from chromadb.utils.embedding_functions import EmbeddingFunction

        class _NullEmbedder(EmbeddingFunction):
            def __call__(self, input):
                # chromadb는 embedding_function이 호출되면 List[List[float]] 기대.
                # upsert에서 embeddings 인자 직접 전달하면 이 함수는 호출되지 않음.
                raise RuntimeError(
                    "structured collection은 embeddings 인자를 직접 전달해야 합니다. "
                    "upsert_structured_docs / query_structured_docs 참조."
                )

        _STRUCT_COLLECTION = client.get_or_create_collection(
            name=STRUCTURED_COLLECTION_NAME,
            metadata={"hnsw:space": "cosine"},
            embedding_function=_NullEmbedder(),
        )
    return _STRUCT_COLLECTION

def reset_structured_collection():
    global _STRUCT_COLLECTION
    _STRUCT_COLLECTION = None

def query_structured_docs(query_text: str, top_k: int = 20, where: dict | None = None) -> list[dict]:
    """
    구조화 컬렉션 조회 결과를 run_search가 쓰는 dict 포맷으로 변환.

    D-012: query_texts 대신 bge-m3로 미리 임베딩한 query_embeddings 전달.
    """
    col = get_structured_collection()

    from src.embed.embedder import embed_texts
    q_emb = embed_texts([query_text])[0]

    kwargs = {
        "query_embeddings": [q_emb],
        "n_results": top_k,
    }
    if where:
        kwargs["where"] = where

    res = col.query(**kwargs)

    ids = (res.get("ids") or [[]])[0]
    docs = (res.get("documents") or [[]])[0]
    metas = (res.get("metadatas") or [[]])[0]
    dists = (res.get("distances") or [[]])[0]

    out = []
    for i in range(len(ids)):
        out.append({
            "id": ids[i],
            "document": docs[i],          # 검색용 텍스트
            "metadata": metas[i] if i < len(metas) else {},  # rag_json 포함
            "distance": dists[i] if i < len(dists) else 0.0,
        })
    return out


_PATH_SPLIT_RE = re.compile(r"\s*(?:>|/|\\|::|\||,)\s*")

def _norm_str(x: Any) -> str:
    if x is None:
        return ""
    s = str(x).strip()
    s = re.sub(r"\s+", " ", s)
    return s

def _norm_key(x: Any) -> str:
    return _norm_str(x).upper()

def _to_float(x: Any) -> float:
    try:
        if x is None or (isinstance(x, float) and pd.isna(x)):
            return 0.0
        s = str(x).strip()
        s = re.sub(r"[^\d\.]+", "", s)  # "1EA" 같은 것도 대응
        return float(s) if s else 0.0
    except Exception:
        return 0.0

# =========================
# Base Snapshot (CLEAN BLOCK)
# =========================
def _detect_columns(df: pd.DataFrame) -> tuple[pd.DataFrame, dict]:
    raw_cols = list(df.columns)
    norm_to_orig = {_norm_key(c): c for c in raw_cols}

    colmap = {}
    rename = {}
    for std, syns in COL_SYNONYMS.items():
        syn_norm = {_norm_key(s) for s in syns}
        for n, orig in norm_to_orig.items():
            if n in syn_norm:
                colmap[std] = orig
                rename[orig] = std
                break

    df2 = df.rename(columns=rename).copy()
    return df2, colmap


# =========================
# Recheck-only Rerank / Regenerate helpers
# =========================
def _tok(s: str) -> set[str]:
    s = (s or "").lower()
    s = re.sub(r"[^a-z0-9가-힣\s]+", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    if not s:
        return set()
    return set(s.split(" "))

def simple_rerank_docs(docs: list[dict], query: str, top_k: int = 10) -> list[dict]:
    """
    2-stage rerank(가벼운 버전):
    - retrieve_docs(벡터)로 이미 뽑힌 docs를
    - query 토큰 겹침(Jaccard)으로 재정렬
    """
    qtok = _tok(query)
    if not qtok:
        return docs[:top_k]

    scored = []
    for d in docs or []:
        text = (d.get("document") or d.get("text") or "")
        dtok = _tok(text)
        inter = len(qtok & dtok)
        union = len(qtok | dtok) or 1
        j = inter / union
        scored.append((j, d))
    scored.sort(key=lambda x: x[0], reverse=True)
    return [d for _, d in scored[:top_k]]

def build_rerank_query(card: dict) -> str:
    """
    재검토 카드 1개를 더 정밀하게 찾기 위한 query 생성
    - 변경점 + 적용경로 + 부품명/품번 + base 비교 힌트 섞기
    """
    change = card.get("change_summary") or ""
    bt = card.get("bom_target") or {}
    path = bt.get("apply_bom_path") or bt.get("base_bom_path") or ""
    parts = card.get("parts") or {}
    main0 = (parts.get("main") or [{}])[0] or {}
    pno = main0.get("part_no") or ""
    pname = main0.get("part_name") or ""
    return f"SRC: {change}\nBOM_PATH: {path}\nPART_NO: {pno}\nPART_NAME: {pname}"

def group_cards_level1(cards: list[dict]) -> dict[str, list[dict]]:
    """
    레벨1(가장 상위) 기준 그룹핑:
    - bom_target.apply_bom_path 또는 base_bom_path의 첫 토큰을 그룹키로 사용
    """
    g = defaultdict(list)
    for c in cards or []:
        bt = c.get("bom_target") or {}
        path = bt.get("apply_bom_path") or bt.get("base_bom_path") or ""
        toks = [t.strip() for t in (path or "").split(">") if t.strip()]
        key = toks[0] if toks else "(NO_PATH)"
        g[key].append(c)
    return dict(g)

def regenerate_for_recheck_cards(
    recheck_cards: list[dict],
    policy: dict,
    top_k_retrieve: int = 40,
    top_k_rerank: int = 8,
) -> list[dict]:
    """
    재검토 카드만 대상으로:
    1) card별 query 생성
    2) retrieve_docs로 재검색(좀 넓게)
    3) simple_rerank로 재정렬(상위만)
    4) reranked docs 기반으로 새 proposal 카드 생성(대안)
    """
    new_cards = []
    seq = 1

    for old in recheck_cards or []:
        q = build_rerank_query(old)

        # (선택) 정책 토큰을 query에 섞어 precision 강화
        # - prefix / product gate 등이 있으면 함께 넣기
        q2 = q
        if policy:
            if policy.get("primary_prefix"):
                q2 += f"\nPREFIX={policy['primary_prefix']}"
            if policy.get("primary_product_gate"):
                q2 += f"\nPRODUCT={policy['primary_product_gate']}"

        docs = retrieve_docs(q2, top_k=top_k_retrieve)  # rag_client.retrieve_docs 사용
        docs2 = simple_rerank_docs(docs, q, top_k=top_k_rerank)

        # docs2 -> proposal 카드로 변환(기존 generate_proposals_from_docs 재사용 가능)
        # 여기서는 old 카드와 연결될 수 있게 meta를 심어줌
        proposals = generate_proposals_from_docs(
            primary_docs=docs2,
            base_snapshot=st.session_state.get("base_snapshot"),
            change_items=[old.get("change_summary", "")]
        )

        for p in proposals:
            p["proposal_id"] = f"RR-{seq:03d}"
            p["meta"] = p.get("meta") or {}
            p["meta"]["recheck_from"] = old.get("proposal_id")
            p["meta"]["rerank_query"] = q
            seq += 1
            new_cards.append(p)

    return new_cards


def _normalize_recheck_doc(d: dict) -> dict:
    return {
        "id": d.get("id", ""),
        "document": d.get("text") or d.get("document") or "",
        "metadata": d.get("meta") or d.get("metadata") or {},
        "distance": d.get("dist") or d.get("distance") or 0.5,
    }


def _recheck_token_overlap(a: str, b: str) -> int:
    at = _tok(a)
    bt = _tok(b)
    return len(at & bt)


def _find_recheck_candidate_parts(candidate_props: list[dict], request: dict, current_part: dict) -> list[dict]:
    lvl1_desc = _norm_merge_text(request.get("lvl1_desc") or "")
    cur_pno = _norm_merge_pno(current_part.get("part_no") or current_part.get("display_pno") or "")
    cur_name = current_part.get("part_name") or current_part.get("desc") or ""
    cur_source = _norm_merge_text(current_part.get("source_doc") or "")
    cur_lvl = str(current_part.get("lvl") or "").strip()
    out = []

    for prop in candidate_props or []:
        prop_lvl1 = _norm_merge_text((prop.get("lvl1") or {}).get("desc") or "")
        if lvl1_desc and prop_lvl1 and lvl1_desc != prop_lvl1:
            continue

        for list_key in ("changed_parts", "indirect_parts"):
            for pt in (prop.get(list_key) or []):
                cand_pno = _norm_merge_pno(pt.get("part_no") or pt.get("display_pno") or "")
                cand_name = pt.get("part_name") or pt.get("desc") or ""
                cand_source = _norm_merge_text(pt.get("source_doc") or "")

                if cand_pno and cur_pno and cand_pno == cur_pno:
                    continue
                if cand_source and cur_source and cand_source == cur_source:
                    continue

                overlap = _recheck_token_overlap(cur_name, cand_name)
                same_lvl = 1 if cur_lvl and str(pt.get("lvl") or "").strip() == cur_lvl else 0
                same_tier = 1 if _norm_merge_text(pt.get("tier") or "") == _norm_merge_text(current_part.get("tier") or "") else 0
                same_action = 1 if _norm_merge_text(pt.get("action") or "") == _norm_merge_text(current_part.get("action") or "") else 0
                score = overlap * 10 + same_lvl * 5 + same_tier * 3 + same_action * 2

                if score <= 0:
                    continue

                out.append({
                    "score": score,
                    "part": dict(pt),
                    "list_key": list_key,
                })

    out.sort(key=lambda x: x["score"], reverse=True)
    return out


def process_pending_recheck_request():
    ss = st.session_state
    req = ss.get("pending_recheck")
    if not req:
        return

    proposals = ss.get("proposals") or []
    p_idx = int(req.get("proposal_idx", -1))
    list_key = req.get("part_list") or ""
    part_idx = int(req.get("part_idx", -1))
    if p_idx < 0 or p_idx >= len(proposals):
        ss["pending_recheck"] = None
        return
    if list_key not in ("changed_parts", "indirect_parts"):
        ss["pending_recheck"] = None
        return

    target_parts = proposals[p_idx].get(list_key) or []
    if part_idx < 0 or part_idx >= len(target_parts):
        ss["pending_recheck"] = None
        return

    target_part = target_parts[part_idx]
    change_summary = req.get("change_summary") or proposals[p_idx].get("change_summary") or " ".join(ss.get("change_items") or [])
    lvl1_desc = req.get("lvl1_desc") or (proposals[p_idx].get("lvl1") or {}).get("desc") or ""
    part_name = req.get("part_name") or target_part.get("part_name") or target_part.get("desc") or ""
    part_no = req.get("part_no") or target_part.get("part_no") or target_part.get("display_pno") or ""

    query = (
        f"CHANGE: {change_summary}\n"
        f"ASSY: {lvl1_desc}\n"
        f"CURRENT_PART_NAME: {part_name}\n"
        f"CURRENT_PART_NO: {part_no}\n"
        f"USER_FEEDBACK: {req.get('user_text', '')}\n"
        "GOAL: alternative part for same change context"
    )

    with st.spinner(f"🔄 {lvl1_desc} 내 {part_name or part_no} 재검토 중..."):
        docs = []
        try:
            legacy_raw = retrieve_docs(query, top_k=40)
            docs.extend([_normalize_recheck_doc(d) for d in (legacy_raw or [])])
        except Exception:
            pass

        try:
            docs.extend(query_structured_docs(query, top_k=30))
        except Exception:
            pass

        docs = simple_rerank_docs(docs, query, top_k=20)
        candidate_props = generate_proposals_from_docs(
            docs,
            ss.get("base_snapshot") or {},
            [change_summary],
        )
        candidates = _find_recheck_candidate_parts(candidate_props, req, target_part)

    if not candidates:
        msg = (
            f"🔎 **{lvl1_desc} > {part_name or part_no}** 재검토를 수행했지만, "
            "현재 변경사유와 Assy 문맥을 만족하는 더 나은 대안을 찾지 못했습니다."
        )
        ss.setdefault("fb_messages", []).append({"role": "assistant", "content": msg})
        ss.setdefault("fb_history", []).append({
            "ts": __import__("datetime").datetime.now().strftime("%H:%M:%S"),
            "action": "RECHECK",
            "items": [f"{lvl1_desc} > {part_name or part_no} : 대안 없음"],
        })
        ss["pending_recheck"] = None
        return

    best = candidates[0]["part"]
    old_name = target_part.get("part_name") or target_part.get("desc") or ""
    old_pno = target_part.get("part_no") or target_part.get("display_pno") or ""
    old_source = target_part.get("source_doc") or ""
    old_action = target_part.get("action") or ""
    old_tier = target_part.get("tier") or ""
    old_lvl = target_part.get("lvl") or ""

    for key in ["part_name", "part_no", "display_pno", "lvl", "qty", "base_type", "rsn", "chg", "tier", "source_doc", "sourcing", "sourcing_reason", "in_base"]:
        if key in best and str(best.get(key, "") or "").strip():
            target_part[key] = best.get(key)

    if not str(target_part.get("action") or "").strip():
        target_part["action"] = old_action
    if not str(target_part.get("tier") or "").strip():
        target_part["tier"] = old_tier
    if not str(target_part.get("lvl") or "").strip():
        target_part["lvl"] = old_lvl

    target_part["recheck_prev_name"] = old_name
    target_part["recheck_prev_pno"] = old_pno
    target_part["recheck_prev_source"] = old_source
    target_part["recheck_note"] = f"재검토 교체: {old_name or old_pno} -> {target_part.get('part_name') or target_part.get('part_no') or ''}"

    ss["proposals"] = proposals
    new_name = target_part.get("part_name") or target_part.get("desc") or ""
    new_pno = target_part.get("part_no") or target_part.get("display_pno") or ""
    new_source = target_part.get("source_doc") or ""
    msg = (
        f"🔄 **재검토 완료**\n\n"
        f"- 위치: {lvl1_desc}\n"
        f"- 기존: {old_name or '(이름없음)'} / {old_pno or '(품번없음)'}\n"
        f"- 대안: {new_name or '(이름없음)'} / {new_pno or '(품번없음)'}\n"
        f"- 근거: {new_source or '재검색 결과'}"
    )
    ss.setdefault("fb_messages", []).append({"role": "assistant", "content": msg})
    ss.setdefault("fb_history", []).append({
        "ts": __import__("datetime").datetime.now().strftime("%H:%M:%S"),
        "action": "RECHECK",
        "items": [f"{lvl1_desc} > {old_name or old_pno} -> {new_name or new_pno}"],
    })
    ss["pending_recheck"] = None


@dataclass
class BaseMasterSnapshot:
    meta: dict
    rows: list[dict]
    part_index: dict
    path_index: dict
    bom_tree: dict
    quality: dict

def build_bom_tree(paths: list[list[str]]) -> dict:
    root: dict = {}
    for toks in paths:
        cur = root
        for t in toks:
            cur = cur.setdefault(t, {})
    return root

## =========================================================
## [블록 A] Base -> Structured RAG Doc Builder
## =========================================================

def _norm_col_name(x: Any) -> str:
    s = "" if x is None else str(x)
    s = s.strip().upper()
    s = re.sub(r"[^A-Z0-9가-힣]+", "", s)
    return s

def build_bom_path_from_module_lvl(df2: pd.DataFrame) -> list[str]:
    """
    Module + Lvl 기반으로 BOM path 생성
    """
    module_col = None
    lvl_col = None

    for c in df2.columns:
        cc = _norm_col_name(c)
        if module_col is None and cc == "MODULE":
            module_col = c
        if lvl_col is None and cc in ["LVL", "LEVEL", "레벨"]:
            lvl_col = c

    if module_col is None or lvl_col is None:
        return [""] * len(df2)

    paths = []
    stack = []

    for _, row in df2.iterrows():
        mod = str(row.get(module_col, "") or "").strip()
        if not mod:
            paths.append("")
            continue

        try:
            lvl = int(float(str(row.get(lvl_col, "1") or "1").strip()))
        except Exception:
            lvl = 1

        if lvl <= 1:
            stack = [mod]
        else:
            stack = stack[:lvl - 1] + [mod]

        paths.append(" > ".join(stack))

    return paths

def pick_col(df: pd.DataFrame, candidates: list[str]) -> str | None:
    """
    Desc. / P/no. 같이 점이 있어도 잡히게 정규화 비교
    """
    norm_map = {_norm_col_name(c): c for c in df.columns}
    for k in candidates:
        nk = _norm_col_name(k)
        if nk in norm_map:
            return norm_map[nk]
    return None

def infer_feature_from_text(text: str) -> str:
    t = (text or "").lower()
    if "camera" in t or "카메라" in t:
        return "CAMERA"
    if "led" in t or "조명" in t or "램프" in t:
        return "LED"
    if "하네스" in t or "harness" in t or "wire" in t or "배선" in t:
        return "HARNESS"
    return ""

def build_structured_docs_from_base(df_base: pd.DataFrame, model: str, dev_grade: str) -> list[dict]:
    """
    BOM DataFrame -> Chroma 검색용 문서 리스트
    BOM 컬럼: Lvl(.1, ..2), Part No, Description, Parent Part No(모), Type
    """
    # BOM 컬럼 찾기
    lvl_col  = pick_col(df_base, ["Lvl", "LVL", "Level", "LEVEL", "레벨"])
    desc_col = pick_col(df_base, ["Description", "DESC", "DESC.", "Desc", "부품명", "품명"])
    pno_col  = pick_col(df_base, ["Part No", "P/NO", "P/NO.", "PNO", "품번", "부품번호"])
    parent_col = pick_col(df_base, ["Parent Part No(모)", "Parent Part No", "모품번"])
    type_col = pick_col(df_base, ["Type", "TYPE", "유형"])

    if lvl_col is None or desc_col is None:
        return []

    docs = []
    # BOM path 스택: Lvl 깊이에 따라 경로 빌드
    path_stack = []  # [(lvl_depth, part_no_or_desc)]

    for i, row in df_base.iterrows():
        lvl_raw = str(row.get(lvl_col, "") or "").strip()
        desc = str(row.get(desc_col, "") or "").strip()
        pno  = str(row.get(pno_col, "") or "").strip() if pno_col else ""
        ptype = str(row.get(type_col, "") or "").strip() if type_col else ""

        if not desc or lvl_raw == "0":
            continue

        # Lvl 깊이 계산: ".1" = 1, "..2" = 2, "...3" = 3
        dot_count = len(lvl_raw) - len(lvl_raw.lstrip("."))
        lvl_depth = max(dot_count, 1)

        # 경로 스택 업데이트
        label = desc[:40]
        path_stack = path_stack[:lvl_depth - 1] + [(lvl_depth, label)]
        bom_path = " > ".join([p[1] for p in path_stack])
        level1 = path_stack[0][1] if path_stack else label

        feature = infer_feature_from_text(desc)

        rag_obj = {
            "doc_id": f"{model}_{dev_grade}_{i}",
            "source": {"model": model, "dev_grade": dev_grade},
            "change": {
                "summary_raw": f"{level1} / {desc}",
                "target_object": level1,
                "action": "BASE",
                "feature": feature,
            },
            "bom": {
                "level1": level1,
                "apply_bom_path": bom_path,
                "base_bom_path": bom_path,
            },
            "parts": {
                "main": [{"part_name": desc, "part_no": pno, "qty": 1, "desc": desc}],
                "sub": [],
            },
            "reason": [],
            "review_points": [],
        }

        # D-012 fix: generate_proposals_from_docs는 doc text에 '[L1]' 마커가
        # 있어야 부품 파싱을 시작 (app.py line 3769). 부품 라인은 '-'로 시작.
        # 원본 doc_packaging.make_index_docs_l1_chunks가 만들던 형식과 정합.
        embedding_text = (
            f"[SRC] runtime_bom | {model}\n"
            f"[MODEL] {model}\n"
            f"[GRADE] {dev_grade}\n"
            f"[DOMAIN] M\n"
            f"[L1] {(pno or 'NOPNO')} | Desc={desc[:80]}\n"
            f"- {bom_path} | Base= New={pno} | {desc[:80]}"
            + (f" | FEATURE={feature}" if feature else "")
            + (f" | TYPE={ptype}" if ptype else "")
            + "\n"
            # 검색용 보조 키워드 (의미 검색에 도움)
            f"OBJECT: {level1}\n"
            f"PART_NAME: {desc}\n"
            f"PART_NO: {pno}\n"
        )
        rag_obj["embedding_text"] = embedding_text

        docs.append({
            "id": rag_obj["doc_id"],
            "document": embedding_text,
            "metadata": {
                "model": model,
                "dev_grade": dev_grade,
                "object": level1,
                "feature": feature,
                "level1": level1,
                "schema": "structured_v3",
                "rag_json": json.dumps(rag_obj, ensure_ascii=False),
            }
        })

    return docs

def upsert_structured_docs(docs: list[dict]) -> int:
    """
    구조화 컬렉션(v3)에 upsert.

    D-012: chromadb 기본 embedding 함수 대신 우리 bge-m3 (Ollama)로 직접
    임베딩해 전달. Windows에서 onnxruntime path 문제로 [Errno 22] 회피.
    """
    if not docs:
        return 0

    col = get_structured_collection()
    ids = [d["id"] for d in docs]
    documents = [d["document"] for d in docs]   # embedding_text
    metadatas = [d["metadata"] for d in docs]   # rag_json 포함

    # bge-m3로 documents 일괄 임베딩 (ENABLE_EMBEDDING은 rag_client에서 자동 활성화)
    from src.embed.embedder import embed_texts
    embeddings = embed_texts(documents)

    col.upsert(  # type: ignore
        ids=ids,
        documents=documents,
        embeddings=embeddings,
        metadatas=metadatas,
    )
    return len(docs)

def make_base_snapshot(df: pd.DataFrame, file_name: str = ""):
    df2, colmap = _detect_columns(df)

    required = ["part_name", "part_no", "module", "lvl"]
    missing_required = [c for c in required if c not in df2.columns]

    bom_path_list = build_bom_path_from_module_lvl(df2)

    rows = []
    path_index = {}
    for i, r in df2.iterrows():
        bp = bom_path_list[i]
        row = {
            "row_id": int(i),
            "part_name": str(r.get("part_name", "")).strip(),
            "part_no": str(r.get("part_no", "")).strip(),
            "module": str(r.get("module", "")).strip(),
            "lvl": str(r.get("lvl", "")).strip(),
            "bom_path": bp,
            "bom_tokens": [x.strip() for x in bp.split(">") if x.strip()],
        }
        rows.append(row)
        path_index.setdefault(bp, []).append(int(i))

    snap = BaseMasterSnapshot(
        meta={"file_name": file_name},
        rows=rows,
        part_index={},
        path_index=path_index,
        bom_tree={},
        quality={
            "file_name": file_name,
            "detected_columns": colmap,
            "missing_required": missing_required,
            "n_rows": len(df2),
            "n_unique_paths": len(set(bom_path_list)),
        },
    )
    return snap

# -------------------------------
# 3) SRC 블록 자동 생성
#  - 사용자는 변경점만 입력
#  - Agent가 SRC 블록 형태로 패키징해서 다음 단계(query/prompt)에 넣음
# -------------------------------
def build_src_block(change_items: list[str]) -> str:
    items = [x.strip() for x in (change_items or []) if (x or "").strip()]
    lines = ["SRC:"]
    for i, it in enumerate(items, 1):
        lines.append(f"- {i}. {it}")
    return "\n".join(lines)


_UI_FORBIDDEN_TOKENS = ("http", "https", "www", "sharepoint", "onedrive", "teams", "drive", "path", "folder")


def _has_forbidden_ui_text(text: Any) -> bool:
    s = str(text or "").lower()
    return any(tok in s for tok in _UI_FORBIDDEN_TOKENS)


def _safe_ui_text(text: Any) -> str:
    s = str(text or "").strip()
    return "" if _has_forbidden_ui_text(s) else s


def _is_meaningful_part_no(text: Any) -> bool:
    s = re.sub(r"\s+", "", str(text or "").strip().upper())
    return bool(s and s not in {"TBD", "N/A", "NA", "NONE", "정보없음"})


_REMARK_BANNED_PATTERNS = (
    "낮은 확신도", "확신도", "low confidence", "confidence", "불확실", "추정",
)


def _sanitize_user_remark_text(text: Any) -> str:
    raw = str(text or "").strip()
    if not raw:
        return ""
    frags = re.split(r"\s*\|\s*", raw)
    kept = []
    for f in frags:
        s = str(f or "").strip()
        if not s:
            continue
        low = s.lower()
        if any(p in low for p in _REMARK_BANNED_PATTERNS):
            continue
        kept.append(s)
    return " | ".join(kept).strip(" |")


_REGION_US_KW = ("북미", "미국", "120v", "60hz", "ul", "csa", "doe", "energy star")
_REGION_EU_KW = ("유럽", "eu", "ce", "230v", "50hz", "iec", "en", "ukca", "erp")


def _detect_region_tag_from_text(text: Any) -> str:
    t = str(text or "").lower()
    has_us = any(k in t for k in _REGION_US_KW)
    has_eu = any(k in t for k in _REGION_EU_KW)
    if has_us and not has_eu:
        return "US"
    if has_eu and not has_us:
        return "EU"
    if has_us and has_eu:
        return "GLOBAL"
    return "UNKNOWN"


def _normalize_region_value(region: Any, hint_text: str = "") -> str:
    raw = str(region or "").strip()
    low = raw.lower()

    if raw == "유럽" or any(k in low for k in ["유럽", "eu", "europe"]):
        return "유럽"
    if raw == "미국" or any(k in low for k in ["미국", "북미", "us", "usa", "north america"]):
        return "미국"
    if raw == "글로벌" or low == "global":
        return "글로벌"

    tag = _detect_region_tag_from_text(hint_text or raw)
    if tag == "EU":
        return "유럽"
    if tag == "US":
        return "미국"
    return "글로벌"


def _infer_project_region() -> str:
    ss = st.session_state
    ppt = ss.get("ppt_extraction_result") or {}
    pm = ppt.get("project_meta") or {}
    hint = " ".join([
        str(pm.get("target_country") or ""),
        str(pm.get("rating") or ""),
        str(ss.get("target_country") or ""),
        str(ss.get("rating") or ""),
        str(ss.get("region") or ""),
    ])
    preferred = (
        ss.get("target_country")
        or pm.get("target_country")
        or ss.get("region")
        or ""
    )
    return _normalize_region_value(preferred, hint)


def _project_region_hint_text() -> str:
    ss = st.session_state
    hint = []
    ppt = ss.get("ppt_extraction_result") or {}
    pm = ppt.get("project_meta") or {}
    hint.append(str(pm.get("target_country") or ""))
    hint.append(str(pm.get("rating") or ""))
    hint.append(str(ss.get("target_country") or ""))
    hint.append(str(ss.get("rating") or ""))
    return " ".join([x for x in hint if x])


def _is_region_compatible(region: str, tag: str) -> bool:
    r = _normalize_region_value(region)
    t = str(tag or "UNKNOWN").strip().upper()
    if r == "유럽":
        return t != "US"
    if r == "미국":
        return t != "EU"
    return True


def _doc_region_tag(doc: dict, project_hint_text: str = "") -> str:
    meta = doc.get("metadata") or {}
    txt = " ".join([
        str(doc.get("document") or ""),
        str(doc.get("text") or ""),
        str(meta),
    ])
    tag = _detect_region_tag_from_text(txt)
    if tag == "UNKNOWN" and project_hint_text:
        return _detect_region_tag_from_text(project_hint_text)
    return tag


def _filter_docs_by_region(docs: list[dict], region: str, project_hint_text: str = "") -> tuple[list[dict], int]:
    out = []
    dropped = 0
    for d in (docs or []):
        tag = _doc_region_tag(d, project_hint_text)
        if _is_region_compatible(region, tag):
            out.append(d)
        else:
            dropped += 1
    return out, dropped


def _part_region_tag(part: dict) -> str:
    txt = " ".join([
        str(part.get("part_name") or part.get("desc") or ""),
        str(part.get("rsn") or ""),
        str(part.get("chg") or ""),
        str(part.get("source_doc") or ""),
        str(part.get("sourcing_reason") or ""),
    ])
    return _detect_region_tag_from_text(txt)


def _filter_parts_by_region(parts: list[dict], region: str) -> tuple[list[dict], int]:
    out = []
    dropped = 0
    for p in (parts or []):
        tag = _part_region_tag(p)
        if _is_region_compatible(region, tag):
            out.append(p)
        else:
            dropped += 1
    return out, dropped


_PLACEHOLDER_PNO_WORDS = ("CONTROLLER", "ASSEMBLY", "ASSY", "PANEL", "SUB", "MODULE", "PART")


def _is_placeholder_part_no(text: Any) -> bool:
    s = str(text or "").strip().upper()
    if not s:
        return False
    s_compact = re.sub(r"[^A-Z0-9]", "", s)
    if not s_compact:
        return True
    if any(w in s for w in _PLACEHOLDER_PNO_WORDS):
        return True
    return bool(re.fullmatch(r"[A-Z]+", s_compact))


def _build_ui_meta_for_change_point(row: dict) -> tuple[str, str]:
    ctype = str(row.get("type") or "").strip()
    base_pno = str(row.get("base_part_no") or "").strip()
    new_pno = str(row.get("new_part_no") or "").strip()
    shared_src = str(row.get("shared_source") or "").strip()
    is_shared = bool(row.get("is_shared_part")) or bool(shared_src) or ("공용" in [str(t).strip() for t in (row.get("tags") or [])])
    concerns = row.get("concerns") or []
    concern0 = str(concerns[0]).strip() if concerns else ""
    desc = str(row.get("description") or "")

    ui_type = ""
    ui_remark = ""

    # =====================================
    # 공용품 우선 처리: 공용이면 채번 필요 로직 무시
    # =====================================
    if is_shared:
        if ctype == "NEW":
            ui_type = "신규(품번확정) · 공용"
            ui_remark = f"공용품 P/N: {new_pno}" if _is_meaningful_part_no(new_pno) else "공용품 | 품번 확인 필요"
        elif ctype == "Changing":
            ui_type = "변경 · 공용"
            if _is_meaningful_part_no(base_pno) and _is_meaningful_part_no(new_pno) and base_pno != new_pno:
                ui_remark = f"Base: {base_pno} → New: {new_pno} (공용)"
            else:
                ui_remark = "공용품"
        else:
            ui_type = "공용"
            ui_remark = "공용품"
    else:
        # =====================================
        # 공용 아닌 경우: 종전 로직
        # =====================================
        is_new_need = (
            ctype == "NEW"
            and (
                not _is_meaningful_part_no(new_pno)
                or "품번추가" in desc
                or "채번" in desc
            )
        )

        if is_new_need:
            ui_type = "신규(채번필요)"
            ui_remark = "채번 필요"
        elif ctype == "Changing":
            if _is_meaningful_part_no(base_pno) and _is_meaningful_part_no(new_pno) and base_pno != new_pno:
                ui_type = "변경(대체/사양)"
                ui_remark = f"Base: {base_pno} → New: {new_pno}"
            else:
                ui_type = "변경(구조/형상)"
                ui_remark = ""
        elif ctype == "NEW":
            ui_type = "신규(품번확정)"
            ui_remark = f"P/N: {new_pno}" if _is_meaningful_part_no(new_pno) else ""

    if concern0:
        ui_type = f"{ui_type} · 확인필요" if ui_type else "확인필요"
        ui_remark = f"{ui_remark} | 확인필요: {concern0}" if ui_remark else f"확인필요: {concern0}"

    ui_type = _safe_ui_text(ui_type)
    ui_remark = _sanitize_user_remark_text(_safe_ui_text(ui_remark))
    return ui_type, ui_remark


def _infer_product_type_from_model(model_text: str) -> str:
    m = str(model_text or "").strip().upper()
    if m.startswith("W"):
        return "Built-in Oven"
    if m.startswith("M"):
        return "Microwave Oven"
    return "Oven/Microwave"


RELATED_PARTS_ANALYSIS_PROMPT_TMPL = """
당신은 오븐/전자레인지 제품 개발 전문가입니다.

[역할]
사용자가 입력한 변경점 목록을 보고, 각 변경점이 파급시킬 수 있는
연관 부품/모듈을 추론하세요.

[입력 정보]
- 제품군: {product_type}
- 출처 모델: {source_model}
- 대상 모델: {target_model}
- 주요 차이: {key_diff}
- 사용자 입력 변경점:
{change_points}

[분석 규칙]
1. 각 변경점에 대해, 해당 변경으로 인해 함께 변경될 가능성이 있는
     연관 부품/모듈을 추론하세요.
2. 추론 근거를 간단히 함께 제시하세요.
3. 주의사항:
     - 사용자가 이미 변경점으로 명시한 항목은 다시 언급하지 마세요.
     - 확실하지 않은 항목은 "가능성 있음" 수준으로 표현하세요.
     - 실제로 연동 변경될 확률이 높은 것만 최대 5개 이내로 제시하세요.
     - 구조물(Cavity, Door Glass, Panel, Frame 등) 사이즈 관련 이슈는 다루지 마세요.

[출력 형식 - 반드시 JSON]
{{
    "analysis": [
        {{
            "change_point": "사용자가 입력한 변경점 원문",
            "related_parts": [
                {{
                    "part_category": "연관 부품/모듈 카테고리명",
                    "reason": "파급 추론 근거 (1문장)",
                    "likelihood": "high" | "medium"
                }}
            ]
        }}
    ],
    "additional_check": "변경점 간 상호 영향이 있으면 여기에 기술. 없으면 null"
}}
""".strip()


def build_related_parts_analysis_prompt(
        change_items: list[str],
        source_model: str,
        target_model: str,
        key_diff: str,
        product_type: str = "",
) -> str:
        ptype = str(product_type or "").strip() or _infer_product_type_from_model(target_model)
        points = [f"- {str(x).strip()}" for x in (change_items or []) if str(x or "").strip()]
        if not points:
                points = ["- (변경점 없음)"]
        return RELATED_PARTS_ANALYSIS_PROMPT_TMPL.format(
                product_type=ptype,
                source_model=str(source_model or ""),
                target_model=str(target_model or ""),
                key_diff=str(key_diff or ""),
                change_points="\n".join(points),
        )


PPT_CHANGE_EXTRACTION_PROMPT_TMPL = """
당신은 LG전자 오븐/전자레인지 개발 프로세스 전문가입니다.
사용자가 업로드한 "개발 유형 및 등급 확정 심의회" PPT에서
변경부품리스트 생성에 필요한 정보를 자동 추출합니다.

═══════════════════════════════════════════
[STEP 1] 개발 PJT 메타정보 추출
═══════════════════════════════════════════
PPT 내 "개발 PJT 개요" 영역에서 아래 필드를 추출하세요:

    - 제품군 (예: Built-in Oven)
    - PJT명 (예: 24인치오븐_Extra grade)
    - Base 모델 (예: WSED7667M)
    - 출시 국가 (예: 유럽)
    - 정격 (예: 220V~240V/50Hz)
    - 용량 (예: 76L)
    - 확정 등급 (예: B)
    - 개발 유형 (예: NPI3.0)

출력 형식:
{{
    "project_meta": {{
        "product_type": "...",
        "project_name": "...",
        "base_model": "...",
        "target_country": "...",
        "rating": "...",
        "capacity": "...",
        "dev_grade": "...",
        "dev_type": "..."
    }}
}}

═══════════════════════════════════════════
[STEP 2] 모듈별 변경점 요약 추출
═══════════════════════════════════════════
"유첨1. 개발 변경점 상세 작성 내용" 중
모듈 요약 테이블(No / Module명 / 주요 변경점)에서 추출하세요.

이 테이블은 보통 Exploded View 이미지와 함께 있으며,
①②③... 번호로 모듈이 구분됩니다.

출력 형식:
{{
    "module_summary": [
        {{
            "no": "①",
            "module": "Cavity",
            "changes": [
                "BLDC Conv. Fan Motor 적용 (AC → BLDC) - 30인치 SKS 공용",
                "Fan Blade 신규 적용 (30인치 SKS 공용)",
                "Fan Cover 금형 신작"
            ]
        }}
    ]
}}

═══════════════════════════════════════════
[STEP 3] 상세 변경점 테이블 추출
═══════════════════════════════════════════
"유첨1" 중 아래 컬럼 구조의 상세 테이블에서 추출하세요:
    구분 | No | Part | 변경 내역(변경 전→변경 후) | 변경 사유 | 걱정점

구분은 다음 카테고리로 나뉩니다:
    - 기구 (4M 변경포함)
    - 제어
    - ThinQ App

출력 형식:
{{
    "detailed_changes": [
        {{
            "category": "기구",
            "no": 1,
            "part": "Conv. Motor",
            "change_detail": "AC motor → BLDC motor",
            "change_reason": "균일 가열 성능 및 요리 시간 단축",
            "concern": "요리 성능 및 온도 정밀도",
            "tags": ["BLDC", "모터"]
        }}
    ]
}}

tags 필드 규칙:
- 변경 내역과 Part명에서 핵심 키워드를 자동 태깅하세요.
- 태그는 나중에 BOM 검색 시 매칭 키워드로 사용됩니다.
- 빈 행(Part, 변경 내역이 모두 비어있는 행)은 무시하세요.

═══════════════════════════════════════════
[STEP 4] 모듈별 상세 변경점 추출 (부품번호 포함)
═══════════════════════════════════════════
"유첨1. 개발 변경점_XXX Assembly" 슬라이드들에서 추출하세요.
이 슬라이드들은 Module | Base | New model | Remark 구조입니다.

출력 형식:
{{
    "module_details": [
        {{
            "module": "Cavity",
            "sub_module": "Convection Fan Motor",
            "base": {{
                "description": "AC Motor + Bracket → CCW",
                "part_no": "EAU65078501"
            }},
            "new": {{
                "description": "BLDC Motor + Bracket (30인치 SKS 공용) → 정/역회전, 가변풍량",
                "part_no": "4810W1N060B",
                "is_shared": true,
                "shared_with": "30인치 SKS"
            }},
            "remark": "Inverter Pro Bake BLDC 모터 개발 - CMR 요리성능 개선"
        }}
    ]
}}

is_shared 판별 규칙:
- "공용", "공용화", "수평전개", "동일 적용" 등의 키워드가 있으면
    is_shared = true, shared_with에 공용 대상 모델/사이즈를 기록
- 해당 키워드가 없으면 is_shared = false

═══════════════════════════════════════════
[STEP 5] 예상 우려점 추출
═══════════════════════════════════════════
"유첨2. 개발 예상 우려점 검토 결과" 테이블에서 추출하세요.
구분: 기구/구조 부품 변경점 | 전장/PCB품 변경점 | 부자재 등 기타 변경점

출력 형식:
{{
    "risk_review": [
        {{
            "category": "기구/구조 부품 변경점",
            "concern": "...",
            "countermeasure": "...",
            "attachment": "..."
        }}
    ]
}}
※ 내용이 비어있는 경우 "내용 없음"으로 표기하세요.

═══════════════════════════════════════════
[STEP 6] 변경부품리스트 생성용 변경점 정리
═══════════════════════════════════════════
위 STEP 2~4의 결과를 종합하여,
변경부품리스트 자동 생성 시스템에 전달할 최종 변경점 목록을
아래 형식으로 정리하세요:

{{
    "change_points_for_bom": [
        {{
            "id": 1,
            "description": "Convection Motor BLDC 적용 (AC → BLDC)",
            "module": "Cavity",
            "type": "스펙변경",
            "base_part_no": "EAU65078501",
            "new_part_no": "4810W1N060B",
            "is_shared_part": true,
            "shared_source": "30인치 SKS",
            "related_parts": [
                "Fan Blade (MDG63965901)",
                "Fan Cover (MCK71660202 → 신규)"
            ],
            "concerns": ["요리 성능 및 온도 정밀도"]
        }}
    ]
}}

type 분류 기준:
- 스펙변경: 기존 부품의 사양이 변경 (예: AC→BLDC, 4.3"→6.8")
- 신규추가: 기존에 없던 부품/모듈 추가 (예: Camera 추가)
- 구조변경: 조립 구조, 체결 방식 변경 (예: Bracket 개조)
- 사이즈변경: 치수/크기 변경

is_shared_part가 true인 부품은
변경부품리스트에서 "공용 적용"으로 표시하고,
신규 개발이 아님을 명시하세요.

═══════════════════════════════════════════
[주의사항]
═══════════════════════════════════════════
1. PPT에서 텍스트를 추출할 때 취소선(~~텍스트~~)이 있으면
     해당 내용은 삭제된 것으로 간주하고 무시하세요.
     취소선 아래 새로 작성된 내용을 최종값으로 사용
2. "변경 없음" 태그가 있는 슬라이드는 해당 모듈의
     기존 내용이 유지됨을 의미합니다.
3. "추가 변경점" 태그가 있는 슬라이드는
     기존 심의 후 추가된 변경점이므로 반드시 포함하세요.
4. 이미지, Feature List 등 비텍스트 첨부는
     존재 여부만 기록하세요 (내용 추출 불필요).
5. 추출한 정보에 대해 절대 임의로 부품번호나 스펙을 생성하지 마세요.
     PPT에 없는 정보는 "TBD" 또는 "정보 없음"으로 표기하세요.

아래 프로젝트 컨텍스트를 참고하세요:
- 제품군(추정): {product_type}
- Base 모델(세션): {source_model}
- Target 모델(세션): {target_model}
- 주요 차이(세션): {key_diff}
- 개발등급(세션): {dev_grade}
- 변경점(세션):
{change_points}
""".strip()


def build_ppt_change_extraction_prompt(
    change_items: list[str],
    source_model: str,
    target_model: str,
    key_diff: str,
    dev_grade: str,
    product_type: str = "",
) -> str:
    ptype = str(product_type or "").strip() or _infer_product_type_from_model(target_model)
    points = [f"- {str(x).strip()}" for x in (change_items or []) if str(x or "").strip()]
    if not points:
        points = ["- (변경점 없음)"]
    return PPT_CHANGE_EXTRACTION_PROMPT_TMPL.format(
        product_type=ptype,
        source_model=str(source_model or ""),
        target_model=str(target_model or ""),
        key_diff=str(key_diff or ""),
        dev_grade=str(dev_grade or ""),
        change_points="\n".join(points),
    )


def _split_lines_text(v: str) -> list[str]:
    txt = str(v or "")
    parts = re.split(r"[\n\r•·]+", txt)
    out = []
    for p in parts:
        s = re.sub(r"\s+", " ", p).strip(" -\t")
        if s:
            out.append(s)
    return out


def _extract_part_no_candidates(text: str) -> list[str]:
    cands = re.findall(r"\b[A-Z0-9]{8,14}\b", str(text or "").upper())
    out = []
    for c in cands:
        if c not in out:
            out.append(c)
    return out


def _strip_part_no_from_desc(text: str, pnos: list[str]) -> str:
    s = str(text or "")
    for p in pnos or []:
        s = re.sub(re.escape(p), "", s, flags=re.IGNORECASE)
    return re.sub(r"\s+", " ", s).strip(" -:|\t")


def _guess_change_type(part: str, detail: str, base_desc: str, new_desc: str, remark: str) -> str:
    t = " ".join([part, detail, base_desc, new_desc, remark]).lower()
    if any(k in t for k in ["inch", "인치", "사이즈", "size", "mm", "치수", "폭", "높이"]):
        return "사이즈변경"
    if any(k in t for k in ["bracket", "체결", "구조", "layout", "개조", "frame", "조립"]):
        return "구조변경"
    if any(k in t for k in ["신규", "추가", "add", "장착", "new"]):
        return "신규추가"
    return "스펙변경"


def _auto_tags(part: str, detail: str, limit: int = 6) -> list[str]:
    txt = f"{part} {detail}".upper()
    found = []
    seed = [
        "BLDC", "AC", "DC", "MOTOR", "FAN", "CAMERA", "HARNESS", "PCB", "UI",
        "BRACKET", "INVERTER", "HEATER", "SENSOR", "DISPLAY", "TOUCH", "THINQ",
        "도어", "카메라", "모터", "하네스", "제어", "앱", "팬", "브라켓", "센서",
    ]
    for s in seed:
        if s in txt and s not in found:
            found.append(s)
    toks = re.findall(r"[A-Z][A-Z0-9\-]{1,15}|[가-힣]{2,12}", txt)
    stop = {"PART", "CHANGE", "MODEL", "BASE", "NEW", "NO", "AND", "THE", "변경", "내역"}
    for t in toks:
        if t in stop or t.isdigit() or t in found:
            continue
        found.append(t)
        if len(found) >= limit:
            break
    return found[:limit]


def _ppt_text_from_shape(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    chunks = []
    tf = shape.text_frame
    for p in tf.paragraphs:
        run_text = []
        for r in p.runs:
            try:
                if getattr(r.font, "strike", False):
                    continue
            except Exception:
                pass
            run_text.append(r.text or "")
        line = "".join(run_text).strip()
        if line:
            chunks.append(line)
    return "\n".join(chunks).strip()


def _norm_header(h: str) -> str:
    return re.sub(r"[^A-Z0-9가-힣]", "", str(h or "").upper())


def _sanitize_product_type(v: Any) -> str:
    s = re.sub(r"\s+", " ", str(v or "")).strip(" -:\t")
    if not s:
        return ""
    bad_markers = ["동일 사업부", "다르면", "으로 봄", "EX:", "예:", "정의", "설명", "기준", "참고", "해당"]
    su = s.upper()
    if any(b.upper() in su for b in bad_markers):
        return ""
    if any(x in s for x in [".", "!", "?", ";", "(ex", "예)"]):
        return ""
    if " " in s and len(s) > 24:
        return ""
    if len(s) > 40:
        return ""
    return s


def _extract_product_type_from_page1(text_p1: str) -> str:
    lines = [re.sub(r"\s+", " ", x).strip() for x in str(text_p1 or "").splitlines() if str(x or "").strip()]
    cands: list[str] = []

    # 1) key:value 또는 key|value 패턴 우선
    p1 = [
        r"(?:^|\|)\s*제품군\s*(?:\||[:：])\s*([^|\n\r]+)",
        r"(?:^|\|)\s*PRODUCT(?:\s*TYPE)?\s*(?:\||[:：])\s*([^|\n\r]+)",
    ]
    for p in p1:
        for m in re.finditer(p, text_p1, flags=re.IGNORECASE):
            v = _sanitize_product_type(m.group(1))
            if v:
                cands.append(v)

    # 2) 테이블 행 형태: ... | 제품군 | 값 | ...
    for ln in lines:
        if "제품군" in ln or "PRODUCT" in ln.upper():
            toks = [t.strip() for t in ln.split("|") if t.strip()]
            for i, t in enumerate(toks):
                tu = t.upper()
                if t == "제품군" or tu in ("PRODUCT", "PRODUCT TYPE"):
                    if i + 1 < len(toks):
                        v = _sanitize_product_type(toks[i + 1])
                        if v:
                            cands.append(v)

    # 2-1) 인접 라인 패턴: "제품군" 다음 줄이 값인 경우
    for i, ln in enumerate(lines):
        lnu = ln.upper()
        is_key_only = _norm_header(ln) in {"제품군", "PRODUCT", "PRODUCTTYPE"}
        if is_key_only or ("제품군" in ln) or ("PRODUCT TYPE" in lnu):
            if i + 1 < len(lines):
                nxt = lines[i + 1]
                # 메타 키로 보이는 라인은 건너뜀
                if _norm_header(nxt) not in {"PJT명", "PROJECTNAME", "BASE모델", "BASEMODEL", "개발등급", "확정등급", "개발유형"}:
                    v = _sanitize_product_type(nxt)
                    if v:
                        cands.append(v)

    # 3) 후보 점수화: 짧고 제품명 키워드가 있는 값 우선
    def _score(v: str) -> int:
        s = str(v or "").strip()
        su = s.upper()
        score = 0
        # BUILT-IN을 최우선으로 처리
        if "BUILT-IN" in su:
            score += 50
        if any(k in su for k in ["OVEN", "MICROWAVE", "레인지", "전자레인지", "오븐"]):
            score += 10
        if len(s) <= 20:
            score += 4
        if " " not in s:
            score += 2
        if any(k in su for k in ["동일사업부", "다르면", "설명", "정의", "EX", "예"]):
            score -= 100
        return score

    cands = [c for c in cands if c]
    if not cands:
        return ""
    cands = sorted(cands, key=lambda x: (_score(x), -len(x)), reverse=True)
    return cands[0]


def _extract_product_type_from_anywhere(text: str) -> str:
    lines = [re.sub(r"\s+", " ", x).strip() for x in str(text or "").splitlines() if str(x or "").strip()]
    cands: list[str] = []

    # 제품군/PRODUCT TYPE 키 기준 키|값 형태만 허용 (자유문장 하드코딩 금지)
    for ln in lines:
        up = ln.upper()
        if ("제품군" not in ln) and ("PRODUCT" not in up):
            continue
        toks = [t.strip() for t in ln.split("|") if t.strip()]
        for i, t in enumerate(toks):
            tu = t.upper()
            if t == "제품군" or tu in ("PRODUCT", "PRODUCT TYPE"):
                if i + 1 < len(toks):
                    v = _sanitize_product_type(toks[i + 1])
                    if v:
                        cands.append(v)

    # 인접 라인 패턴 허용: KEY 라인 다음 라인 값
    for i, ln in enumerate(lines):
        lnu = ln.upper()
        is_key_only = _norm_header(ln) in {"제품군", "PRODUCT", "PRODUCTTYPE"}
        if is_key_only or ("제품군" in ln) or ("PRODUCT TYPE" in lnu):
            if i + 1 < len(lines):
                v = _sanitize_product_type(lines[i + 1])
                if v:
                    cands.append(v)

    def _score(v: str) -> int:
        s = str(v or "").strip()
        su = str(v or "").upper()
        score = 0
        if any(k in su for k in ["OVEN", "MICROWAVE", "BUILT-IN", "BUILT IN", "RANGE", "COOKTOP", "REFRIGERATOR", "WASHER", "DRYER", "DISHWASHER"]):
            score += 10
        if len(s) <= 24:
            score += 3
        if any(k in su for k in ["동일사업부", "다르면", "설명", "정의", "EX", "예", "기준", "참고"]):
            score -= 100
        return score

    cands = [c for c in cands if c]
    if not cands:
        return ""
    cands.sort(key=lambda x: (_score(x), -len(x)), reverse=True)
    return cands[0]


def _classify_slide_role(slide_lines: list[str], table_headers: list[str]) -> dict:
    text = "\n".join(slide_lines or [])
    hdr = " | ".join(table_headers or [])
    blob = f"{text}\n{hdr}".upper()

    overview_hits = 0
    detail_hits = 0
    ref_hits = 0

    if "개발 PJT 개요".upper() in blob:
        overview_hits += 3
    if "심의 결과".upper() in blob or "확정".upper() in blob:
        overview_hits += 3
    for k in ["개발 유형", "개발 등급", "제품군", "BASE MODEL", "BASE 모델", "PJT명", "PROJECT NAME"]:
        if k.upper() in blob:
            overview_hits += 1

    for k in ["유첨", "개발 변경점 상세", "CAVITY", "DOOR", "CONTROLLER ASSEMBLY", "구분", "변경 내역", "BASE", "NEW"]:
        if k.upper() in blob:
            detail_hits += 1

    for k in ["개발 등급 분류 기준", "운영 기준", "부표", "기준", "절차", "정의", "상세", "참고"]:
        if k.upper() in blob:
            ref_hits += 1

    # 최우선 규칙: 개요 키워드 동시 출현은 신뢰도 최고
    if "개발 PJT 개요".upper() in blob and "심의 결과".upper() in blob:
        return {"role": "OVERVIEW", "reason": "overview_high_conf"}

    if overview_hits >= max(4, detail_hits + 1) and overview_hits >= ref_hits:
        return {"role": "OVERVIEW", "reason": "overview_score"}

    if detail_hits >= max(2, ref_hits + 1):
        return {"role": "DETAIL", "reason": "detail_score"}

    if ref_hits >= 2:
        return {"role": "REFERENCE", "reason": "reference_score"}

    return {"role": "UNKNOWN", "reason": "weak_signal"}


def _extract_project_meta(lines: list[str], ctx: dict, page1_lines: list[str] | None = None) -> dict:
    text = "\n".join(lines)
    text_p1 = "\n".join(page1_lines or [])
    text_for_ptype = f"{text_p1}\n{text}"
    row = {
        "product_type": "정보 없음",
        "project_name": "정보 없음",
        "base_model": str(ctx.get("source_model") or "정보 없음"),
        "target_country": "정보 없음",
        "rating": "정보 없음",
        "capacity": "정보 없음",
        "dev_grade": str(ctx.get("dev_grade") or "정보 없음"),
        "dev_type": "정보 없음",
    }
    kv_patterns = {
        "project_name": [r"PJT명\s*[:：]\s*(.+)", r"PROJECT\s*NAME\s*[:：]\s*(.+)"],
        "base_model": [r"BASE\s*모델\s*[:：]\s*([A-Z0-9]{8,14})", r"BASE\s*MODEL\s*[:：]\s*([A-Z0-9]{8,14})"],
        "target_country": [r"출시\s*국가\s*[:：]\s*(.+)", r"COUNTRY\s*[:：]\s*(.+)"],
        "rating": [r"정격\s*[:：]\s*(.+)", r"RATING\s*[:：]\s*(.+)"],
        "capacity": [r"용량\s*[:：]\s*(.+)", r"CAPACITY\s*[:：]\s*(.+)"],
        "dev_grade": [r"확정\s*등급\s*[:：]\s*([A-Z0-9\.\-]+)", r"등급\s*[:：]\s*([A-Z0-9\.\-]+)"],
        "dev_type": [r"개발\s*유형\s*[:：]\s*([A-Z0-9\.\-]+)", r"DEV\s*TYPE\s*[:：]\s*([A-Z0-9\.\-]+)"],
    }

    def _clean_meta_value(v: Any) -> str:
        s = re.sub(r"\s+", " ", str(v or "")).strip(" -:\t")
        if not s:
            return ""
        if s.upper() in {"TBD", "N/A", "NA", "-", "X", "O"}:
            return ""
        if any(k in s for k in ["정보 없음", "내용 없음"]):
            return ""
        return s

    def _extract_from_pipe_lines(src_lines: list[str], key_aliases: list[str]) -> str:
        alias_norm = [_norm_header(a) for a in key_aliases]
        for ln in src_lines or []:
            toks = [re.sub(r"\s+", " ", t).strip() for t in str(ln or "").split("|") if str(t or "").strip()]
            if len(toks) < 2:
                continue
            for i, t in enumerate(toks):
                tn = _norm_header(t)
                if any(a and (a == tn or a in tn) for a in alias_norm):
                    if i + 1 < len(toks):
                        cand = _clean_meta_value(toks[i + 1])
                        if cand:
                            return cand
        return ""

    for k, pats in kv_patterns.items():
        for p in pats:
            m = re.search(p, text, flags=re.IGNORECASE)
            if m and str(m.group(1)).strip():
                row[k] = str(m.group(1)).strip()
                break

    # 표 형태(키|값)에서 메타를 추출한다. page1을 최우선으로 보고 없으면 전체 텍스트를 본다.
    p1_lines = [x for x in str(text_p1 or "").splitlines() if str(x or "").strip()]
    all_src_lines = [x for x in str(text or "").splitlines() if str(x or "").strip()]
    table_keys = {
        "project_name": ["PJT명", "PROJECT NAME", "프로젝트명"],
        "base_model": ["BASE 모델", "BASE MODEL", "BASE"],
        "target_country": ["출시 국가", "COUNTRY", "적용 국가"],
        "rating": ["정격", "RATING"],
        "capacity": ["용량", "CAPACITY"],
        "dev_grade": ["확정 등급", "개발 등급", "등급"],
        "dev_type": ["개발 유형", "DEV TYPE", "유형"],
    }
    for k, aliases in table_keys.items():
        if row.get(k) not in ("", "정보 없음"):
            continue
        cand = _extract_from_pipe_lines(p1_lines, aliases) or _extract_from_pipe_lines(all_src_lines, aliases)
        if cand:
            row[k] = cand

    # product_type는 1페이지 템플릿만 source-of-truth로 사용한다.
    p1_ptype = _extract_product_type_from_page1(text_p1)
    if p1_ptype:
        row["product_type"] = p1_ptype

    # 1페이지에서 못 잡으면 전체 텍스트에서 Built-in Oven/제품군 키를 보강 추출
    if row["product_type"] in ("", "정보 없음"):
        any_ptype = _extract_product_type_from_anywhere(text_for_ptype)
        if any_ptype:
            row["product_type"] = any_ptype

    if row["product_type"] == "정보 없음":
        # target_model이 비어있는 경우가 많아서 source_model/base_model까지 순차 fallback
        model_hint = str(
            ctx.get("target_model")
            or ctx.get("source_model")
            or row.get("base_model")
            or ""
        )
        row["product_type"] = str(_sanitize_product_type(ctx.get("product_type")) or _infer_product_type_from_model(model_hint))
    return row


def extract_change_review_from_pptx_bytes(pptx_bytes: bytes, ctx: dict | None = None) -> dict:
    ctx = ctx or {}
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("python-pptx 패키지가 필요합니다. 설치 후 다시 시도하세요.") from e

    import io
    prs = Presentation(io.BytesIO(pptx_bytes))

    discarded: list[dict] = []
    all_lines: list[str] = []
    all_table_text_lines: list[str] = []
    slide_lines_map: dict[int, list[str]] = {}
    slide_table_headers_map: dict[int, list[str]] = {}
    slide_table_text_map: dict[int, list[str]] = {}
    table_rows: list[dict] = []
    image_presence: list[dict] = []

    def _flat_cells(rows: list[list[str]]) -> list[str]:
        out = []
        for r in rows or []:
            for c in r or []:
                t = re.sub(r"\s+", " ", str(c or "")).strip()
                if t:
                    out.append(t)
        return out

    def _classify_table_candidate(rows: list[list[str]], slide_lines: list[str], slide: int, table_id: int) -> dict:
        if not rows:
            return {"type": "IGNORE", "why": "empty_table"}

        raw_header = [re.sub(r"\s+", " ", str(x or "")).strip() for x in (rows[0] or [])]
        header_norm = [_norm_header(x) for x in raw_header]
        cells = _flat_cells(rows)
        all_text = " ".join(raw_header + cells + (slide_lines or []))

        # 블랙리스트 키워드: 테이블 성격 자체 제외
        blk = ["적용여부", "과전압", "PDR", "ODR", "EVENT", "NPI", "ACTIVITY"]
        all_up = all_text.upper()
        if any(k.upper() in all_up for k in blk):
            return {"type": "IGNORE", "why": "blacklist_keyword"}

        # O/X 체크표 비율이 높은 경우 제외
        chk_tokens = {"O", "X", "OX", "(O,X)", "(O,X)", "(O, X)"}
        non_empty = [c for c in cells if str(c).strip()]
        chk_cnt = 0
        for c in non_empty:
            n = re.sub(r"\s+", "", str(c).upper())
            if n in chk_tokens:
                chk_cnt += 1
        if non_empty and (chk_cnt / max(len(non_empty), 1)) >= 0.20 and len(non_empty) >= 4:
            return {"type": "IGNORE", "why": "checkbox_ratio_high"}

        hraw_join = " | ".join(raw_header).upper()
        hnorm_join = "|".join(header_norm)
        hraw_compact = re.sub(r"[^A-Z0-9가-힣]", "", hraw_join.upper())
        hnorm_compact = re.sub(r"[^A-Z0-9가-힣]", "", hnorm_join.upper())

        def _has_any(src: str, kws: list[str]) -> bool:
            s = src.upper()
            return any(k.upper() in s for k in kws)

        def _has_any_compact(src_compact: str, kws: list[str]) -> bool:
            return any(re.sub(r"[^A-Z0-9가-힣]", "", k.upper()) in src_compact for k in kws)

        # MAIN_CHANGE_TABLE
        has_change = _has_any(hraw_join, ["변경내역", "변경 내역", "CHANGE"]) or _has_any_compact(hraw_compact, ["변경내역"])
        has_before_after = (
            _has_any(hraw_join, ["변경 전", "변경 후", "->", "→", "BEFORE", "AFTER"])
            or _has_any_compact(hraw_compact, ["변경전", "변경후"])
            or ("변경전" in hnorm_join or "변경후" in hnorm_join)
        )
        has_reason = _has_any(hraw_join, ["변경사유", "변경 사유", "REASON"]) or _has_any_compact(hraw_compact, ["변경사유"])
        has_concern = _has_any(hraw_join, ["걱정점", "걱정 점", "우려", "CONCERN"]) or _has_any_compact(hraw_compact, ["걱정점", "우려"])
        has_part = _has_any(hraw_join, ["PART", "부품"]) or _has_any_compact(hraw_compact, ["PART", "부품"])
        if has_change and has_before_after and has_reason and has_concern and has_part:
            return {"type": "MAIN_CHANGE_TABLE", "why": "whitelist_main_change"}

        # DETAIL_CHANGE_TABLE
        has_base = _has_any(hraw_join, ["BASE"])
        has_new = _has_any(hraw_join, ["NEW"])
        has_mod_or_remark = _has_any(hraw_join, ["MODULE", "모듈", "REMARK", "비고"])
        if has_base and has_new and has_mod_or_remark:
            return {"type": "DETAIL_CHANGE_TABLE", "why": "whitelist_detail_change"}

        # RISK_TABLE
        has_risk = _has_any(hraw_join, ["우려", "걱정", "CONCERN"]) or _has_any_compact(hraw_compact, ["우려", "걱정"])
        has_counter = _has_any(hraw_join, ["대책", "대응", "검증", "방안", "COUNTER"]) or _has_any_compact(hraw_compact, ["대책", "대응", "검증방안"])
        if has_risk and has_counter:
            return {"type": "RISK_TABLE", "why": "whitelist_risk"}

        return {"type": "IGNORE", "why": "not_whitelisted"}

    for s_idx, slide in enumerate(prs.slides, 1):
        slide_lines = []
        slide_table_headers: list[str] = []
        slide_table_texts: list[str] = []
        has_non_text_attach = False
        table_seq = 0
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False):
                txt = _ppt_text_from_shape(sh)
                if txt:
                    ls = _split_lines_text(txt)
                    slide_lines.extend(ls)
                    all_lines.extend(ls)
            if getattr(sh, "shape_type", None) in (13, 14):
                has_non_text_attach = True
            if getattr(sh, "has_table", False):
                table_seq += 1
                t = sh.table
                rows = []
                for r in t.rows:
                    row = []
                    for c in r.cells:
                        row.append(re.sub(r"\s+", " ", c.text or "").strip())
                    rows.append(row)
                if rows:
                    for rr in rows:
                        rrj = " | ".join([re.sub(r"\s+", " ", str(x or "")).strip() for x in rr if str(x or "").strip()])
                        if rrj:
                            slide_table_texts.append(rrj)
                            all_table_text_lines.append(rrj)
                    hdr_join = " | ".join([re.sub(r"\s+", " ", str(x or "")).strip() for x in (rows[0] or []) if str(x or "").strip()])
                    if hdr_join:
                        slide_table_headers.append(hdr_join)
                    table_rows.append({
                        "slide": s_idx,
                        "table_id": f"S{s_idx}-T{table_seq}",
                        "rows": rows,
                        "slide_lines": list(slide_lines),
                    })
        slide_lines_map[s_idx] = list(slide_lines)
        slide_table_headers_map[s_idx] = list(slide_table_headers)
        slide_table_text_map[s_idx] = list(slide_table_texts)
        image_presence.append({"slide": s_idx, "has_attachment": bool(has_non_text_attach)})

    slide_roles: dict[int, str] = {}
    slide_role_details: list[dict] = []
    for s_idx in sorted(slide_lines_map.keys()):
        cls = _classify_slide_role(slide_lines_map.get(s_idx) or [], slide_table_headers_map.get(s_idx) or [])
        role = str(cls.get("role") or "UNKNOWN")
        slide_roles[s_idx] = role
        slide_role_details.append({"slide": s_idx, "role": role, "reason": cls.get("reason", "")})

    overview_scope_lines: list[str] = []
    for s_idx, role in slide_roles.items():
        if role == "OVERVIEW":
            overview_scope_lines.extend(slide_lines_map.get(s_idx) or [])
            overview_scope_lines.extend(slide_table_text_map.get(s_idx) or [])
            overview_scope_lines.extend(slide_table_headers_map.get(s_idx) or [])

    # 개요/결론 슬라이드가 없으면 약한 fallback: 첫 슬라이드 사용
    if not overview_scope_lines and slide_lines_map:
        first_idx = sorted(slide_lines_map.keys())[0]
        overview_scope_lines.extend(slide_lines_map.get(first_idx) or [])
        overview_scope_lines.extend(slide_table_text_map.get(first_idx) or [])
        overview_scope_lines.extend(slide_table_headers_map.get(first_idx) or [])

    # 제품군/프로젝트 메타는 첫 슬라이드 템플릿을 source-of-truth로 우선 사용한다.
    page1_scope_lines: list[str] = []
    if slide_lines_map:
        first_idx = sorted(slide_lines_map.keys())[0]
        page1_scope_lines.extend(slide_lines_map.get(first_idx) or [])
        page1_scope_lines.extend(slide_table_text_map.get(first_idx) or [])
        page1_scope_lines.extend(slide_table_headers_map.get(first_idx) or [])

    # project_meta는 OVERVIEW + table cell 텍스트 중심으로 읽고, 전체 텍스트는 보조로 합친다.
    meta_scope_lines: list[str] = []
    meta_scope_lines.extend(overview_scope_lines)
    meta_scope_lines.extend(all_table_text_lines)
    meta_scope_lines.extend(all_lines)

    project_meta = _extract_project_meta(meta_scope_lines, ctx, page1_lines=page1_scope_lines or overview_scope_lines)

    module_summary: list[dict] = []
    detailed_changes: list[dict] = []
    module_details: list[dict] = []
    risk_review: list[dict] = []
    module_summary_hints: list[dict] = []

    for tb in table_rows:
        rows = tb["rows"]
        slide_no = int(tb.get("slide") or 0)
        slide_role = slide_roles.get(slide_no, "UNKNOWN")

        # 부표/기준 슬라이드는 결정 정보에서 제외
        if slide_role == "REFERENCE":
            flat = _flat_cells(rows)
            discarded.append({
                "why": "reference_slide_scope",
                "slide": tb.get("slide"),
                "table_id": tb.get("table_id"),
                "sample_text": " | ".join(flat[:6]) if flat else "",
            })
            continue

        cls = _classify_table_candidate(rows, tb.get("slide_lines") or [], tb.get("slide"), tb.get("table_id"))
        ttype = cls.get("type", "IGNORE")

        if ttype == "IGNORE":
            # DETAIL 슬라이드의 "Module명 | 주요 변경점" 요약표는 module 매핑 힌트로만 사용
            if slide_role == "DETAIL" and rows and rows[0]:
                hdr0 = [_norm_header(x) for x in rows[0]]
                has_mod_col = any(("MODULE" in h) or ("모듈" in h) for h in hdr0)
                has_change_col = any(("주요변경점" in h) or ("변경내역" in h) or ("CHANGE" in h) for h in hdr0)
                if has_mod_col and has_change_col:
                    m_i = next((i for i, h in enumerate(hdr0) if ("MODULE" in h or "모듈" in h)), 0)
                    c_i = next((i for i, h in enumerate(hdr0) if ("주요변경점" in h or "변경내역" in h or "CHANGE" in h)), min(1, len(hdr0)-1))
                    cur_mod = ""
                    for rr in rows[1:]:
                        mod = (rr[m_i] if m_i < len(rr) else "").strip() or cur_mod
                        cur_mod = mod or cur_mod
                        chg = (rr[c_i] if c_i < len(rr) else "").strip()
                        if mod and chg:
                            module_summary_hints.append({
                                "module": mod,
                                "text": chg,
                                "src": {"slide": tb.get("slide"), "table_id": tb.get("table_id")},
                            })

            flat = _flat_cells(rows)
            discarded.append({
                "why": cls.get("why", "ignored"),
                "slide": tb.get("slide"),
                "table_id": tb.get("table_id"),
                "sample_text": " | ".join(flat[:6]) if flat else "",
            })
            continue

        hdr = [_norm_header(x) for x in rows[0]]

        # 상세 변경 추출은 상세 슬라이드에서만 허용
        if slide_role != "DETAIL" and ttype in ("MAIN_CHANGE_TABLE", "DETAIL_CHANGE_TABLE", "RISK_TABLE"):
            flat = _flat_cells(rows)
            discarded.append({
                "why": "non_detail_slide_for_change_table",
                "slide": tb.get("slide"),
                "table_id": tb.get("table_id"),
                "sample_text": " | ".join(flat[:6]) if flat else "",
            })
            continue

        # MAIN_CHANGE_TABLE => STEP3 only
        if ttype == "MAIN_CHANGE_TABLE":
            c_i = next((i for i, h in enumerate(hdr) if ("구분" in h or "CATEGORY" in h)), 0)
            n_i = next((i for i, h in enumerate(hdr) if "NO" in h), 1)
            p_i = next((i for i, h in enumerate(hdr) if ("PART" in h or "부품" in h)), 2)
            d_i = next((i for i, h in enumerate(hdr) if ("변경내역" in h or "CHANGE" in h)), 3)
            r_i = next((i for i, h in enumerate(hdr) if ("변경사유" in h or "REASON" in h)), 4)
            g_i = next((i for i, h in enumerate(hdr) if ("걱정" in h or "CONCERN" in h)), 5)
            cur_cat = ""
            for r in rows[1:]:
                cat = (r[c_i] if c_i < len(r) else "").strip() or cur_cat
                cur_cat = cat or cur_cat
                no = (r[n_i] if n_i < len(r) else "").strip()
                part = (r[p_i] if p_i < len(r) else "").strip()
                detail = (r[d_i] if d_i < len(r) else "").strip()
                reason = (r[r_i] if r_i < len(r) else "").strip()
                concern = (r[g_i] if g_i < len(r) else "").strip()
                if not part and not detail:
                    discarded.append({
                        "why": "empty_part_and_detail_row",
                        "slide": tb.get("slide"),
                        "table_id": tb.get("table_id"),
                        "sample_text": " | ".join([str(x) for x in r[:6]]),
                    })
                    continue
                n_val: Any = no
                try:
                    n_val = int(re.sub(r"[^0-9]", "", no)) if re.search(r"[0-9]", no) else no
                except Exception:
                    n_val = no
                detailed_changes.append({
                    "category": cat or "정보 없음",
                    "discipline": cat or "정보 없음",
                    "no": n_val,
                    "part": part or "정보 없음",
                    "change_detail": detail or "정보 없음",
                    "change_reason": reason or "정보 없음",
                    "concern": concern or "내용 없음",
                    "tags": _auto_tags(part, detail),
                    "_src": {"slide": tb.get("slide"), "table_id": tb.get("table_id")},
                })
            continue

        # DETAIL_CHANGE_TABLE => STEP4 only
        if ttype == "DETAIL_CHANGE_TABLE":
            m_i = next((i for i, h in enumerate(hdr) if ("MODULE" in h or "모듈" in h)), 0)
            b_i = next((i for i, h in enumerate(hdr) if "BASE" in h), 1)
            n_i = next((i for i, h in enumerate(hdr) if "NEW" in h), 2)
            rm_i = next((i for i, h in enumerate(hdr) if ("REMARK" in h or "비고" in h)), min(3, len(hdr)-1))
            for r in rows[1:]:
                mod = (r[m_i] if m_i < len(r) else "").strip()
                base_txt = (r[b_i] if b_i < len(r) else "").strip()
                new_txt = (r[n_i] if n_i < len(r) else "").strip()
                remark = (r[rm_i] if rm_i < len(r) else "").strip()
                if not mod and not base_txt and not new_txt:
                    continue

                # 상세표 행 안전망: 체크표류 차단
                mod_up = mod.upper()
                chk_desc = re.sub(r"\s+", "", f"{new_txt or base_txt}").upper()
                if "과전압" in mod_up or chk_desc in {"O", "X", "(O,X)", "OX"}:
                    discarded.append({
                        "why": "detail_row_blacklist_or_checkbox",
                        "slide": tb.get("slide"),
                        "table_id": tb.get("table_id"),
                        "sample_text": f"{mod} | {base_txt} | {new_txt}",
                    })
                    continue

                base_pnos = _extract_part_no_candidates(base_txt)
                new_pnos = _extract_part_no_candidates(new_txt)
                mix_txt = f"{new_txt} {remark}".lower()
                is_shared = any(k in mix_txt for k in ["공용", "공용화", "수평전개", "동일 적용"])
                sw = ""
                if is_shared:
                    m_sw = re.search(r"(\d{2}\s*인치\s*[A-Z0-9가-힣 ]+|[A-Z0-9가-힣 ]+\s*공용)", f"{new_txt} {remark}", flags=re.IGNORECASE)
                    if m_sw:
                        sw = re.sub(r"\s+", " ", m_sw.group(1)).strip()
                module_details.append({
                    "module": mod or "정보 없음",
                    "sub_module": mod or "정보 없음",
                    "base": {
                        "description": _strip_part_no_from_desc(base_txt, base_pnos) or (base_txt or "정보 없음"),
                        "part_no": base_pnos[0] if base_pnos else "TBD",
                    },
                    "new": {
                        "description": _strip_part_no_from_desc(new_txt, new_pnos) or (new_txt or "정보 없음"),
                        "part_no": new_pnos[0] if new_pnos else "TBD",
                        "is_shared": bool(is_shared),
                        "shared_with": sw,
                    },
                    "remark": remark or "정보 없음",
                    "_src": {"slide": tb.get("slide"), "table_id": tb.get("table_id")},
                })
            continue

        # RISK_TABLE => STEP5 only
        if ttype == "RISK_TABLE":
            c_i = next((i for i, h in enumerate(hdr) if ("구분" in h or "CATEGORY" in h)), 0)
            o_i = next((i for i, h in enumerate(hdr) if ("우려" in h or "CONCERN" in h)), 1)
            m_i = next((i for i, h in enumerate(hdr) if ("대책" in h or "COUNTER" in h)), 2)
            a_i = next((i for i, h in enumerate(hdr) if ("첨부" in h or "ATTACH" in h)), min(3, len(hdr)-1))
            for r in rows[1:]:
                cat = (r[c_i] if c_i < len(r) else "").strip() or "정보 없음"
                concern = (r[o_i] if o_i < len(r) else "").strip() or "내용 없음"
                cm = (r[m_i] if m_i < len(r) else "").strip() or "내용 없음"
                at = (r[a_i] if a_i < len(r) else "").strip() or "내용 없음"
                if cat == "정보 없음" and concern == "내용 없음" and cm == "내용 없음":
                    continue
                risk_review.append({
                    "category": cat,
                    "concern": concern,
                    "countermeasure": cm,
                    "attachment": at,
                })

    # STEP2 (요약): 상세 변경점을 모듈/카테고리 기준으로 합성
    mod_map: dict[str, list[str]] = {}
    for dc in detailed_changes:
        mod = str(dc.get("category") or "정보 없음").strip()
        msg = str(dc.get("change_detail") or "").strip()
        if not msg:
            continue
        mod_map.setdefault(mod, [])
        if msg not in mod_map[mod]:
            mod_map[mod].append(msg)
    module_summary = [
        {"no": str(i + 1), "module": k, "changes": v}
        for i, (k, v) in enumerate(mod_map.items())
    ]

    # STEP6: MAIN_CHANGE_TABLE(detailed_changes) 기반 생성 + module_details 보완
    def _norm_txt(s: Any) -> str:
        return re.sub(r"\s+", " ", str(s or "").strip().upper())

    def _rule_based_module_from_change(dc: dict) -> str:
        scope = _norm_txt(f"{dc.get('part') or ''} {dc.get('change_detail') or ''} {dc.get('change_reason') or ''}")
        if not scope:
            return ""

        # 제품군 공통 규칙: 키워드 점수 높은 모듈을 우선 사용
        rules = {
            "Door Assembly": ["DOOR", "CAMERA", "HARNESS", "HINGE", "TRAY", "FRAME"],
            "Cavity": ["CAVITY", "CONV", "FAN", "MOTOR", "BLADE", "INLET", "OUTLET", "AIR"],
            "Controller Assembly": ["CONTROLLER", "LCD", "PANEL", "PCB", "UI", "OS", "S/W", "SW"],
        }
        best_mod = ""
        best_score = 0
        for mod, kws in rules.items():
            score = 0
            for kw in kws:
                if kw in scope:
                    score += 2
            # 강한 앵커 키워드 가중치
            if mod == "Door Assembly" and any(k in scope for k in ["HARNESS", "TRAY", "HINGE", "CAMERA"]):
                score += 2
            if mod == "Cavity" and any(k in scope for k in ["CONV", "FAN", "MOTOR", "BLADE"]):
                score += 2
            if mod == "Controller Assembly" and any(k in scope for k in ["LCD", "PCB", "PANEL", "UI", "OS"]):
                score += 2
            if score > best_score:
                best_mod = mod
                best_score = score
        return best_mod if best_score >= 3 else ""

    def _pick_module_from_summary(dc: dict) -> str:
        part = _norm_txt(dc.get("part"))
        detail = _norm_txt(dc.get("change_detail"))
        scope = f"{part} {detail}"
        stop_tokens = {"ASSY", "ASSEMBLY", "PART", "MODULE", "변경", "신규", "추가", "기구", "제어"}
        weak_tokens = {"COVER", "구조", "적용", "변경", "추가", "형상", "HOLE", "SIZE"}
        # 파트명 토큰을 더 신뢰하고, 일반 단어는 제외
        part_toks = [x for x in re.split(r"[^A-Z0-9가-힣]+", part) if len(x) >= 3 and x not in stop_tokens and x not in weak_tokens]
        detail_toks = [x for x in re.split(r"[^A-Z0-9가-힣]+", detail) if len(x) >= 4 and x not in stop_tokens and x not in weak_tokens]
        toks = part_toks + detail_toks
        best_mod = ""
        best_score = 0
        for it in module_summary_hints:
            mod = str(it.get("module") or "").strip()
            txt = _norm_txt(it.get("text"))
            if not mod or not txt:
                continue
            score = 0
            for tk in toks:
                if tk in txt:
                    score += 2
            if score > best_score:
                best_mod = mod
                best_score = score
        return best_mod if best_score >= 2 and len(toks) > 0 else ""

    def _pick_best_detail(dc: dict, preferred_module: str = "") -> dict | None:
        part = _norm_txt(dc.get("part"))
        detail = _norm_txt(dc.get("change_detail"))
        scope = f"{part} {detail}"
        best = None
        best_score = -1
        stop_tokens = {"ASSY", "ASSEMBLY", "PART", "MODULE", "변경", "신규", "추가", "기구", "제어"}
        
        # 부품명을 토큰화 (2자 이상)
        part_toks = [x for x in re.split(r"[^A-Z0-9가-힣]+", part) if len(x) >= 2 and x not in stop_tokens]
        
        for md in module_details:
            md_module = _norm_txt(md.get("module"))
            blob = " ".join([
                md_module,
                _norm_txt(md.get("sub_module")),
                _norm_txt((md.get("base") or {}).get("description")),
                _norm_txt((md.get("new") or {}).get("description")),
                _norm_txt(md.get("remark")),
            ])
            score = 0
            if preferred_module and _norm_txt(preferred_module) == md_module:
                score += 4
            # module명이 본문에 직접 등장하면 가중치 부여
            mod_toks = [x for x in re.split(r"[^A-Z0-9가-힣]+", md_module) if len(x) >= 3 and x not in stop_tokens]
            for mt in mod_toks:
                if mt in scope:
                    score += 2
            
            # 부품명 토큰: 3자 이상 +3, 2자 +1 (강화)
            for tok in part_toks:
                if tok in blob:
                    score += (3 if len(tok) >= 3 else 1)
            
            # 변경상세 토큰: 3자 이상만
            detail_toks = [x for x in re.split(r"[^A-Z0-9가-힣]+", detail) if len(x) >= 3 and x not in stop_tokens]
            for tok in detail_toks:
                if tok in blob:
                    score += 1
            
            if score > best_score:
                best = md
                best_score = score
        return best if best_score >= 1 else None

    def _infer_change_type(detail_text: Any, reason_text: Any, part_text: Any) -> str:
        blob = " ".join([str(detail_text or ""), str(reason_text or ""), str(part_text or "")]).upper()
        if any(k in blob for k in ["삭제", "제거", "미적용", "DELETED", "REMOVE"]):
            return "삭제"
        if any(k in blob for k in ["추가", "신규", "ADD", "NEW"]):
            return "NEW"
        return "Changing"

    change_points_for_bom = []
    for i, dc in enumerate(detailed_changes, 1):
        rule_mod = _rule_based_module_from_change(dc)
        summary_mod = _pick_module_from_summary(dc)
        preferred_mod = rule_mod or summary_mod
        md = _pick_best_detail(dc, preferred_module=preferred_mod)
        part = str(dc.get("part") or "정보 없음")
        detail = str(dc.get("change_detail") or "정보 없음")
        reason = str(dc.get("change_reason") or "정보 없음")
        concern = str(dc.get("concern") or "내용 없음")

        base_pno = "TBD"
        new_pno = "TBD"
        is_shared = False
        shared_src = ""
        related_parts: list[str] = []
        remark = ""
        discipline = str(dc.get("discipline") or dc.get("category") or "정보 없음").strip() or "정보 없음"
        module_name = "정보 없음"

        if preferred_mod:
            module_name = preferred_mod

        if md:
            md_module = str(md.get("module") or "").strip()
            if md_module:
                module_name = md_module
            base_pno = str((md.get("base") or {}).get("part_no") or "TBD")
            new_pno = str((md.get("new") or {}).get("part_no") or "TBD")
            is_shared = bool((md.get("new") or {}).get("is_shared"))
            shared_src = str((md.get("new") or {}).get("shared_with") or "")
            remark = str(md.get("remark") or "")
            rp = str(md.get("sub_module") or "").strip()
            if rp and _norm_txt(rp) != _norm_txt(part):
                related_parts.append(rp)

        # 구조 Assembly를 못 찾았을 때만 discipline(기구/제어)을 fallback으로 사용
        if module_name in ("", "정보 없음"):
            module_name = discipline

        desc = f"{part}: {detail}" if part and detail else (detail or part or "정보 없음")
        ctype = _infer_change_type(detail, reason, part)
        tags = list(dc.get("tags") or [])
        if is_shared and "공용" not in tags:
            tags.append("공용")

        row = {
            "id": i,
            "description": desc,
            "module": module_name,
            "discipline": discipline,
            "type": ctype,
            "base_part_no": base_pno,
            "new_part_no": new_pno,
            "is_shared_part": is_shared,
            "shared_source": shared_src,
            "related_parts": related_parts or ["정보 없음"],
            "concerns": [concern] if concern else ["내용 없음"],
            "tags": tags,
            "evidence": {
                "part": part,
                "change_reason": reason,
                "detail_src": dc.get("_src") or {},
                "discipline_src": "main_change_table.category",
                "module_detail_src": (md or {}).get("_src") if md else {},
            },
        }
        ui_type, ui_remark = _build_ui_meta_for_change_point(row)
        row["ui_type"] = ui_type
        row["ui_remark"] = ui_remark
        change_points_for_bom.append(row)

    # 최후 안전망: O/X 체크표 잔여 행 제거
    cp_filtered = []
    for r in change_points_for_bom:
        mod = str(r.get("module") or "")
        desc = re.sub(r"\s+", "", str(r.get("description") or "").upper())
        if "과전압" in mod:
            discarded.append({
                "why": "cp_drop_blacklist_module",
                "slide": None,
                "table_id": "final_cp",
                "sample_text": f"{mod} | {r.get('description','')}",
            })
            continue
        if desc in {"O", "X", "(O,X)", "OX"}:
            discarded.append({
                "why": "cp_drop_checkbox_token",
                "slide": None,
                "table_id": "final_cp",
                "sample_text": f"{mod} | {r.get('description','')}",
            })
            continue
        cp_filtered.append(r)

    return {
        "project_meta": project_meta,
        "module_summary": module_summary,
        "detailed_changes": detailed_changes,
        "module_details": module_details,
        "risk_review": risk_review,
        "change_points_for_bom": cp_filtered,
        "attachment_presence": image_presence,
        "slide_roles": slide_role_details,
        "discarded": discarded,
    }


def to_chroma_where(filters: dict | None) -> dict | None:
    """
    Chroma where 포맷으로 변환
    - 1개 조건: {"product": {"$eq": "W"}}
    - 다중 조건: {"$and": [{"product":{"$eq":"W"}}, {"prefix":{"$eq":"WSED"}}]}
    """
    if not filters:
        return None

    items = [(k, v) for k, v in filters.items() if v not in (None, "", [])]
    if not items:
        return None

    if len(items) == 1:
        k, v = items[0]
        return {k: {"$eq": v}}

    return {"$and": [{k: {"$eq": v}} for k, v in items]}


# [DB-1] retrieve_docs 래퍼 (filters 지원/미지원 모두 대응)
def retrieve_with_filters(query: str, top_k: int, filters: dict | None = None):
    """
    retrieve_docs가 filters(where)를 지원하면 where 포맷으로 넣고,
    형식/버전 문제로 실패하면 query에 토큰 섞어서 우회.
    """
    where = to_chroma_where(filters)

    # 1) where로 시도
    try:
        if where:
            return retrieve_docs(query, top_k=top_k, filters=where)
        else:
            return retrieve_docs(query, top_k=top_k)
    except (TypeError, ValueError):
        # 2) 우회: query에 토큰 삽입 + filters 없이 검색
        q2 = query
        if filters:
            for k, v in filters.items():
                if v:
                    q2 += f"\n{str(k).upper()}={v}"
        return retrieve_docs(q2, top_k=top_k)
    
# [DB-2] 문서 텍스트에서 모델/등급 추출
DOC_MODEL_RE = re.compile(r"\b([WL][A-Z0-9]{8})\b", re.I)
DOC_GRADE_RE = re.compile(r"(?:개발\s*등급|등급)\s*[:/ ]\s*([ABCD])\b|\(([ABCD])\)", re.I)

def doc_extract_model_and_grade(doc_text: str) -> tuple[str, str]:
    t = doc_text or ""
    m = DOC_MODEL_RE.search(t)
    model = m.group(1).upper() if m else ""
    g = DOC_GRADE_RE.search(t)
    grade = (g.group(1) or g.group(2)).upper() if g else ""
    return model, grade

## =========================================================
## [블록 B] Search + Evidence + Proposal
## =========================================================

import time
from collections import Counter, defaultdict

# -----------------------------
# 변경점 -> intent
# -----------------------------
OBJECT_ALIASES = {
    "DOOR": ["도어", "door"],
    "CONTROL": ["제어", "control", "panel", "조작"],
    "OVEN": ["오븐", "oven", "캐비티", "cavity"],
}

MAIN_PART_ALIASES = {
    "CAMERA": ["카메라", "camera", "cam"],
    "LED": ["led", "조명", "램프"],
    "HARNESS": ["하네스", "harness", "wire", "배선"],
    "PROBE": ["probe", "프로브", "3p", "3p probe", "탐침"],
}

def parse_change_intent(change_items: list[str]) -> dict:
    """
    변경점 -> target_object / action / main_change_keyword
    """
    text = " ".join(change_items or []).strip().lower()

    target_object = ""
    for obj, aliases in OBJECT_ALIASES.items():
        if any(a in text for a in aliases):
            target_object = obj
            break

    action = "MODIFY"
    if any(k in text for k in ["추가", "add", "신규", "장착"]):
        action = "ADD"
    elif any(k in text for k in ["삭제", "제거", "remove"]):
        action = "REMOVE"
    elif any(k in text for k in ["변경", "교체", "modify", "replace"]):
        action = "MODIFY"

    main_change_keyword = ""
    for part, aliases in MAIN_PART_ALIASES.items():
        if any(a in text for a in aliases):
            main_change_keyword = part
            break

    return {
        "raw_text": " ".join(change_items or []),
        "target_object": target_object,
        "action": action,
        "main_change_keyword": main_change_keyword,
    }

def _doc_score_to_sim(d: dict) -> float:
    dist = d.get("distance", None)
    try:
        if dist is not None:
            dist = float(dist)
            return (1.0 - dist) if 0.0 < dist < 1.0 else 0.7
    except Exception:
        pass
    return 0.7

def run_search(change_items: list[str], target_model: str, dev_grade: str,
               top_k_primary: int = 60, top_k_secondary: int = 40) -> dict:
    """
    1) Legacy chroma (변경이력) - retrieve_docs() 우선
    2) Structured chroma (Base BOM) - 보조
    """
    intent = parse_change_intent(change_items)
    raw_text = intent["raw_text"].lower()

    # ✅ FEATURE 키워드만 추출 (OBJECT/ACTION 제외)
    feature_keywords = set()
    for alias_key, aliases in MAIN_PART_ALIASES.items():
        for a in aliases:
            if a in raw_text:
                feature_keywords.add(a)
                feature_keywords.add(alias_key.lower())
    if not feature_keywords:
        stop_words = {"에","를","을","의","에서","추가","삭제","변경","교체","add","remove","modify"}
        for w in re.split(r'[\s,]+', raw_text):
            w = w.strip()
            if len(w) >= 2 and w not in stop_words:
                feature_keywords.add(w)

    # ✅ 원자재/소모품 필터
    MATERIAL_PATTERNS = ["resin","sheet,steel","coil,steel","tape","paint","powder",
                         "sealant","grease","flux","solder","foam","film","수지","도료","테이프","그리스"]
    def is_material_part(desc):
        d = (desc or "").lower()
        return any(p in d for p in MATERIAL_PATTERNS)

    def is_direct_related(desc, rsn):
        d = (desc or "").lower()
        r = (rsn or "").lower()
        text = f"{d} {r}"
        return any(kw in text for kw in feature_keywords)

    query_lines = []
    if intent["target_object"]:
        query_lines.append(f"OBJECT: {intent['target_object']}")
    if intent["main_change_keyword"]:
        query_lines.append(f"FEATURE: {intent['main_change_keyword']}")
    query_lines.append(f"ACTION: {intent['action']}")
    query_lines.append(f"CHANGE: {intent['raw_text']}")
    query = "\n".join(query_lines)

    t0 = time.time()
    seen_ids = set()
    primary_docs = []
    min_target_docs = 22
    search_debug = {
        "legacy_filtered": 0,
        "legacy_unfiltered_added": 0,
        "structured_filtered": 0,
        "structured_unfiltered_added": 0,
        "sim_threshold": 0.35,
        "pre_sim_docs": 0,
        "post_sim_docs": 0,
    }

    # 모델/등급 필터를 강하게 적용해 다른 라인업 문서 유입을 줄인다.
    legacy_filters = {}
    if str(target_model or "").strip():
        legacy_filters["model"] = str(target_model).strip()
    if str(dev_grade or "").strip():
        legacy_filters["dev_grade"] = str(dev_grade).strip()

    # ✅ 1) Legacy chroma 검색 (변경이력 - 핵심!)
    try:
        legacy_top_k = max(14, min(int(top_k_primary), 30))
        legacy_raw = retrieve_with_filters(intent["raw_text"], top_k=legacy_top_k, filters=legacy_filters)
        search_debug["legacy_filtered"] = len(legacy_raw or [])
        for d in (legacy_raw or []):
            doc_id = d.get("id", "")
            if doc_id in seen_ids:
                continue
            seen_ids.add(doc_id)
            primary_docs.append({
                "id": doc_id,
                "document": d.get("text") or d.get("document") or "",
                "metadata": d.get("meta") or d.get("metadata") or {},
                "distance": d.get("dist") or d.get("distance") or 0.5,
            })

        # 결과가 목표치보다 적으면 단계적으로 완화 조회를 수행한다.
        if len(primary_docs) < min_target_docs:
            # 1) 모델만 유지(등급 해제)
            model_only_filters = {"model": str(target_model).strip()} if str(target_model or "").strip() else None
            legacy_relaxed = retrieve_with_filters(
                intent["raw_text"],
                top_k=max(legacy_top_k, 36),
                filters=model_only_filters,
            ) if model_only_filters else []

            added = 0
            for d in (legacy_relaxed or []):
                doc_id = d.get("id", "")
                if doc_id in seen_ids:
                    continue
                seen_ids.add(doc_id)
                primary_docs.append({
                    "id": doc_id,
                    "document": d.get("text") or d.get("document") or "",
                    "metadata": d.get("meta") or d.get("metadata") or {},
                    "distance": d.get("dist") or d.get("distance") or 0.5,
                })
                added += 1

            # 2) 그래도 부족하면 완전 해제 조회
            if len(primary_docs) < min_target_docs:
                legacy_relaxed2 = retrieve_docs(intent["raw_text"], top_k=max(legacy_top_k, 44))
                for d in (legacy_relaxed2 or []):
                    doc_id = d.get("id", "")
                    if doc_id in seen_ids:
                        continue
                    seen_ids.add(doc_id)
                    primary_docs.append({
                        "id": doc_id,
                        "document": d.get("text") or d.get("document") or "",
                        "metadata": d.get("meta") or d.get("metadata") or {},
                        "distance": d.get("dist") or d.get("distance") or 0.5,
                    })
                    added += 1

            search_debug["legacy_unfiltered_added"] = added
    except Exception as e:
        st.session_state.setdefault("debug_err", {})
        st.session_state["debug_err"]["legacy_error"] = repr(e)

    # ✅ 2) Structured chroma 검색 (보조)
    try:
        struct_where = {}
        if str(target_model or "").strip():
            struct_where["model"] = str(target_model).strip()
        if str(dev_grade or "").strip():
            struct_where["dev_grade"] = str(dev_grade).strip()
        structured_top_k = max(20, min(int(top_k_primary), 60))
        structured_docs = query_structured_docs(
            query,
            top_k=structured_top_k,
            where=to_chroma_where(struct_where) if struct_where else None,
        )
        search_debug["structured_filtered"] = len(structured_docs or [])
        for d in (structured_docs or []):
            doc_id = d.get("id", "")
            if doc_id in seen_ids:
                continue
            seen_ids.add(doc_id)
            primary_docs.append(d)

        # structured도 너무 적으면 where 없이 보강 조회
        if len(primary_docs) < min_target_docs:
            model_only_where = {"model": str(target_model).strip()} if str(target_model or "").strip() else {}
            relaxed_structured = query_structured_docs(
                query,
                top_k=max(24, structured_top_k),
                where=to_chroma_where(model_only_where) if model_only_where else None,
            )
            added = 0
            for d in (relaxed_structured or []):
                doc_id = d.get("id", "")
                if doc_id in seen_ids:
                    continue
                seen_ids.add(doc_id)
                primary_docs.append(d)
                added += 1

            if len(primary_docs) < min_target_docs:
                relaxed_structured2 = query_structured_docs(query, top_k=max(28, structured_top_k), where=None)
                for d in (relaxed_structured2 or []):
                    doc_id = d.get("id", "")
                    if doc_id in seen_ids:
                        continue
                    seen_ids.add(doc_id)
                    primary_docs.append(d)
                    added += 1
            search_debug["structured_unfiltered_added"] = added
    except Exception as e:
        st.session_state.setdefault("debug_err", {})
        st.session_state["debug_err"]["structured_error"] = repr(e)

    secondary_docs = []
    if is_reference_trigger(change_items):
        try:
            struct_where = {}
            if str(target_model or "").strip():
                struct_where["model"] = str(target_model).strip()
            if str(dev_grade or "").strip():
                struct_where["dev_grade"] = str(dev_grade).strip()
            secondary_docs = query_structured_docs(
                query + "\nREFERENCE",
                top_k=max(14, min(int(top_k_secondary), 40)),
                where=to_chroma_where(struct_where) if struct_where else None,
            )
        except Exception as e:
            st.session_state.setdefault("debug_err", {})
            st.session_state["debug_err"]["secondary"] = repr(e)

    # 거리 점수가 매우 낮은 문서는 파싱 전에 제거한다.
    # (legacy/structured 모두 공통 적용)
    def _keep_doc(d: dict, threshold: float) -> bool:
        sim = _doc_score_to_sim(d)
        return sim >= threshold

    search_debug["pre_sim_docs"] = len(primary_docs or [])
    sim_threshold = 0.30
    primary_docs_kept = [d for d in (primary_docs or []) if _keep_doc(d, sim_threshold)]
    # 결과가 너무 적으면 임계치를 완화해 recall을 보강한다.
    if len(primary_docs_kept) < min_target_docs:
        sim_threshold = 0.18
        primary_docs_kept = [d for d in (primary_docs or []) if _keep_doc(d, sim_threshold)]
    if len(primary_docs_kept) < max(10, min_target_docs // 2):
        sim_threshold = 0.12
        primary_docs_kept = [d for d in (primary_docs or []) if _keep_doc(d, sim_threshold)]
    # 그래도 너무 적으면 상위 유사도 순으로 최소 개수를 보장한다.
    if len(primary_docs_kept) < max(4, min_target_docs // 2):
        ranked = sorted((primary_docs or []), key=lambda x: _doc_score_to_sim(x), reverse=True)
        primary_docs_kept = ranked[:max(min_target_docs, len(primary_docs_kept))]

    primary_docs = primary_docs_kept
    secondary_docs = [d for d in (secondary_docs or []) if _keep_doc(d, max(0.2, sim_threshold - 0.05))]
    search_debug["sim_threshold"] = sim_threshold
    search_debug["post_sim_docs"] = len(primary_docs or [])

    hint_text = _project_region_hint_text()
    selected_region = _normalize_region_value(st.session_state.get("region"), hint_text)
    primary_docs, dropped_primary = _filter_docs_by_region(primary_docs, selected_region, hint_text)
    secondary_docs, dropped_secondary = _filter_docs_by_region(secondary_docs, selected_region, hint_text)
    st.session_state["region_filter_last"] = {
        "region": selected_region,
        "dropped_primary": dropped_primary,
        "dropped_secondary": dropped_secondary,
    }

    return {
        "policy_runtime": {"elapsed_sec": round(time.time() - t0, 3)},
        "primary_docs": primary_docs,
        "secondary_docs": secondary_docs,
        "query": query,
        "search_debug": search_debug,
    }

def evidence_from_doc(doc: dict, query: str = "") -> dict:
    """
    metadata['rag_json'] 우선 사용
    """
    meta = doc.get("metadata") or {}
    rag_json = meta.get("rag_json", "")
    rag_obj = None

    if rag_json:
        try:
            rag_obj = json.loads(rag_json)
        except Exception:
            rag_obj = None

    if rag_obj is None:
        txt = (doc.get("document") or doc.get("text") or doc.get("embedding_text") or "")
        rag_obj = unwrap_rag_json(txt)

    if rag_obj is None:
        return {
            "_invalid": True,
            "meta": {"doc_id": doc.get("id", "")},
            "scores": {"sim": 0.0},
        }

    return {
        "meta": {"doc_id": rag_obj.get("doc_id", "")},
        "object": {
            "main_object": (rag_obj.get("change") or {}).get("target_object", ""),
            "action": (rag_obj.get("change") or {}).get("action", ""),
        },
        "bom": rag_obj.get("bom") or {},
        "payload": {
            "feature": (rag_obj.get("change") or {}).get("feature", ""),
            "main_parts": (rag_obj.get("parts") or {}).get("main") or [],
            "sub_parts": (rag_obj.get("parts") or {}).get("sub") or [],
        },
        "reason": rag_obj.get("reason") or [],
        "review_points": rag_obj.get("review_points") or [],
        "scores": {"sim": _doc_score_to_sim(doc)},
        "_raw": rag_obj,
    }

def _vote_best_path(evs: list[dict]) -> tuple[str, float]:
    c = Counter()
    for ev in evs:
        p = (ev.get("bom") or {}).get("apply_bom_path", "") or ""
        if p:
            c[p] += 1
    if not c:
        return "", 0.0
    best, cnt = c.most_common(1)[0]
    return best, cnt / (sum(c.values()) or 1)

def _is_main_part_match(ev: dict, intent: dict) -> bool:
    kw = intent.get("main_change_keyword", "")
    if not kw:
        return False

    payload = ev.get("payload") or {}
    feat = (payload.get("feature") or "").upper()
    if feat == kw:
        return True

    for p in payload.get("main_parts") or []:
        pn = (p.get("part_name") or "").upper()
        if kw in pn:
            return True

    return False

def _vote_main_part(evs: list[dict]) -> dict | None:
    counter = Counter()
    sample = {}
    for ev in evs:
        for p in (ev.get("payload") or {}).get("main_parts") or []:
            name = (p.get("part_name") or "").strip()
            if name:
                counter[name] += 1
                sample.setdefault(name, p)
    if not counter:
        return None

    best_name, _ = counter.most_common(1)[0]
    return sample[best_name]

def _collect_sub_parts(all_evs: list[dict], intent: dict, selected_main: dict | None, selected_path: str) -> list[dict]:
    """
    같은 객체/경로 근처의 연관 부품을 서브 변경품 후보로
    """
    counter = Counter()
    sample = {}

    target_object = intent.get("target_object", "")
    selected_main_name = (selected_main or {}).get("part_name", "")
    selected_prefix = selected_path.split(">")[0].strip() if selected_path else target_object

    for ev in all_evs:
        obj = (ev.get("object") or {}).get("main_object", "")
        path = (ev.get("bom") or {}).get("apply_bom_path", "") or ""
        level1 = (ev.get("bom") or {}).get("level1", "") or ""

        if target_object and obj != target_object:
            continue

        if selected_prefix and not (path.startswith(selected_prefix) or level1 == selected_prefix):
            continue

        for p in (ev.get("payload") or {}).get("main_parts") or []:
            name = (p.get("part_name") or "").strip()
            if not name:
                continue
            if selected_main_name and name == selected_main_name:
                continue
            counter[name] += 1
            sample.setdefault(name, p)

    out = []
    for name, _ in counter.most_common(5):
        p = sample[name]
        out.append({
            "part_name": p.get("part_name", ""),
            "part_no": p.get("part_no", ""),
            "qty": p.get("qty", 1),
            "desc": p.get("desc", ""),
        })
    return out


def _parse_legacy_doc(doc, intent):
    txt = doc.get('document') or doc.get('text') or ''
    kw = (intent.get('main_change_keyword') or '').lower()
    obj = (intent.get('target_object') or '').lower()
    raw_text = (intent.get('raw_text') or '').lower()

    kw_terms = set()
    if kw:
        kw_terms.add(kw)
        for a in (MAIN_PART_ALIASES.get(kw.upper()) or []):
            aa = str(a or '').strip().lower()
            if aa:
                kw_terms.add(aa)

    # 변경점 원문의 모든 유효 토큰을 keyword로 포함 (별칭 테이블 의존 제거)
    _kw_stop = {
        "에","를","을","의","에서","로","으로","및","하여","관련","대한",
        "추가","삭제","변경","교체","적용","개선","신규","기능","개발","사양",
        "add","remove","modify","change","for","the","and","with","new",
    }
    for _tok in re.findall(r"[a-z0-9]+(?:[.\-][a-z0-9]+)*|[가-힣]{2,}", raw_text):
        _tok = _tok.strip()
        if len(_tok) >= 2 and _tok not in _kw_stop:
            kw_terms.add(_tok)

    def _norm_alias_text(s: str) -> str:
        return re.sub(r"[^a-z0-9가-힣]+", "", str(s or "").lower())

    results = []
    for line in txt.splitlines():
        line = line.strip()
        if not line.startswith('-'):
            continue
        ll = line.lower()
        ll_norm = _norm_alias_text(ll)

        has_kw = any(_norm_alias_text(t) in ll_norm for t in kw_terms) if kw_terms else False
        has_obj = bool(obj and obj in ll)
        if not (has_kw or has_obj):
            continue

        new_pno = ''
        base_pno = ''
        m = re.search(r'New=([A-Z0-9\-]{5,20})', line, re.I)
        if m:
            new_pno = re.sub(r'[^A-Z0-9]', '', m.group(1).upper())
        mb = re.search(r'Base=([A-Z0-9\-]{5,20})', line, re.I)
        if mb:
            base_pno = re.sub(r'[^A-Z0-9]', '', mb.group(1).upper())

        desc = ''
        m_desc = re.search(r'Desc=(.+?)(?:\||$)', line, re.I)
        if m_desc:
            desc = m_desc.group(1).strip()
        else:
            pts = [p.strip() for p in line.split('|')]
            for p in pts:
                if not p or p.startswith('-'):
                    continue
                if re.match(r'^(Base|New|RSN|CHG|row|L\d)=', p, re.I):
                    continue
                desc = p
                break

        bom_path = ''
        pm = re.search(r'L\d+:([A-Z0-9]+)\s*>\s*L\d+:([A-Z0-9]+)', line, re.I)
        if pm:
            bom_path = pm.group(1) + ' > ' + pm.group(2)
        else:
            pm2 = re.search(r'L\d+:([A-Z0-9]+)', line, re.I)
            if pm2:
                bom_path = pm2.group(1)
        part_no = new_pno or base_pno
        if not part_no:
            continue
        if 'nan' in bom_path.lower():
            continue
        results.append({'part_name': desc, 'part_no': part_no, 'bom_path': bom_path})
    return results


    def reorder_parts_by_base_skeleton(parts, base_df, l1_pno="", l1_desc=""):
        """
        Base BOM skeleton의 실제 DFS 순서를 기준으로 flat 추천 부품을 재정렬한다.
        반환: (matched_parts, unmatched_parts)
        """
        def _clean(v):
            s = "" if v is None else str(v).strip()
            return "" if s.lower() in ("", "nan", "none") else s

        def _norm_pno(v):
            return re.sub(r"[^A-Z0-9]", "", _clean(v).upper())

        def _norm_desc(v):
            return re.sub(r"\s+", " ", _clean(v).upper())

        def _lvl_depth(v):
            s = _clean(v)
            if not s:
                return None
            if s.startswith("."):
                return s.count(".")
            try:
                return int(s)
            except Exception:
                return None

        if base_df is None or not isinstance(base_df, pd.DataFrame) or len(base_df) == 0:
            return list(parts or []), []

        pno_col = pick_col(base_df, ["P/NO", "P/NO.", "Part No", "품번"])
        desc_col = pick_col(base_df, ["Description", "DESC", "부품명"])
        lvl_col = pick_col(base_df, ["Lvl", "LVL", "Level", "레벨"])

        if not (lvl_col and (pno_col or desc_col)):
            return list(parts or []), []

        target_pno = _norm_pno(l1_pno)
        target_desc = _norm_desc(l1_desc)

        # 1) L1 시작 행 찾기 (.1 우선, 품번 → 부품명)
        start_idx = None
        if pno_col and target_pno:
            for i, row in base_df.iterrows():
                if _norm_pno(row.get(pno_col, "")) == target_pno and _clean(row.get(lvl_col, "")) == ".1":
                    start_idx = i
                    break

        if start_idx is None and desc_col and target_desc:
            for i, row in base_df.iterrows():
                lvl_v = _clean(row.get(lvl_col, ""))
                if lvl_v != ".1":
                    continue
                rd = _norm_desc(row.get(desc_col, ""))
                if target_desc and (target_desc in rd or rd in target_desc):
                    start_idx = i
                    break

        if start_idx is None:
            return list(parts or []), []

        base_depth = _lvl_depth(base_df.loc[start_idx, lvl_col])
        if base_depth is None:
            return list(parts or []), []

        subtree_rows = []
        for j in range(start_idx + 1, len(base_df)):
            row = base_df.iloc[j]
            depth = _lvl_depth(row.get(lvl_col, ""))
            if depth is None:
                continue
            if depth <= base_depth:
                break

            subtree_rows.append({
                "row_index": j,
                "pno_norm": _norm_pno(row.get(pno_col, "")) if pno_col else "",
                "desc_norm": _norm_desc(row.get(desc_col, "")) if desc_col else "",
                "depth": depth,
            })

        if not subtree_rows:
            return list(parts or []), []

        # subtree 내 순서를 고정 인덱스로 만든다.
        subtree_order_map = {}
        for order, node in enumerate(subtree_rows):
            if node["pno_norm"] and node["pno_norm"] not in subtree_order_map:
                subtree_order_map[node["pno_norm"]] = (order, node["depth"])
            if node["desc_norm"] and node["desc_norm"] not in subtree_order_map:
                subtree_order_map[node["desc_norm"]] = (order, node["depth"])

        matched_parts = []
        unmatched_parts = []

        for p in (parts or []):
            if not isinstance(p, dict):
                continue
            cp = dict(p)
            pnos = [
                cp.get("db_new_pno"), cp.get("db_base_pno"),
                cp.get("part_no"), cp.get("display_pno"),
            ]
            matched = None
            for pn in pnos:
                pn_norm = _norm_pno(pn)
                if pn_norm and pn_norm in subtree_order_map:
                    matched = subtree_order_map[pn_norm]
                    break

            if matched is None:
                desc_norm = _norm_desc(cp.get("desc") or cp.get("part_name") or "")
                if desc_norm and desc_norm in subtree_order_map:
                    matched = subtree_order_map[desc_norm]

            if matched is None:
                unmatched_parts.append(cp)
                continue

            order, depth = matched
            cp["skeleton_order"] = order
            if not _clean(cp.get("db_lvl")):
                cp["db_lvl"] = "." * depth + str(depth)
            if not _clean(cp.get("lvl")):
                cp["lvl"] = cp["db_lvl"]
            matched_parts.append(cp)

        matched_parts.sort(key=lambda x: x.get("skeleton_order", 999999))
        return matched_parts, unmatched_parts

def reorder_parts_by_base_skeleton(parts, base_df, l1_pno="", l1_desc=""):
    """Reorder flat parts by Base BOM subtree DFS order."""
    def _clean(v):
        s = "" if v is None else str(v).strip()
        return "" if s.lower() in ("", "nan", "none") else s

    def _pno(v):
        return re.sub(r"[^A-Z0-9]", "", _clean(v).upper())

    def _desc(v):
        return re.sub(r"\s+", " ", _clean(v).upper())

    def _depth(v):
        s = _clean(v)
        if not s:
            return None
        if s.startswith("."):
            return s.count(".")
        try:
            return int(s)
        except Exception:
            return None

    if base_df is None or not isinstance(base_df, pd.DataFrame) or len(base_df) == 0:
        return list(parts or []), []

    pno_col = pick_col(base_df, ["P/NO", "P/NO.", "Part No", "품번"])
    desc_col = pick_col(base_df, ["Description", "DESC", "부품명"])
    lvl_col = pick_col(base_df, ["Lvl", "LVL", "Level", "레벨"])
    if not (lvl_col and (pno_col or desc_col)):
        return list(parts or []), []

    target_pno = _pno(l1_pno)
    target_desc = _desc(l1_desc)
    start_idx = None
    if pno_col and target_pno:
        for i, row in base_df.iterrows():
            if _pno(row.get(pno_col, "")) == target_pno and _clean(row.get(lvl_col, "")) == ".1":
                start_idx = i
                break
    if start_idx is None and desc_col and target_desc:
        for i, row in base_df.iterrows():
            if _clean(row.get(lvl_col, "")) != ".1":
                continue
            rd = _desc(row.get(desc_col, ""))
            if target_desc in rd or rd in target_desc:
                start_idx = i
                break
    if start_idx is None:
        return list(parts or []), []

    base_depth = _depth(base_df.loc[start_idx, lvl_col])
    if base_depth is None:
        return list(parts or []), []

    subtree_rows = []
    for j in range(start_idx + 1, len(base_df)):
        row = base_df.iloc[j]
        d = _depth(row.get(lvl_col, ""))
        if d is None:
            continue
        if d <= base_depth:
            break
        subtree_rows.append({
            "order": len(subtree_rows),
            "pno": _pno(row.get(pno_col, "")) if pno_col else "",
            "desc": _desc(row.get(desc_col, "")) if desc_col else "",
            "depth": d,
        })
    if not subtree_rows:
        return list(parts or []), []

    order_map = {}
    for n in subtree_rows:
        if n["pno"] and n["pno"] not in order_map:
            order_map[n["pno"]] = (n["order"], n["depth"])
        if n["desc"] and n["desc"] not in order_map:
            order_map[n["desc"]] = (n["order"], n["depth"])

    matched_parts, unmatched_parts = [], []
    for p in (parts or []):
        if not isinstance(p, dict):
            continue
        cp = dict(p)
        matched = None
        for k in ("db_new_pno", "db_base_pno", "part_no", "display_pno"):
            pn = _pno(cp.get(k))
            if pn and pn in order_map:
                matched = order_map[pn]
                break
        if matched is None:
            dn = _desc(cp.get("desc") or cp.get("part_name") or "")
            if dn and dn in order_map:
                matched = order_map[dn]
        if matched is None:
            unmatched_parts.append(cp)
            continue
        sk_order, depth = matched
        cp["skeleton_order"] = sk_order
        if not _clean(cp.get("db_lvl")):
            cp["db_lvl"] = "." * depth + str(depth)
        if not _clean(cp.get("lvl")):
            cp["lvl"] = cp["db_lvl"]
        matched_parts.append(cp)

    matched_parts.sort(key=lambda x: x.get("skeleton_order", 999999))
    return matched_parts, unmatched_parts


def post_retrieval_dedup(rows):
    """
    Step 2.5: Post-Retrieval Dedup
    - 1순위: 실제 품번 동일 시 동일 부품
    - 2순위: 품번이 없으면 (부품명 + 부품유형) 동일 시 동일 부품
    - 레벨/사유/출처는 판단 기준에서 제외하고 메타 리스트로 보존
    """
    def _clean_text(v):
        s = "" if v is None else str(v).strip()
        return "" if s.lower() in ("", "nan", "none") else s

    def _first(d, keys):
        for k in keys:
            if k in d:
                vv = _clean_text(d.get(k))
                if vv:
                    return vv
        return ""

    def _extract_real_pno(d):
        for k in ("part_no", "display_pno", "db_new_pno", "db_base_pno", "품번", "Part No", "P/NO", "P/NO."):
            raw = _clean_text(d.get(k))
            if not raw:
                continue
            if any(tok in _norm_merge_text(raw) for tok in ("채번", "TBD", "NEED", "미정")):
                continue
            pn = _norm_merge_pno(raw)
            if pn:
                return pn
        return ""

    def _norm_name_type_key(d):
        name = _norm_merge_text(
            _first(d, ["part_name", "desc", "부품명", "Description", "Part Name(자)", "Part Name"])
        )
        ptype = _norm_merge_text(
            _first(d, ["base_type", "type", "유형", "Part Type", "part_type"])
        )
        return (name, ptype)

    def _action_rank(d):
        amap = {"MODIFY": 0, "ADD": 1, "CHECK": 2}
        raw = _norm_merge_text(_first(d, ["action", "변경유형"]))
        if raw in ("변경",):
            raw = "MODIFY"
        elif raw in ("추가",):
            raw = "ADD"
        elif raw in ("⚠️확인필요", "확인필요"):
            raw = "CHECK"
        return amap.get(raw, 99)

    def _lvl_depth(d):
        lv = _first(d, ["lvl", "db_lvl", "레벨", "Lvl", "Level"])
        if not lv:
            return 9999
        return str(lv).count(".")

    def _dedup_values(items):
        out = []
        seen = set()
        for v in items:
            vv = _clean_text(v)
            if not vv:
                continue
            kk = _norm_merge_text(vv)
            if kk in seen:
                continue
            seen.add(kk)
            out.append(vv)
        return out

    normalized_rows = []
    for r in (rows or []):
        if r is None:
            continue
        if isinstance(r, dict):
            normalized_rows.append(dict(r))
            continue
        to_dict = getattr(r, "to_dict", None)
        if callable(to_dict):
            normalized_rows.append(dict(to_dict()))

    before_cnt = len(normalized_rows)
    if before_cnt == 0:
        st.session_state["_post_retrieval_dedup_stats"] = {"before": 0, "after": 0, "merged": 0}
        return []

    grouped = {}
    order_keys = []
    for idx, row in enumerate(normalized_rows):
        pno = _extract_real_pno(row)
        if pno:
            dkey = ("PNO", pno)
        else:
            name_key, type_key = _norm_name_type_key(row)
            dkey = ("NAME_TYPE", name_key, type_key)
        if dkey not in grouped:
            grouped[dkey] = []
            order_keys.append(dkey)
        grouped[dkey].append(row)

    merged_rows = []
    for gk in order_keys:
        items = grouped[gk]

        rep = sorted(
            items,
            key=lambda x: (_action_rank(x), _lvl_depth(x), -sum(1 for k in x.keys() if _clean_text(x.get(k)))),
        )[0]

        action_vals = _dedup_values([_first(it, ["action", "변경유형"]) for it in items])
        lvl_vals = _dedup_values([_first(it, ["lvl", "db_lvl", "레벨", "Lvl", "Level"]) for it in items])
        source_vals = _dedup_values([_first(it, ["source_doc", "출처", "source", "file_name"]) for it in items])
        l1_vals = _dedup_values([_first(it, ["l1_desc", "L1", "level1", "level1_desc"]) for it in items])
        rsn_vals = _dedup_values([_first(it, ["rsn", "chg", "변경사유", "reason"]) for it in items])

        out = dict(rep)

        # action 우선순위 대표값 반영
        out["action"] = _first(rep, ["action", "변경유형"])

        # 레벨은 가장 얕은 값 선택
        shallow = sorted(items, key=lambda x: (_lvl_depth(x), _first(x, ["lvl", "db_lvl", "레벨", "Lvl", "Level"])))[0]
        best_lvl = _first(shallow, ["lvl", "db_lvl", "레벨", "Lvl", "Level"])
        if best_lvl:
            out["lvl"] = best_lvl
            out["db_lvl"] = best_lvl

        # 동일 키 병합 메타 보존
        merged_meta = []
        seen_meta = set()
        for it in items:
            meta_item = {
                "source_doc": _first(it, ["source_doc", "출처", "source", "file_name"]),
                "l1_desc": _first(it, ["l1_desc", "L1", "level1", "level1_desc"]),
                "reason": _first(it, ["rsn", "chg", "변경사유", "reason"]),
                "action": _first(it, ["action", "변경유형"]),
                "lvl": _first(it, ["lvl", "db_lvl", "레벨", "Lvl", "Level"]),
            }
            mk = (
                _norm_merge_text(meta_item["source_doc"]),
                _norm_merge_text(meta_item["l1_desc"]),
                _norm_merge_text(meta_item["reason"]),
                _norm_merge_text(meta_item["action"]),
                _norm_merge_text(meta_item["lvl"]),
            )
            if mk in seen_meta:
                continue
            seen_meta.add(mk)
            merged_meta.append(meta_item)

        out["_merged_sources"] = merged_meta

        # 기존 호환용 메타 유지
        out["meta_part_nos"] = _dedup_values([_extract_real_pno(it) for it in items])
        out["meta_levels"] = lvl_vals
        out["meta_actions"] = action_vals
        out["meta_reasons"] = rsn_vals
        out["meta_sources"] = source_vals
        out["meta_l1_desc"] = l1_vals
        out["meta_merged_count"] = len(items)

        if not _clean_text(out.get("source_doc")) and source_vals:
            out["source_doc"] = source_vals[0]
        if not _clean_text(out.get("rsn")) and rsn_vals:
            out["rsn"] = rsn_vals[0]

        out["rsn_key"] = _norm_merge_text(out.get("rsn") or out.get("chg"))
        merged_rows.append(out)

    after_cnt = len(merged_rows)
    st.session_state["_post_retrieval_dedup_stats"] = {
        "before": before_cnt,
        "after": after_cnt,
        "merged": before_cnt - after_cnt,
    }
    return merged_rows

def generate_proposals_from_docs(primary_docs, base_snapshot, change_items):
    """
    ✅ v4: RSN 기반 필터링 + CORE/CASCADE 분류 + 출처 태깅

    Step 0: 의도 파싱 → feature 키워드(카메라) + target 키워드(도어)
    Step 1: DB 문서에서 부품 파싱 + source_doc 태깅
    Step 2: RSN 기반 필터 (CORE / CASCADE / EXCLUDE)
    Step 3: base BOM 대조 (품번 우선, in_base 확인 + lvl/qty 보충)
    Step 4: L1 대표행 + 중복 제거
    Step 5: proposal dict 생성
    """
    intent = parse_change_intent(change_items)
    raw_text = intent["raw_text"].lower()

    # ═══════════════════════════════════════════════════
    # Step 0: feature / target 키워드 추출
    #   "도어에 카메라 추가"
    #   → feature_keywords = {"카메라","camera","CAMERA"}
    #   → target_keywords  = {"도어","door","DOOR"}
    # ═══════════════════════════════════════════════════
    feature_keywords = set()
    target_keywords = set()

    def _extract_core_keywords(text: str) -> set[str]:
        # Keep noun-like Korean/English/number tokens for loose RSN matching.
        stop_words = {
            "에", "를", "을", "의", "에서", "로", "으로", "및", "또는", "그리고",
            "추가", "삭제", "변경", "교체", "적용", "개선", "관련", "대응", "검토",
            "기능", "신규", "디자인", "개발", "사양", "spec",
            "add", "remove", "modify", "change", "for", "and", "the", "with",
            "new", "design", "develop", "development", "function", "feature", "spec",
        }
        toks = set()
        for w in re.findall(r"[A-Za-z0-9]+(?:[._\-][A-Za-z0-9]+)?|[가-힣]{2,}", str(text or "").lower()):
            ww = w.strip()
            if not ww or ww in stop_words:
                continue
            toks.add(ww)
        return toks

    for alias_key, aliases in MAIN_PART_ALIASES.items():
        for a in aliases:
            if a in raw_text:
                feature_keywords.add(a)
                feature_keywords.add(alias_key.lower())

    for alias_key, aliases in OBJECT_ALIASES.items():
        for a in aliases:
            if a in raw_text:
                target_keywords.add(a)
                target_keywords.add(alias_key.lower())

    # FEATURE가 명시되면 해당 alias만 핵심 키워드로 사용한다.
    # (일반 단어 확장으로 인한 오탐 방지)
    core_keywords = set(k.lower() for k in feature_keywords)
    if not core_keywords:
        core_keywords.update(_extract_core_keywords(intent.get("raw_text") or ""))

    # 흔한 일반어는 CORE 판정 키워드에서 제외
    generic_noise = {
        "추가", "변경", "신규", "디자인", "개발", "기능", "사양",
        "add", "modify", "new", "design", "develop", "function", "feature", "spec",
    }
    core_keywords = {k for k in core_keywords if k and k not in generic_noise and len(k) >= 2}
    target_keywords = set(k.lower() for k in target_keywords)

    def _kw_match(text: str, kw: str) -> bool:
        txt = str(text or "").lower()
        k = str(kw or "").lower().strip()
        if not txt or not k:
            return False
        # 영문/숫자 키워드는 단어 경계로 엄격 매칭
        if re.fullmatch(r"[a-z0-9_\-]+", k):
            return re.search(rf"(?<![a-z0-9]){re.escape(k)}(?![a-z0-9])", txt) is not None
        # 한글/혼합 키워드는 부분 포함 매칭
        return k in txt

    # ── 원자재/소모품 필터 (이건 무조건 제외) ──
    MATERIAL_PATTERNS = [
        "resin","sheet,steel","coil,steel","tape","paint","powder",
        "sealant","grease","flux","solder","foam","film",
        "수지","도료","테이프","그리스",
    ]
    def is_material_part(desc):
        d = (desc or "").lower()
        return any(p in d for p in MATERIAL_PATTERNS)

    # ═══════════════════════════════════════════════════
    # RSN 기반 분류 함수
    #   CORE:    RSN에 feature 키워드 포함 ("카메라 모듈 개발")
    #   CASCADE: RSN에 target 키워드 포함 ("도어 무게 변경으로 인한~")
    #   EXCLUDE: RSN 없음/nan/관련 키워드 없음 ("신규 디자인 핸들")
    # ═══════════════════════════════════════════════════
    def classify_by_rsn(rsn, desc="", chg=""):
        r = (rsn or chg or "").strip().lower()
        d = (desc or "").lower()
        c = (chg or "").lower()

        # D-012: RSN/CHG가 비어있어도 desc에 변경점 키워드가 매치되면
        # CORE/CASCADE로 인정. 원본은 RSN 없으면 무조건 EXCLUDE라 BOM
        # 부품/마스터처럼 변경 이력 컬럼이 없는 doc은 모두 떨어졌음.

        # 1) CORE: rsn / desc / chg 중 어디든 core 키워드 매치
        if any(_kw_match(r, kw) for kw in core_keywords if r):
            return "CORE"
        if any(_kw_match(d, kw) for kw in core_keywords):
            return "CORE"
        if any(_kw_match(c, kw) for kw in core_keywords):
            return "CORE"

        # 2) CASCADE: target 키워드 매치
        if any(_kw_match(r, kw) for kw in target_keywords if r):
            return "CASCADE"
        if any(_kw_match(d, kw) for kw in target_keywords):
            return "CASCADE"
        if any(_kw_match(c, kw) for kw in target_keywords):
            return "CASCADE"

        # 3) RSN/CHG 모두 비어있는 doc도 desc에 query_tokens 매치되면 CASCADE
        # (변경 이력 없는 BOM 부품을 사용자 변경점 키워드로 회수)
        if (not r) and any(_kw_match(d, kw) for kw in query_tokens):
            return "CASCADE"

        return "EXCLUDE"

    # 변경점 텍스트와의 직접 연관성 검증 (정밀도 우선)
    query_tokens = _extract_core_keywords(intent.get("raw_text") or "")

    def _is_related_to_change(desc: str, rsn: str, chg: str = "") -> bool:
        blob = f"{str(desc or '').lower()} {str(rsn or '').lower()} {str(chg or '').lower()}"

        # feature가 명확하면 우선 feature 일치를 본다.
        # 단, 미일치라고 즉시 탈락시키지 않고 target/query 토큰 fallback을 허용한다.
        if core_keywords:
            if any(_kw_match(blob, kw) for kw in core_keywords):
                return True

        # target/질의토큰과 최소 1개 이상 겹치면 연관으로 본다.
        if any(_kw_match(blob, kw) for kw in target_keywords):
            return True
        if any(_kw_match(blob, kw) for kw in query_tokens):
            return True
        return False

    # ═══════════════════════════════════════════════════
    # Step 1: DB 문서에서 부품 파싱
    #   기존과 동일한 파싱 + source_doc/tier 태깅 추가
    # ═══════════════════════════════════════════════════
    db_parts = []

    for doc in (primary_docs or []):
        txt = doc.get("document") or doc.get("text") or ""
        if not txt or "[L1]" not in txt:
            continue

        doc_meta = doc.get("metadata") or {}

        source_doc = (doc_meta.get("source_file")
                      or doc_meta.get("source")
                      or doc_meta.get("src")
                      or doc_meta.get("file_name")
                      or doc_meta.get("id")
                      or "")

        current_l1_desc = ""
        for line in txt.splitlines():
            line = line.strip()

            l1_match = re.match(
                r'\[L1\]\s*[A-Z0-9]+\s*\|\s*Desc=(.+)', line, re.I
            )
            if l1_match:
                current_l1_desc = l1_match.group(1).strip()
                continue

            if not line.startswith("-"):
                continue

            fields = [f.strip() for f in line.split("|")]
            if len(fields) < 3:
                continue

            # ── 수정: 태그 기반 파싱 ──
            full_line = "|".join(fields)

            # 품번: 전체 라인에서 검색
            m_base = re.search(r'Base=([A-Z0-9]+)', full_line, re.I)
            m_new  = re.search(r'New=([A-Z0-9]+)', full_line, re.I)
            base_pno = m_base.group(1) if m_base else ""
            new_pno  = m_new.group(1) if m_new else ""

            # DB 라인의 경로(L1/L2/...)에서 레벨을 추출해 보존한다.
            db_lvl = ""
            lvl_tokens = re.findall(r'L\s*(\d+)\s*:', fields[0] if fields else "", re.I)
            if lvl_tokens:
                try:
                    lv_num = int(lvl_tokens[-1])
                    if lv_num > 0:
                        db_lvl = "." * lv_num + str(lv_num)
                except Exception:
                    db_lvl = ""

            # desc: Desc= 태그 우선, 없으면 3번째 필드 fallback
            desc = ""
            m_desc = re.search(r'Desc=(.+?)(?:\||$)', full_line, re.I)
            if m_desc:
                desc = m_desc.group(1).strip()
            else:
                # fallback: 태그 아닌 필드 중 첫 번째
                for f in fields[1:]:
                    f = f.strip()
                    if not re.match(r'^(Base|New|RSN|CHG|row|L\d)=', f, re.I) and f:
                        desc = f
                        break

            if not desc:
                continue

            # RSN, CHG: 전체 라인에서 검색
            rsn = ""
            chg = ""
            m_rsn = re.search(r'RSN=(.+?)(?:\||$)', full_line, re.I)
            m_chg = re.search(r'CHG=(.+?)(?:\||$)', full_line, re.I)
            if m_rsn:
                rsn = m_rsn.group(1).strip()
            if m_chg:
                chg = m_chg.group(1).strip()

            # action 판정 (DB 기록 기반)
            if not base_pno and new_pno:
                action = "ADD"
            elif base_pno and new_pno and base_pno != new_pno:
                action = "MODIFY"
            elif base_pno and new_pno and base_pno == new_pno:
                # Same part-no changes (spec/reason updates) should still be considered
                # when they are directly related to the user change intent.
                if _is_related_to_change(desc, rsn, chg):
                    action = "MODIFY"
                else:
                    action = "KEEP"
            elif base_pno and not new_pno:
                action = "DELETE"
            else:
                action = "KEEP"

            if action == "KEEP":
                continue
            if is_material_part(desc):
                continue

            # 1차 판정: KEEP/자재 제외 후 RSN 분류를 기록 (EXCLUDE도 일단 보관)
            tier = classify_by_rsn(rsn, desc, chg)
            if tier != "EXCLUDE" and not _is_related_to_change(desc, rsn, chg):
                tier = "EXCLUDE"
            rsn_key = _norm_merge_text(rsn or chg)
            db_parts.append({
                "desc": desc, "action": action,
                "l1_desc": current_l1_desc,
                "db_base_pno": base_pno, "db_new_pno": new_pno,
                "db_lvl": db_lvl,
                "rsn": rsn, "chg": chg,
                "tier": tier,
                "rsn_key": rsn_key,
                "source_doc": source_doc,   # ✅ 신규
            })

    # 2차 판정: RSN 세트 복구
    passed_rsn_keys = {
        p.get("rsn_key", "")
        for p in db_parts
        if p.get("tier") in ("CORE", "CASCADE") and p.get("rsn_key")
    }
    for p in db_parts:
        if p.get("tier") == "EXCLUDE" and p.get("rsn_key") and p.get("rsn_key") in passed_rsn_keys:
            p["tier"] = "CASCADE"

    # CORE가 포함된 RSN 세트는 전부 CORE_GROUP으로 승격
    rsn_groups = defaultdict(list)
    for p in db_parts:
        rk = p.get("rsn_key") or ""
        if rk:
            rsn_groups[rk].append(p)

    for _, group in rsn_groups.items():
        has_core = any(g.get("tier") == "CORE" for g in group)
        if not has_core:
            continue
        for g in group:
            if g.get("tier") in ("CASCADE",):
                g["tier"] = "CORE_GROUP"

    # 3차 판정: 같은 source_doc + 같은 L1 블록 내에서
    # CORE/CORE_GROUP이 1건이라도 있으면, 같은 블록의 EXCLUDE를 CASCADE로 복구한다.
    # (문맥상 같은 변경 묶음인데 키워드 매칭이 약해 탈락한 경우 보강)
    doc_l1_groups = defaultdict(list)
    for p in db_parts:
        k = (
            _norm_merge_text(p.get("source_doc") or ""),
            _norm_merge_text(p.get("l1_desc") or ""),
        )
        doc_l1_groups[k].append(p)

    recovered_doc_l1 = 0
    for _, group in doc_l1_groups.items():
        has_core_like = any(g.get("tier") in ("CORE", "CORE_GROUP") for g in group)
        if not has_core_like:
            continue
        for g in group:
            if g.get("tier") == "EXCLUDE":
                rsn_txt = str(g.get("rsn") or "").strip().lower()
                if rsn_txt and rsn_txt != "nan":
                    g["tier"] = "CASCADE"
                    recovered_doc_l1 += 1

    # 4차 판정(적응형): 같은 source_doc + 같은 L1 내부에서만 제한적으로 복구
    # 전역 EXCLUDE 복구는 무관 부품 확산을 유발하므로 제거한다.
    kept_after_l1 = [p for p in db_parts if p.get("tier") != "EXCLUDE"]
    raw_cnt = len(db_parts)
    min_keep = max(16, int(raw_cnt * 0.40)) if raw_cnt > 0 else 0
    recovered_adaptive = 0

    if len(kept_after_l1) < min_keep:
        for _, group in doc_l1_groups.items():
            has_kept = any(g.get("tier") != "EXCLUDE" for g in group)
            if not has_kept:
                continue
            for g in group:
                if g.get("tier") != "EXCLUDE":
                    continue
                rsn_txt = str(g.get("rsn") or g.get("chg") or "").strip().lower()
                if rsn_txt and rsn_txt != "nan":
                    g["tier"] = "CASCADE"
                    recovered_adaptive += 1

    if recovered_doc_l1 or recovered_adaptive:
        st.session_state["_recover_doc_l1"] = recovered_doc_l1
        st.session_state["_recover_adaptive"] = recovered_adaptive

    # D-012 진단: tier 분포 + intent를 _dbg.log에 박음
    try:
        from collections import Counter as _DbgCounter
        from pathlib import Path as _DbgPath
        from datetime import datetime as _DbgDt
        _dbg_path = _DbgPath(__file__).resolve().parent / "_dbg.log"
        _tier_dist = _DbgCounter(p.get("tier", "?") for p in db_parts)
        with open(_dbg_path, "a", encoding="utf-8") as _df:
            _df.write(
                f"[{_DbgDt.now().isoformat(timespec='seconds')}] "
                f"  proposal_parser: db_parts_total={len(db_parts)} "
                f"tier_dist={dict(_tier_dist)} "
                f"core_kw={sorted(core_keywords)[:10]} "
                f"target_kw={sorted(target_keywords)[:10]} "
                f"query_tokens={sorted(query_tokens)[:10]}\n"
            )
            # 첫 3개 row 샘플
            for p in db_parts[:3]:
                _df.write(
                    f"    sample: tier={p.get('tier')} desc={p.get('desc','')[:50]!r} "
                    f"rsn={p.get('rsn','')[:50]!r} chg={p.get('chg','')[:50]!r}\n"
                )
    except Exception:
        pass

    db_parts = [p for p in db_parts if p.get("tier") != "EXCLUDE"]
    # ── 임시 디버그: 품번 없는 부품 확인 (확인 후 삭제!) ──
    _no_pno = [p for p in db_parts if not p.get("db_new_pno") and not p.get("db_base_pno")]
    if _no_pno:
        st.session_state["_debug_no_pno"] = [
            {"desc": p["desc"], "rsn": p.get("rsn",""), "tier": p.get("tier","")}
            for p in _no_pno
        ]

    # Step 2.5: Post-Retrieval Dedup
    # 같은 부품이 출처만 달라 중복된 경우 병합하고, 출처/사유는 메타 리스트로 보존한다.
    unique_parts = post_retrieval_dedup(db_parts)

    if not unique_parts:
        return []

    # ═══════════════════════════════════════════════════
    # Step 3: base BOM 대조
    #   목적: in_base 확인 + lvl/qty 보충
    #   매칭 우선순위: ① DB new_pno → ② DB base_pno → ③ desc fallback
    #   부품 정보(품번/desc)는 DB에서 온 그대로 유지!
    # ═══════════════════════════════════════════════════
    base_df = st.session_state.get("base_df")
    base_pno_index = {}    # 품번 key → base row
    base_desc_index = {}   # desc key → base row list
    desc_col = pno_col = lvl_col = parent_col = type_col = None

    if isinstance(base_df, pd.DataFrame) and len(base_df) > 0:
        desc_col = pick_col(base_df, ["Description","DESC","부품명","품명"])
        pno_col  = pick_col(base_df, ["Part No","P/NO","품번"])
        lvl_col  = pick_col(base_df, ["Lvl","LVL","Level","LEVEL","레벨"])
        parent_col = pick_col(base_df, ["Parent Part No(모)","Parent Part No"])
        type_col = pick_col(base_df, ["Type","TYPE"])

        if pno_col:
            for _, row in base_df.iterrows():
                pn = str(row.get(pno_col, "") or "").strip()
                d  = str(row.get(desc_col, "") or "").strip() if desc_col else ""
                info = {
                    "part_no": pn, "desc": d,
                    "lvl": str(row.get(lvl_col, "") or "").strip() if lvl_col else "",
                    "parent": str(row.get(parent_col, "") or "").strip() if parent_col else "",
                    "type": str(row.get(type_col, "") or "").strip() if type_col else "",
                    "qty": str(row.get("Qty","") or row.get("QTY","") or "1").strip(),
                }
                if pn:
                    pn_key = re.sub(r'[^A-Z0-9]','', pn.upper())
                    base_pno_index[pn_key] = info
                if d:
                    base_desc_index.setdefault(d.upper(), []).append(info)

    def map_to_base(p):
        """품번 우선 매칭 → in_base 확인 + lvl/qty 보충"""
        db_new  = p.get("db_new_pno") or ""
        db_base = p.get("db_base_pno") or ""
        new_key  = re.sub(r'[^A-Z0-9]','', db_new.upper()) if db_new else ""
        base_key = re.sub(r'[^A-Z0-9]','', db_base.upper()) if db_base else ""

        bm = None
        matched_pno = ""

        # ① new 품번으로 base 조회
        if new_key and new_key in base_pno_index:
            bm = base_pno_index[new_key]
            matched_pno = db_new
        # ② base 품번으로 조회
        elif base_key and base_key in base_pno_index:
            bm = base_pno_index[base_key]
            matched_pno = db_base
        # ③ desc fallback (정확 매칭만)

        if bm:
            # in_base + ADD → 확인필요로 변경
            action = p.get("action", "")
            if action == "ADD":
                action = "CHECK"
            db_lvl = str(p.get("db_lvl") or "").strip()
            lvl_val = db_lvl if db_lvl else bm["lvl"]
            lvl_src = "DB" if db_lvl else "BASE"
            return {**p,
                    "part_no": matched_pno, "action": action,
                    "base_type": bm["type"], "in_base": True,
                "lvl": lvl_val, "lvl_source": lvl_src,
                "qty": bm["qty"]}
        else:
            # base에 없음 → 신규 부품 (DB 정보 그대로)
            action = p.get("action", "ADD")
            db_lvl = str(p.get("db_lvl") or "").strip()
            return {**p,
                    "part_no": db_new or db_base or "",
                    "base_type": "", "in_base": False,
                "action": action,
                "lvl": db_lvl,
                "lvl_source": "DB" if db_lvl else "",
                "qty": "1"}

    mapped_core    = [map_to_base(p) for p in unique_parts if p["tier"] in ("CORE", "CORE_GROUP")]
    mapped_cascade = [map_to_base(p) for p in unique_parts if p["tier"] == "CASCADE"]

    # ═══════════════════════════════════════════════════
    # Step 3.5: 채번 필요 부품 레벨 추정
    # ═══════════════════════════════════════════════════
    def _estimate_lvl(parts_in_group):
        lvls = []
        for p in parts_in_group:
            lv = (p.get("lvl") or "").strip()
            if lv and lv != "0":
                lvls.append(lv)
        if not lvls:
            return "..2"
        c = Counter(lvls)
        return c.most_common(1)[0][0]

    _all_mapped = mapped_core + mapped_cascade
    _lvl_by_l1 = defaultdict(list)
    for mp in _all_mapped:
        _lvl_by_l1[mp.get("l1_desc") or "(미분류)"].append(mp)

    for l1_key, group in _lvl_by_l1.items():
        est_lvl = _estimate_lvl(group)
        for mp in group:
            if not (mp.get("lvl") or "").strip():
                mp["lvl"] = est_lvl
                mp["lvl_estimated"] = True
                mp["lvl_source"] = "ESTIMATED"

    # ═══════════════════════════════════════════════════
    # Step 3.6: 공용품/참고품번/신규 분류
    # ═══════════════════════════════════════════════════
    def classify_sourcing(p):
        act = (p.get("action") or "").upper()
        if act in ("DELETE", "REMOVE"):
            p["sourcing"] = ""
            p["sourcing_reason"] = ""
            return p

        chg = (p.get("chg") or "").lower()
        db_pno = p.get("db_new_pno") or p.get("db_base_pno") or ""
        in_base = p.get("in_base", False)

        if "공용" in chg:
            # 공용품: DB 품번을 표시 품번으로
            if db_pno and not in_base:
                p["display_pno"] = db_pno
            p["sourcing"] = "공용검토"
            p["sourcing_reason"] = chg
        elif db_pno and not in_base:
            # 참고품번: DB에 품번 있지만 base에 없음
            p["display_pno"] = db_pno
            p["sourcing"] = "참고품번"
            p["sourcing_reason"] = chg if chg else p.get("rsn","")
        elif not db_pno and not in_base:
            # 신규
            p["sourcing"] = "신규"
            p["sourcing_reason"] = ""
        else:
            p["sourcing"] = ""
            p["sourcing_reason"] = ""
        return p

    mapped_core = [classify_sourcing(p) for p in mapped_core]
    mapped_cascade = [classify_sourcing(p) for p in mapped_cascade]

    def _merge_same_part_rows(parts: list[dict]) -> list[dict]:
        # Keep one row per same part/action/L1 and prefer stronger tier.
        # Recall 보호: 품번이 다르면 절대 병합하지 않고,
        # 품번이 없는 경우에만 (부품명+유형+action) 기준으로 병합한다.
        rank = {"CORE": 3, "CORE_GROUP": 2, "CASCADE": 1}
        sourcing_rank = {"참고품번": 3, "공용검토": 2, "신규": 1, "": 0}
        merged: dict[tuple, dict] = {}

        for p in parts or []:
            key = _part_identity_key(p, l1_desc=p.get("l1_desc") or "")
            if key not in merged:
                merged[key] = dict(p)
                continue

            cur = merged[key]
            cur_rank = rank.get(str(cur.get("tier") or ""), 0)
            new_rank = rank.get(str(p.get("tier") or ""), 0)
            
            # tier 우선순위로 선택
            if new_rank > cur_rank:
                chosen = dict(p)
            elif new_rank < cur_rank:
                chosen = dict(cur)
            else:
                # tier가 같으면 sourcing으로 결정 (참고품번 > 공용검토 > 신규)
                cur_src = str(cur.get("sourcing") or "")
                new_src = str(p.get("sourcing") or "")
                cur_src_rank = sourcing_rank.get(cur_src, 0)
                new_src_rank = sourcing_rank.get(new_src, 0)
                chosen = dict(p) if new_src_rank > cur_src_rank else dict(cur)
            
            other = cur if chosen is not cur else p

            # Keep trace fields from both rows when they differ.
            for fld in ("rsn", "chg", "source_doc", "sourcing_reason"):
                a = str(chosen.get(fld) or "").strip()
                b = str(other.get(fld) or "").strip()
                if not a:
                    chosen[fld] = b
                elif b and b not in a:
                    chosen[fld] = f"{a} | {b}"

            merged[key] = chosen

        return list(merged.values())

    mapped_core = _merge_same_part_rows(mapped_core)
    mapped_cascade = _merge_same_part_rows(mapped_cascade)

    # CORE/CORE_GROUP와 CASCADE를 리스트 간 교차로 한 번 더 병합한다.
    # 같은 부품이 두 리스트에 동시에 있으면 CORE 계열을 우선 유지한다.
    def _merge_across_core_cascade(core_parts: list[dict], cascade_parts: list[dict]) -> tuple[list[dict], list[dict]]:
        tier_rank = {"CORE": 3, "CORE_GROUP": 2, "CASCADE": 1}
        src_rank = {"참고품번": 3, "공용검토": 2, "신규": 1, "": 0}
        merged: dict[tuple, dict] = {}

        for p in (core_parts or []) + (cascade_parts or []):
            key = _part_identity_key(p, l1_desc=p.get("l1_desc") or "")
            if key not in merged:
                merged[key] = dict(p)
                continue

            cur = merged[key]
            cur_t = tier_rank.get(_norm_merge_text(cur.get("tier") or ""), 0)
            new_t = tier_rank.get(_norm_merge_text(p.get("tier") or ""), 0)
            if new_t > cur_t:
                chosen = dict(p)
                other = cur
            elif new_t < cur_t:
                chosen = dict(cur)
                other = p
            else:
                cur_s = src_rank.get(str(cur.get("sourcing") or ""), 0)
                new_s = src_rank.get(str(p.get("sourcing") or ""), 0)
                chosen = dict(p) if new_s > cur_s else dict(cur)
                other = cur if chosen is not cur else p

            for fld in ("rsn", "chg", "source_doc", "sourcing_reason"):
                a = str(chosen.get(fld) or "").strip()
                b = str(other.get(fld) or "").strip()
                if not a:
                    chosen[fld] = b
                elif b and b not in a:
                    chosen[fld] = f"{a} | {b}"

            merged[key] = chosen

        out_core = [p for p in merged.values() if _norm_merge_text(p.get("tier") or "") in ("CORE", "CORE_GROUP")]
        out_cascade = [p for p in merged.values() if _norm_merge_text(p.get("tier") or "") == "CASCADE"]
        return out_core, out_cascade

    mapped_core, mapped_cascade = _merge_across_core_cascade(mapped_core, mapped_cascade)

    selected_region = _infer_project_region()
    _core_before_region = len(mapped_core)
    _cascade_before_region = len(mapped_cascade)
    mapped_core, dropped_core = _filter_parts_by_region(mapped_core, selected_region)
    mapped_cascade, dropped_cascade = _filter_parts_by_region(mapped_cascade, selected_region)

    # 과도한 지역 필터로 추천 수가 급감하면 자동 완화(Recall 보호)
    _before_total = _core_before_region + _cascade_before_region
    _after_total = len(mapped_core) + len(mapped_cascade)
    if _before_total >= 12 and _after_total < max(8, int(_before_total * 0.55)):
        mapped_core = [p for p in (mapped_core + dropped_core)]
        mapped_cascade = [p for p in (mapped_cascade + dropped_cascade)]
        dropped_core = []
        dropped_cascade = []
        st.session_state["_region_filter_relaxed"] = True
    else:
        st.session_state["_region_filter_relaxed"] = False

    st.session_state["part_region_filter_last"] = {
        "region": selected_region,
        "dropped_core": dropped_core,
        "dropped_cascade": dropped_cascade,
    }

    # ═══ Step 3.7: Base BOM skeleton 기반 트리 재정렬 ═══
    base_df = st.session_state.get("base_df")
    if isinstance(base_df, pd.DataFrame) and len(base_df) > 0:
        _lc = pick_col(base_df, ["Lvl", "LVL", "Level", "레벨"])
        _pc = pick_col(base_df, ["P/NO", "P/NO.", "Part No", "품번"])
        _dc = pick_col(base_df, ["Description", "DESC", "부품명"])

        def _find_l1_pno(l1_key: str) -> tuple[str, str]:
            if not (_lc and _pc and _dc):
                return "", l1_key
            lk = _norm_merge_text(l1_key or "")
            for _, _row in base_df.iterrows():
                if str(_row.get(_lc, "") or "").strip() != ".1":
                    continue
                rd = _norm_merge_text(_row.get(_dc, "") or "")
                if lk and (lk in rd or rd in lk):
                    return str(_row.get(_pc, "") or "").strip(), str(_row.get(_dc, "") or "").strip()
            return "", l1_key

        _core_by_l1 = defaultdict(list)
        for p in mapped_core:
            _core_by_l1[p.get("l1_desc") or "(미분류)"].append(p)

        reordered_core = []
        for l1_key, parts_group in _core_by_l1.items():
            l1_pno, l1_desc_val = _find_l1_pno(l1_key)
            matched, unmatched = reorder_parts_by_base_skeleton(parts_group, base_df, l1_pno, l1_desc_val)
            reordered_core.extend(matched)
            reordered_core.extend(unmatched)
        mapped_core = reordered_core

        _cascade_by_l1 = defaultdict(list)
        for p in mapped_cascade:
            _cascade_by_l1[p.get("l1_desc") or "(미분류)"].append(p)

        reordered_cascade = []
        for l1_key, parts_group in _cascade_by_l1.items():
            l1_pno, l1_desc_val = _find_l1_pno(l1_key)
            matched, unmatched = reorder_parts_by_base_skeleton(parts_group, base_df, l1_pno, l1_desc_val)
            reordered_cascade.extend(matched)
            reordered_cascade.extend(unmatched)
        mapped_cascade = reordered_cascade

    # ═══════════════════════════════════════════════════
    # Step 4: L1 대표행 + 그룹핑
    # ═══════════════════════════════════════════════════
    base_lvl1_list = []
    if isinstance(base_df, pd.DataFrame) and len(base_df) > 0 and lvl_col and desc_col and pno_col:
        for _, row in base_df.iterrows():
            if str(row.get(lvl_col,"") or "").strip() == ".1":
                base_lvl1_list.append({
                    "part_no": str(row.get(pno_col,"") or "").strip(),
                    "desc":    str(row.get(desc_col,"") or "").strip(),
                })

    def find_base_lvl1(l1_desc):
        l1u = (l1_desc or "").upper()
        for bl in base_lvl1_list:
            if bl["desc"].upper() == l1u:
                return bl
            if l1u and l1u in bl["desc"].upper():
                return bl
            if bl["desc"].upper() and bl["desc"].upper() in l1u:
                return bl
        return None

    lvl1_groups = defaultdict(list)
    for mp in mapped_core:
        lvl1_groups[mp.get("l1_desc") or "(미분류)"].append(mp)

    # ═══════════════════════════════════════════════════
    # Step 5: Proposal 생성
    # ═══════════════════════════════════════════════════
    proposals = []
    seq = 1

    for l1_desc, parts in lvl1_groups.items():
        base_lvl1 = find_base_lvl1(l1_desc)
        base_l1_pno  = base_lvl1["part_no"] if base_lvl1 else ""
        base_l1_desc = base_lvl1["desc"] if base_lvl1 else l1_desc

        changed_descs_upper = {mp["desc"].upper() for mp in parts}
        cascade_descs_upper = {mp["desc"].upper() for mp in mapped_cascade
                               if mp.get("l1_desc") == l1_desc}
        all_part_descs = changed_descs_upper | cascade_descs_upper
        skip_lvl1_row = base_l1_desc.upper() in all_part_descs

        # ── changed_parts: CORE ──
        changed_parts = []
        for mp in parts:
            changed_parts.append({
                "part_name":  mp["desc"],
                "part_no":    mp.get("part_no") or "",
                "display_pno": mp.get("display_pno") or (mp.get("part_no") if mp["in_base"] else "(채번 필요)"),
                "action":     mp["action"],
                "in_base":    mp["in_base"],
                "rsn":        mp.get("rsn",""),
                "chg":        mp.get("chg",""),
                "base_type":  mp.get("base_type",""),
                "lvl":        mp.get("lvl",""),
                "qty":        mp.get("qty","1"),
                "tier":       mp.get("tier") or "CORE",
                "source_doc": mp.get("source_doc",""),
                "sourcing":       mp.get("sourcing",""),         # ✅ 추가
                "sourcing_reason": mp.get("sourcing_reason",""), # ✅ 추가
                "skeleton_order": mp.get("skeleton_order", 999999),
            })

        # ── indirect_parts: CASCADE (같은 L1) ──
        #    L1 대표행과 동일 부품은 중복이니 제거
        indirect_for_l1 = [mp for mp in mapped_cascade
                           if mp.get("l1_desc") == l1_desc]
        indirect_display = []
        for mp in indirect_for_l1:
            # L1 대표행과 동일하면 skip
            if mp.get("part_no") and mp["part_no"] == base_l1_pno:
                continue
            if mp["desc"].upper() == base_l1_desc.upper():
                continue
            indirect_display.append({
                "part_name":  mp["desc"],
                "part_no":    mp.get("part_no") or "",
                "display_pno": mp.get("part_no") if mp["in_base"] else "(채번 필요)",
                "action":     mp["action"],
                "in_base":    mp["in_base"],
                "rsn":        mp.get("rsn",""),
                "chg":        mp.get("chg",""),
                "base_type":  mp.get("base_type",""),
                "lvl":        mp.get("lvl",""),       # ✅ 추가
                "qty":        mp.get("qty","1"),       # ✅ 추가
                "tier":       mp.get("tier") or "CASCADE",
                "source_doc": mp.get("source_doc",""),
                "skeleton_order": mp.get("skeleton_order", 999999),
            })

        # NOTE:
        # 과거에는 L1 미귀속 CASCADE를 첫 proposal에 합류시켰는데,
        # 이 로직이 "Cover,Heater 하위에 PCB" 같은 잘못된 귀속을 만들었다.
        # 따라서 orphan 강제 합류는 비활성화한다.

        # ── existing_parts: base L1 하위 중 변경 안 된 부품 ──
        existing_parts = []
        if base_l1_pno and isinstance(base_df, pd.DataFrame) and parent_col and desc_col and pno_col:
            children = base_df[base_df[parent_col].astype(str).str.strip() == base_l1_pno]
            changed_descs = {p["desc"].upper() for p in parts + indirect_for_l1}
            changed_pnos = set()
            for p in parts + indirect_for_l1:
                pn = p.get("part_no") or ""
                if pn:
                    changed_pnos.add(re.sub(r'[^A-Z0-9]','', pn.upper()))

            for _, row in children.iterrows():
                cd   = str(row.get(desc_col,"") or "").strip()
                cpno = str(row.get(pno_col,"") or "").strip()
                cpno_key = re.sub(r'[^A-Z0-9]','', cpno.upper())
                if cd.upper() in changed_descs or cpno_key in changed_pnos:
                    continue
                existing_parts.append({
                    "part_name": cd, "part_no": cpno,
                    "action": "KEEP", "in_base": True,
                    "base_type": str(row.get(type_col,"") or "").strip() if type_col else "",
                    "lvl": str(row.get(lvl_col,"") or "").strip() if lvl_col else "",
                    "qty": str(row.get("Qty","") or row.get("QTY","") or "1").strip(),
                })

        # ── 출처 문서 목록 ──
        source_docs = sorted(set(
            mp.get("source_doc","") for mp in parts + indirect_for_l1
            if mp.get("source_doc")
        ))

        proposals.append({
            "proposal_id": f"P-{seq:03d}",
            "status": "PENDING",
            "change_summary": intent["raw_text"],
            "target": {
                "main_object": intent.get("target_object",""),
                "action": intent.get("action","ADD"),
            },
            "lvl1": {
                "desc": base_l1_desc,
                "part_no": base_l1_pno,
                "in_base": bool(base_lvl1),
                "user_action": "",
                "skip": skip_lvl1_row,
            },
            "changed_parts": changed_parts,
            "indirect_parts": indirect_display,
            "existing_parts": existing_parts,
            "confidence": round(
                len(changed_parts) / max(len(mapped_core), 1), 2
            ),
            "source_docs": source_docs,
            "ref_models": list(set(
                (doc.get("metadata") or {}).get("model","")
                for doc in (primary_docs or [])
                if (doc.get("metadata") or {}).get("model","")
            ))[:3],
        })
        seq += 1

    proposals.sort(key=lambda x: x["confidence"], reverse=True)
    return proposals

# =========================================================
# PATCH 1) OK/재검토 상태 머신용 session + helper
# =========================================================

RECHECK_MODES = ["PART_ONLY", "PATH_ONLY", "BOTH"]

def init_review_state():
    ss = st.session_state
    ss.setdefault("proposals", [])           # 카드 리스트
    ss.setdefault("decisions", {})           # pid -> {"status": "OK/RECHECK", "mode": "..."}
    ss.setdefault("evidence_pool", {})       # pid -> evidences(list). (카드별 evidence 재사용)
    ss.setdefault("cards_for_review", [])

def set_decision(pid: str, status: str, mode: str):
    st.session_state["decisions"][pid] = {
        "status": (status or "OK").upper(),
        "mode": (mode or "BOTH").upper(),
    }

def all_reviewed(proposals: list[dict]) -> bool:
    dec = st.session_state.get("decisions") or {}
    for p in proposals or []:
        pid = p.get("proposal_id")
        if not pid or pid not in dec:
            return False
        if dec[pid]["status"] not in ["OK", "RECHECK"]:
            return False
    return True


def build_recheck_query_from_card(card: dict) -> str:
    """
    원본에 build_rerank_query()가 있는데(스켈레톤) 
    여기선 재검토 모드 분기까지 포함한 최소 템플릿을 제공.
    """
    mode = (card.get("recheck_mode") or "BOTH").upper()
    mode = mode if mode in RECHECK_MODES else "BOTH"

    obj = (card.get("target") or {}).get("main_object","")
    feat = (card.get("payload") or {}).get("feature","")
    path = (card.get("bom_target") or {}).get("apply_bom_path","")
    main0 = ((card.get("payload") or {}).get("main_parts") or [{}])[0]
    pname = main0.get("part_name","")
    pno = main0.get("part_no","")

    if mode == "PART_ONLY":
        return f"OBJECT:{obj}\nFEATURE:{feat}\nBOM_PATH:{path}\nSEARCH:PART_VARIATION"
    if mode == "PATH_ONLY":
        return f"OBJECT:{obj}\nFEATURE:{feat}\nPART_NAME:{pname}\nPART_NO:{pno}\nSEARCH:PATH_VARIATION"
    return f"OBJECT:{obj}\nFEATURE:{feat}\nCHANGE:{card.get('change_summary','')}\nSEARCH:FULL_RECHECK"

# =========================
# Proposal vs Base Compare + REMOVE 후보 생성
# =========================
import re
from typing import Any, Dict, List, Tuple, Optional

def _n(s: Any) -> str:
    s = "" if s is None else str(s)
    s = re.sub(r"\s+", " ", s.strip().upper())
    return s

def _nk(s: Any) -> str:
    # 키/품번 비교용: 특수문자 제거
    s = _n(s)
    return re.sub(r"[^A-Z0-9]+", "", s)

def _np(p: Any) -> str:
    # 경로 표준화: "A>B" "A / B" 모두 "A > B"
    p = "" if p is None else str(p)
    p = p.replace("/", ">").replace("\\", ">")
    p = re.sub(r"\s*>\s*", " > ", p.strip())
    p = re.sub(r"\s+", " ", p)
    return p

def _path_tokens(p: Any) -> List[str]:
    p = _np(p)
    return [t.strip() for t in p.split(">") if t.strip()]

def _path_prefixes(tokens: List[str]) -> List[str]:
    out, cur = [], []
    for t in tokens:
        cur.append(t)
        out.append(" > ".join(cur))
    return out

# --- Base 인덱스 생성 (빠른 조회)
def index_base_snapshot(base_snapshot: Dict[str, Any]) -> Dict[str, Any]:
    rows = base_snapshot.get("rows") or []
    by_part_no: Dict[str, List[dict]] = {}
    by_part_name: Dict[str, List[dict]] = {}
    by_bom_path: Dict[str, List[dict]] = {}

    for r in rows:
        pno = _nk(r.get("part_no"))
        pname = _nk(r.get("part_name"))
        path = _np(r.get("bom_path"))

        if pno:
            by_part_no.setdefault(pno, []).append(r)
        if pname:
            by_part_name.setdefault(pname, []).append(r)
        if path:
            by_bom_path.setdefault(path, []).append(r)

    return {
        "rows": rows,
        "by_part_no": by_part_no,
        "by_part_name": by_part_name,
        "by_bom_path": by_bom_path,
    }

# --- base row 매칭 (품번 우선, 다음 품명+경로)
def match_base_part(
    base_idx: Dict[str, Any],
    part_no: str = "",
    part_name: str = "",
    bom_path_hint: str = "",
) -> Tuple[Optional[dict], float, str]:
    pno = _nk(part_no)
    pname = _nk(part_name)
    hint = _np(bom_path_hint)

    # 1) 품번 exact
    if pno and pno in base_idx["by_part_no"]:
        cand = base_idx["by_part_no"][pno]
        if hint:
            pref = set(_path_prefixes(_path_tokens(hint)))
            best, best_score = None, -1
            for r in cand:
                rpath = _np(r.get("bom_path"))
                score = 2 if (rpath in pref) else 1
                if score > best_score:
                    best, best_score = r, score
            return best, (0.98 if best_score == 2 else 0.95), "PART_NO_MATCH"
        return cand[0], 0.95, "PART_NO_MATCH"

    # 2) 품명 + 경로 prefix
    if pname and pname in base_idx["by_part_name"] and hint:
        cand = base_idx["by_part_name"][pname]
        pref = set(_path_prefixes(_path_tokens(hint)))
        for r in cand:
            if _np(r.get("bom_path")) in pref:
                return r, 0.85, "PART_NAME+PATH_PREFIX_MATCH"

    # 3) 품명만
    if pname and pname in base_idx["by_part_name"]:
        return base_idx["by_part_name"][pname][0], 0.70, "PART_NAME_MATCH"

    # 4) 경로만
    if hint and hint in base_idx["by_bom_path"]:
        return base_idx["by_bom_path"][hint][0], 0.55, "PATH_MATCH_ONLY"

    return None, 0.0, "NO_MATCH"

def _iter_proposal_parts(proposal: Dict[str, Any]) -> List[Tuple[str, dict]]:
    out = []
    parts = (proposal or {}).get("parts") or {}
    for role in ("main", "sub"):
        for p in (parts.get(role) or []):
            out.append((role, p))
    return out

# --- proposal 1개를 base와 비교 → base_match + diffs + type 보정
def compare_proposal_to_base(base_snapshot: Dict[str, Any], proposal: Dict[str, Any]) -> Dict[str, Any]:
    base_idx = index_base_snapshot(base_snapshot)

    bom_target = (proposal or {}).get("bom_target") or {}
    apply_path = _np(bom_target.get("apply_bom_path", ""))
    base_hint = _np(bom_target.get("base_bom_path", ""))

    diffs = []
    any_exists, all_exists = False, True

    for role, p in _iter_proposal_parts(proposal):
        pno = p.get("part_no", "")
        pname = p.get("part_name", "")
        hint_path = apply_path or base_hint or p.get("bom_path", "")

        base_row, conf, why = match_base_part(base_idx, pno, pname, hint_path)

        exists = base_row is not None
        any_exists = any_exists or exists
        all_exists = all_exists and exists

        p["base_match"] = {
            "exists_in_base": exists,
            "confidence": conf,
            "reason": why,
            "base_row_id": base_row.get("row_id") if base_row else None,
            "base_part_no": base_row.get("part_no") if base_row else None,
            "base_part_name": base_row.get("part_name") if base_row else None,
            "base_bom_path": base_row.get("bom_path") if base_row else None,
        }

        if not exists:
            diffs.append({
                "op": "ADD",
                "role": role,
                "part_no": pno,
                "part_name": pname,
                "apply_bom_path": apply_path or base_hint,
                "note": "Base에 없음 → 추가",
            })
        else:
            changed = []
            if _nk(pno) and _nk(base_row.get("part_no")) and _nk(pno) != _nk(base_row.get("part_no")):
                changed.append("part_no")
            if _nk(pname) and _nk(base_row.get("part_name")) and _nk(pname) != _nk(base_row.get("part_name")):
                changed.append("part_name")
            if apply_path and _np(base_row.get("bom_path")) and apply_path != _np(base_row.get("bom_path")):
                changed.append("bom_path")

            if changed:
                diffs.append({
                    "op": "MODIFY",
                    "role": role,
                    "part_no": pno,
                    "part_name": pname,
                    "base_part_no": base_row.get("part_no"),
                    "base_part_name": base_row.get("part_name"),
                    "base_bom_path": base_row.get("bom_path"),
                    "apply_bom_path": apply_path or base_hint,
                    "changed_fields": changed,
                    "note": "Base 존재 → 변경",
                })
            else:
                diffs.append({
                    "op": "KEEP",
                    "role": role,
                    "part_no": pno,
                    "part_name": pname,
                    "base_part_no": base_row.get("part_no"),
                    "base_bom_path": base_row.get("bom_path"),
                    "note": "Base와 동일/유사",
                })

    # 카드 타입 보정
    if not any_exists:
        ptype = "ADD"
    elif all_exists:
        ptype = "MODIFY"
    else:
        ptype = "ADD"

    proposal["proposal_type"] = proposal.get("proposal_type") or ptype
    proposal["base_comparison"] = {
        "any_exists_in_base": any_exists,
        "all_exists_in_base": all_exists,
        "apply_bom_path": apply_path,
        "base_bom_path_hint": base_hint,
    }
    proposal["diffs"] = diffs
    return proposal

def compare_all_proposals(base_snapshot: Dict[str, Any], proposals: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    return [compare_proposal_to_base(base_snapshot, p) for p in (proposals or [])]


# =========================
# REMOVE 후보 생성 (보수적으로!)
# =========================
def compute_remove_candidates(
    base_snapshot: Dict[str, Any],
    proposals_compared: List[Dict[str, Any]],
    only_ok: bool = False,
) -> List[Dict[str, Any]]:
    """
    ✅ 안전한(보수적) REMOVE 생성 규칙:
    - 명시적으로 '대체'된 base 품번만 REMOVE 후보로 만듦
      (MODIFY diff에서 base_part_no가 있고, 신규 part_no와 다를 때)
    - 또는 proposal 자체가 REMOVE 타입이면 그대로 반영 가능(추후)
    - base 전체에서 '안 쓰일 것 같음' 추측으로 삭제는 안 함 (오탐 방지)
    """
    replaced_base_pnos = set()
    scope_paths = set()

    for pr in proposals_compared or []:
        if only_ok:
            if (pr.get("decision") or {}).get("status") != "OK":
                continue

        # scope: 적용/힌트 경로
        bc = pr.get("base_comparison") or {}
        if bc.get("apply_bom_path"):
            scope_paths.add(_np(bc["apply_bom_path"]))
        if bc.get("base_bom_path_hint"):
            scope_paths.add(_np(bc["base_bom_path_hint"]))

        # 대체(replace)된 base 품번 수집
        for d in (pr.get("diffs") or []):
            if d.get("op") == "MODIFY":
                bpn = d.get("base_part_no")
                npn = d.get("part_no")
                if _nk(bpn) and _nk(npn) and _nk(bpn) != _nk(npn):
                    replaced_base_pnos.add(_nk(bpn))

    if not replaced_base_pnos:
        return []

    # base rows에서 replaced 품번 찾아 REMOVE 카드 만들기
    base_rows = base_snapshot.get("rows") or []
    remove_cards = []
    seq = 1

    for r in base_rows:
        bpn = _nk(r.get("part_no"))
        if not bpn or bpn not in replaced_base_pnos:
            continue

        # scope_path가 있으면 그 범위 내만(더 안전)
        rpath = _np(r.get("bom_path"))
        if scope_paths:
            ok = False
            for sp in scope_paths:
                if sp and rpath.startswith(sp):
                    ok = True
                    break
            if not ok:
                continue

        remove_cards.append({
            "proposal_id": f"R-{seq:03d}",
            "change_id": "REMOVE-AUTO",
            "change_summary": "Base 대비 대체/삭제 후보(자동)",
            "proposal_type": "REMOVE",
            "title": f"대체로 인한 Base 부품 삭제 후보: {r.get('part_no')} / {r.get('part_name')}",
            "bom_target": {
                "base_bom_path": rpath,
                "apply_bom_path": rpath,
                "confidence": 0.60
            },
            "parts": {
                "main": [{
                    "part_no": r.get("part_no"),
                    "part_name": r.get("part_name"),
                    "description": "",
                    "qty": r.get("qty", 1),
                    "is_new": False
                }],
                "sub": []
            },
            "base_comparison": {
                "exists_in_base": True,
                "base_part_no": r.get("part_no"),
                "base_bom_path": rpath
            },
            "rationale": {
                "why": ["다른 제안(MODIFY)에서 해당 Base 품번이 다른 품번으로 대체됨"],
                "evidence": []
            },
            "review_points": ["정말 삭제가 맞는지(대체 관계 확인)", "수량/경로 중복 존재 여부 확인"],
            "decision": {"status": "PENDING", "comment": ""},
            "meta": {"generated_by": "RULE_ENGINE", "confidence_score": 0.60}
        })
        seq += 1

    return remove_cards

# ✅ build_recommendations_grouped 결과(dict) → app.py proposals 리스트 포맷 변환
def _convert_rec_to_proposals(main_list: list, sub_list: list) -> list:
    results = []
    for i, p in enumerate(main_list or []):
        pn   = getattr(p, "part_no", "") if hasattr(p, "part_no") else p.get("part_no", "")
        desc = getattr(p, "description", "") if hasattr(p, "description", "") else p.get("description", "")
        path = getattr(p, "bom_path", "") if hasattr(p, "bom_path", "") else p.get("bom_path", "")
        results.append({
            "proposal_id": f"F-MAIN-{i+1:03d}",
            "status": "PENDING",
            "recheck_mode": "BOTH",
            "confidence": 0.5,
            "change_summary": "",
            "proposal_type": "MODIFY",
            "title": f"{pn} {desc}".strip(),
            "target": {"main_object": desc, "action": "MODIFY"},
            "payload": {
                "feature": desc,
                "main_parts": [{"part_no": pn, "description": desc}],
                "sub_parts": [],
            },
            "bom_target": {"apply_bom_path": path, "path_confidence": 0.5},
            "reason": ["RAG 유사사례 기반 fallback 추천"],
            "review_points": ["부품 적합성 확인", "BOM 경로 확인"],
            "parts": {
                "main": [{"part_no": pn, "description": desc, "bom_path": path}]
            },
        })
    for i, p in enumerate(sub_list or []):
        pn   = getattr(p, "part_no", "") if hasattr(p, "part_no") else p.get("part_no", "")
        desc = getattr(p, "description", "") if hasattr(p, "description", "") else p.get("description", "")
        path = getattr(p, "bom_path", "") if hasattr(p, "bom_path", "") else p.get("bom_path", "")
        results.append({
            "proposal_id": f"F-SUB-{i+1:03d}",
            "status": "PENDING",
            "recheck_mode": "BOTH",
            "confidence": 0.4,
            "change_summary": "",
            "proposal_type": "MODIFY",
            "title": f"{pn} {desc}".strip(),
            "target": {"main_object": desc, "action": "MODIFY"},
            "payload": {
                "feature": desc,
                "main_parts": [],
                "sub_parts": [{"part_no": pn, "description": desc}],
            },
            "bom_target": {"apply_bom_path": path, "path_confidence": 0.4},
            "reason": ["RAG 유사사례 기반 fallback 추천 (sub)"],
            "review_points": ["부품 적합성 확인", "BOM 경로 확인"],
            "parts": {
                "main": [{"part_no": pn, "description": desc, "bom_path": path}]
            },
        })
    return results

def _convert_rag_docs_to_proposals(raw_docs: list, change_text: str = "") -> list:
    """
    retrieve_docs 결과에서 부품 후보 추출.
    doc_packaging.py 저장 포맷 기준:
      - row 줄: "- row=N L1 | 품번 | 품명 | 수량 | 단위 | ..."
      - meta 키: new_part_no, base_part_no, description, bom_path 등
    """
    import re
    results = []
    seen_pn = set()

    for doc in (raw_docs or []):
        text  = doc.get("text") or ""
        meta  = doc.get("meta") or {}
        doc_id = doc.get("id", "")

        # ── ① meta에서 직접 품번/품명 추출 (가장 신뢰도 높음) ──
        for pn_key, desc_key in [
            ("new_part_no", "description"),
            ("part_no",     "description"),
            ("base_part_no","description"),
        ]:
            pn   = str(meta.get(pn_key) or "").strip().upper()
            desc = str(meta.get(desc_key) or meta.get("desc") or "").strip()
            bom  = str(meta.get("bom_path") or meta.get("path") or "").strip()
            if pn and len(pn) >= 5 and pn not in seen_pn:
                seen_pn.add(pn)
                results.append(_make_proposal(pn, desc, bom, doc_id, change_text, len(results)))
                break  # 같은 doc에서 pn_key 중복 방지

        # ── ② text 본문 파싱 (PART_NO: xxx 형태 + pipe 형태 둘 다 지원) ──
        for ln in text.splitlines():
            ln = ln.strip()
            if not ln:
                continue

            # 포맷 A: "PART_NO: MHL49207705"
            m_pno = re.match(r"PART_NO\s*:\s*(.+)", ln, re.I)
            if m_pno:
                raw_pn = re.sub(r"[^A-Z0-9]", "", m_pno.group(1).upper())
                if 5 <= len(raw_pn) <= 16 and raw_pn not in seen_pn:
                    # 같은 doc 내에서 PART_NAME 찾기
                    desc = ""
                    bom  = ""
                    for ln2 in text.splitlines():
                        m2 = re.match(r"PART_NAME\s*:\s*(.+)", ln2.strip(), re.I)
                        if m2:
                            desc = m2.group(1).strip()
                        m3 = re.match(r"BOM_PATH\s*:\s*(.+)", ln2.strip(), re.I)
                        if m3:
                            bom = m3.group(1).strip()
                    seen_pn.add(raw_pn)
                    results.append(_make_proposal(raw_pn, desc, bom, doc_id, change_text, len(results)))
                continue  # PART_NO 줄은 처리 완료

            # 포맷 B: "- row=N L1 | 품번 | 품명 | ..."
            if ln.startswith("- row=") and "|" in ln:
                parts = [p.strip() for p in ln.split("|")]
                if len(parts) >= 3:
                    raw_pn = re.sub(r"[^A-Z0-9]", "", parts[1].upper())
                    pn   = raw_pn if 5 <= len(raw_pn) <= 16 else ""
                    desc = parts[2] if len(parts) > 2 else ""
                    bom  = parts[5] if len(parts) > 5 else ""
                    uid  = pn or desc
                    if uid and uid not in seen_pn:
                        seen_pn.add(uid)
                        results.append(_make_proposal(pn, desc, bom, doc_id, change_text, len(results)))


def _make_proposal(pn: str, desc: str, bom: str, doc_id: str, change_text: str, idx: int) -> dict:
    label = f"{pn} | {desc}".strip(" |") if pn or desc else doc_id
    return {
        "proposal_id":  f"RAG-{idx+1:03d}",
        "status":       "PENDING",
        "recheck_mode": "BOTH",
        "confidence":   0.5,
        "change_summary": change_text,
        "proposal_type": "MODIFY",
        "title":        label,
        "target":       {"main_object": desc or pn, "action": "MODIFY"},
        "payload": {
            "feature":    desc,
            "main_parts": [{"part_no": pn, "description": desc}],
            "sub_parts":  [],
        },
        "bom_target": {"apply_bom_path": bom, "path_confidence": 0.5},
        "reason":        [f"RAG 유사사례({doc_id}) 기반 추천"],
        "review_points": ["부품 적합성 확인", "BOM 경로 확인"],
        "parts": {
            "main": [{"part_no": pn, "description": desc, "bom_path": bom}]
        },
    }

# =========================================================
# 부품 Sourcing 자동 분류 (보수적 디폴트)
# =========================================================

_COMMON_USE_KW = {"MODULE", "모듈", "SENSOR", "센서"}

def classify_sourcing(part_name: str, part_no: str = "", action: str = "") -> str:
    """
    부품명 기반 sourcing 자동 분류
    - 공용검토: 모듈 단품류 (스펙 동일 시 그대로 사용 가능)
    - 신규: 그 외 전부 (보수적 디폴트)
    - "": KEEP/REMOVE 등 분류 불필요
    """
    act = (action or "").upper()
    if act in ("KEEP", "REMOVE", "DELETE", "삭제", "유지", ""):
        return ""

    name_up = (part_name or "").strip().upper()
    pno = (part_no or "").strip()

    # 품번 자체가 없으면 무조건 신규
    if not pno or pno.lower() == "nan":
        return "신규"

    # 공용검토: 모듈 단품 (ASS'Y 붙으면 제외)
    is_module = any(kw in name_up for kw in _COMMON_USE_KW)
    is_assy = any(kw in name_up for kw in ("ASSY", "ASS'Y", "ASSEMBLY"))

    if is_module and not is_assy:
        return "공용검토"

    # 디폴트: 신규 (보수적)
    return "신규"

## ── Helper: proposals → L1 그룹별 ──
def _build_proposal_df(proposals, base_df):
    ACTION_MAP = {"ADD": "추가", "MODIFY": "변경", "REMOVE": "삭제", "DELETE": "삭제", "KEEP": "유지", "CHECK": "⚠️확인필요"}
    base_map = {}
    base_order_by_pno = {}
    base_order_by_desc = {}
    if isinstance(base_df, pd.DataFrame) and len(base_df) > 0:
        pc = pick_col(base_df, ["P/NO","P/NO.","P/no.","Part No","품번","부품번호","PNO","PARTNO"])
        dc = pick_col(base_df, ["Description","DESC","DESC.","Desc","부품명","품명"])
        lc = pick_col(base_df, ["Lvl","LVL","Level","LEVEL","레벨"])
        qc = pick_col(base_df, ["Qty","QTY","수량","개수","수량(EA)"])
        tc = pick_col(base_df, ["CMDT","공정","분류","Type","TYPE","유형"])
        for row_idx, row in base_df.iterrows():
            desc_key = re.sub(r'\s+',' ',str(row.get(dc,"")).strip().upper()) if dc else ""
            if desc_key and desc_key not in base_order_by_desc:
                base_order_by_desc[desc_key] = int(row_idx)
            if pc:
                pn = re.sub(r'[^A-Z0-9]','',str(row.get(pc,"")).strip().upper())
                if pn:
                    if pn not in base_order_by_pno:
                        base_order_by_pno[pn] = int(row_idx)
                    base_map[pn] = {
                        "lvl":str(row.get(lc,"")).strip() if lc else "",
                        "qty":str(row.get(qc,"")).strip() if qc else "",
                        "cmdt":str(row.get(tc,"")).strip() if tc else "",
                        "row_idx": int(row_idx),
                    }
    def _sk(v):
        s = str(v).strip()
        if not s: return 998
        if s.startswith("."): return s.count(".")
        try: return float(s)
        except: return 999
    def _base_order(part_name="", part_no=""):
        pkey = re.sub(r'[^A-Z0-9]','',str(part_no or '').strip().upper())
        dkey = re.sub(r'\s+',' ',str(part_name or '').strip().upper())
        if pkey and pkey in base_order_by_pno:
            return base_order_by_pno[pkey]
        if dkey and dkey in base_order_by_desc:
            return base_order_by_desc[dkey]
        return 10**9
    def _extract_target_objects_local(change_text):
        t = (change_text or "").lower()
        objs = []
        mapping = {
            "door": ["door", "도어"],
            "cavity": ["cavity", "캐비티"],
            "panel": ["panel", "패널"],
            "camera": ["camera", "카메라"],
            "rack": ["rack", "랙"],
            "tray": ["tray", "트레이"],
            "control": ["control", "제어", "조작"],
            "motor": ["motor", "모터", "bldc"],
            "steam": ["steam", "스팀"],
            "hinge": ["hinge", "힌지"],
            "harness": ["harness", "하네스", "wire", "배선"],
        }
        for vals in mapping.values():
            if any(v in t for v in vals):
                for v in vals:
                    if v not in objs:
                        objs.append(v)
        return objs
    target_objects = _extract_target_objects_local(" ".join(st.session_state.get("change_items") or []))
    def _group_priority(header_desc, child_df):
        if not target_objects:
            return 1
        h = str(header_desc or "").lower()
        if any(obj in h for obj in target_objects):
            return 0
        if isinstance(child_df, pd.DataFrame) and not child_df.empty and "부품명" in child_df.columns:
            names = " ".join(str(v).lower() for v in child_df["부품명"].fillna("").tolist())
            if any(obj in names for obj in target_objects):
                return 0
        return 1
    def _mr(pt):
            a = pt.get("action") or ""
            pno = pt.get("display_pno") or pt.get("part_no") or ""
            if pno.lower() == "nan": pno = ""
            pk = re.sub(r'[^A-Z0-9]','', pno.upper())
            qty=str(pt.get("qty") or ""); cmdt=pt.get("base_type") or ""
            bi = base_map.get(pk,{})
            if bi and bi.get("lvl", ""):
                lvl = bi["lvl"]
            else:
                lvl = pt.get("lvl") or ""
            if not qty and bi: qty=bi.get("qty","")
            if not cmdt and bi: cmdt=bi.get("cmdt","")

            part_name = pt.get("part_name") or pt.get("desc") or ""
            src = pt.get("sourcing") or ""
            src_rsn = pt.get("sourcing_reason") or ""

            # ── sourcing 미지정 시 자동 분류 ──
            if not src:
                src = classify_sourcing(part_name, pno, a)

            # ── 비고 & 품번 표시 결정 ──
            display_pno = pno
            if src == "공용검토":
                bigo = f"🔵 공용검토 · {src_rsn}" if src_rsn else "🔵 공용검토"
                # 품번 유지 (DB 품번 그대로 → 사용자가 검토)
            elif src == "신규":
                if pno:
                    bigo = f"🟡 참고: {pno}"
                    display_pno = "(채번 필요)"
                else:
                    bigo = ""
                    display_pno = "(채번 필요)"
            elif src == "참고품번":
                bigo = f"🟡 참고: {pno}" if pno else ""
                display_pno = "(채번 필요)"
            else:
                bigo = ""

            review_tag = _sanitize_user_remark_text(pt.get("review_tag") or "")
            if review_tag:
                bigo = (bigo + " | " if bigo else "") + review_tag

            bigo = _sanitize_user_remark_text(bigo)

            return {
                "변경유형": ACTION_MAP.get(a.upper(), a),
                "부품명": part_name,
                "품번": display_pno,
                "__dedup_pno": pt.get("part_no") or pt.get("display_pno") or "",
                "레벨": lvl,
                "수량": qty if qty else "",
                "유형": cmdt,
                "변경사유": pt.get("rsn") or "",
                "분류": pt.get("tier") or "",
                "비고": bigo,
                "출처": pt.get("source_doc") or "",
                "skeleton_order": pt.get("skeleton_order", 999999)
            }
    
    CC = ["변경유형","부품명","품번","레벨","수량","유형","변경사유","분류","비고","출처"]

    def _dedup_display_rows(rows: list[dict], l1_desc: str = "") -> list[dict]:
        # 화면 표시용 최종 중복 제거: 같은 부품/유형은 1건으로 합친다.
        tier_rank = {"CORE": 3, "CORE_GROUP": 2, "CASCADE": 1, "": 0}

        def _key(r: dict) -> tuple:
            return _part_identity_key(
                {
                    "part_no": r.get("__dedup_pno") or r.get("품번") or "",
                    "part_name": r.get("부품명") or "",
                    "base_type": r.get("유형") or "",
                    "action": r.get("변경유형") or "",
                    "lvl": r.get("레벨") or "",
                    "source_doc": r.get("출처") or "",
                    "rsn": r.get("변경사유") or "",
                },
                l1_desc=l1_desc,
            )

        merged: dict[tuple, dict] = {}
        for r in rows or []:
            k = _key(r)
            if k not in merged:
                merged[k] = dict(r)
                continue

            cur = merged[k]
            cur_t = tier_rank.get(_norm_merge_text(cur.get("분류") or ""), 0)
            new_t = tier_rank.get(_norm_merge_text(r.get("분류") or ""), 0)
            if new_t > cur_t:
                chosen = dict(r)
                other = cur
            elif new_t < cur_t:
                chosen = dict(cur)
                other = r
            else:
                # 같은 분류면 정보가 더 많은(참고품번/출처/사유) 행 유지
                def _score(x: dict) -> int:
                    s = 0
                    if str(x.get("비고") or "").strip():
                        s += 1
                    if "참고" in str(x.get("비고") or ""):
                        s += 2
                    if str(x.get("출처") or "").strip():
                        s += 1
                    if str(x.get("변경사유") or "").strip():
                        s += 1
                    return s
                chosen = dict(r) if _score(r) > _score(cur) else dict(cur)
                other = cur if chosen is not cur else r

            for fld in ("변경사유", "비고", "출처"):
                a = str(chosen.get(fld) or "").strip()
                b = str(other.get(fld) or "").strip()
                if not a:
                    chosen[fld] = b
                elif b and b not in a:
                    chosen[fld] = f"{a} | {b}"

            # 추가/변경/확인필요가 섞인 경우 표시는 "변경"으로 통일
            if _action_bucket(chosen.get("변경유형") or "") == "CHANGE":
                chosen["변경유형"] = "변경"

            # 품번은 '(채번 필요)'보다 실제 참고 품번이 있으면 비고로 보존된 행을 우선.
            a_pno = str(chosen.get("품번") or "").strip()
            b_pno = str(other.get("품번") or "").strip()
            if (not a_pno or a_pno == "(채번 필요)") and b_pno and b_pno != "(채번 필요)":
                chosen["품번"] = b_pno

            merged[k] = chosen

        # ── dedup 후 skeleton_order 제거 (내부용) ──
        for r in merged.values():
            r.pop("skeleton_order", None)
            r.pop("__dedup_pno", None)
        return list(merged.values())

    groups = []
    for p in proposals:
        cr = []
        seq = 0

        def _lvl_depth(v: str) -> int:
            s = str(v or "").strip()
            if not s:
                return 999
            if s.startswith('.'):
                return s.count('.')
            m = re.search(r'(\d+)$', s)
            if m:
                try:
                    return int(m.group(1))
                except Exception:
                    return 999
            return 999

        for pt in (p.get("changed_parts") or []):
            bord = _base_order(pt.get("part_name") or pt.get("desc") or "", pt.get("part_no") or pt.get("display_pno") or "")
            depth = _lvl_depth(pt.get("lvl") or "")
            sk_order = pt.get("skeleton_order", 999999)
            cr.append((bord, depth, seq, _mr(pt), sk_order))
            seq += 1
        for pt in (p.get("indirect_parts") or []):
            bord = _base_order(pt.get("part_name") or pt.get("desc") or "", pt.get("part_no") or pt.get("display_pno") or "")
            depth = _lvl_depth(pt.get("lvl") or "")
            sk_order = pt.get("skeleton_order", 999999)
            cr.append((bord, depth, seq, _mr(pt), sk_order))
            seq += 1
        # skeleton_order만 기준으로 정렬한다. 레벨 숫자 정렬은 사용하지 않는다.
        cr.sort(key=lambda x: (
            x[4],
            x[2],
        ))
        lv1 = p.get("lvl1") or {}
        lp = lv1.get("part_no") or ""; ld = lv1.get("desc") or "(미분류)"
        lk = re.sub(r'[^A-Z0-9]','',lp.upper())
        lb = base_map.get(lk,{})
        rows = [r for _, _, _, r, _ in cr]
        rows = _dedup_display_rows(rows, l1_desc=ld)
        hi = {"desc":ld,"part_no":lp,"lvl":lb.get("lvl","") or ".1","type":lb.get("cmdt","") or lv1.get("base_type",""),"n_children":len(rows),"skip":lv1.get("skip",False),"source":f"Base BOM ({st.session_state.get('bom_model','')})".strip()}
        cdf = pd.DataFrame(rows,columns=CC) if rows else pd.DataFrame(columns=CC)
        groups.append({"header":hi,"children":cdf,"_priority":_group_priority(ld, cdf),"_order":_base_order(ld, lp)})
    groups.sort(key=lambda g:(g.get("_priority", 1), g.get("_order", 10**9), _sk(g["header"].get("lvl","")), str(g["header"].get("desc",""))))
    for g in groups:
        g.pop("_priority", None)
        g.pop("_order", None)

    seen_pno = set()
    seen_name = set()
    deduped_groups = []
    for g in groups:
        cdf = g.get("children")
        if not isinstance(cdf, pd.DataFrame) or cdf.empty:
            continue

        keep_idx = []
        for idx, row in cdf.iterrows():
            pno_key = re.sub(r'[^A-Z0-9]', '', str(row.get("품번", "")).upper())
            if pno_key:
                if pno_key in seen_pno:
                    continue
                seen_pno.add(pno_key)
                keep_idx.append(idx)
                continue

            name_key = str(row.get("부품명", "")).strip().upper()
            if name_key:
                if name_key in seen_name:
                    continue
                seen_name.add(name_key)
            keep_idx.append(idx)

        new_cdf = cdf.loc[keep_idx].reset_index(drop=True)
        if len(new_cdf) == 0:
            continue

        g["children"] = new_cdf
        g["header"]["n_children"] = len(new_cdf)
        deduped_groups.append(g)

    groups = deduped_groups
    return groups

# ── 상위 Assy 행 삽입 후처리 ──────────────────────────────────
def _inject_assy_into_groups(groups, base_df):
    """
    _build_proposal_df 결과(groups)에 상위 Assy 행을 역추적해서 삽입
    1) 그룹 header(최상위 Assy) → children 첫 행
    2) 변경부품 ~ header 사이 중간 Assy → 역추적해서 삽입
    """
    if not isinstance(base_df, pd.DataFrame) or len(base_df) == 0:
        return groups

    # ── Base BOM 컬럼 탐색 ──
    pno_col    = pick_col(base_df, ["P/NO","P/NO.","P/no.","Part No","품번","부품번호","PNO","PARTNO"])
    desc_col   = pick_col(base_df, ["Description","DESC","DESC.","Desc","부품명","품명"])
    lvl_col    = pick_col(base_df, ["Lvl","LVL","Level","LEVEL","레벨"])
    parent_col = pick_col(base_df, ["Parent Part No(모)","Parent Part No","모품번"])
    qty_col    = pick_col(base_df, ["Qty","QTY","수량","개수","수량(EA)"])
    type_col   = pick_col(base_df, ["CMDT","공정","분류","Type","TYPE","유형"])

    if not pno_col:
        return groups

    # ── Base BOM 인덱스 구축 ──
    pno_to_info = {}        # norm_pno → {부품명, 품번, 레벨, 수량, 유형}
    child_to_parent = {}    # norm_child → norm_parent

    if parent_col:
        # ▸ Parent 컬럼이 있으면 직접 사용
        for _, row in base_df.iterrows():
            pn = re.sub(r'[^A-Z0-9]', '', str(row.get(pno_col, '')).strip().upper())
            if not pn:
                continue
            pno_to_info[pn] = {
                "부품명": str(row.get(desc_col, '')).strip() if desc_col else '',
                "품번":   pn,
                "레벨":   str(row.get(lvl_col, '')).strip() if lvl_col else '',
                "수량":   str(row.get(qty_col, '')).strip() if qty_col else '1',
                "유형":   str(row.get(type_col, '')).strip() if type_col else '',
            }
            pp = re.sub(r'[^A-Z0-9]', '', str(row.get(parent_col, '')).strip().upper())
            if pp and pp != pn:
                child_to_parent[pn] = pp
    else:
        # ▸ Parent 컬럼 없으면 Lvl 구조(.1, ..2, ...3)로 부모 추론
        stack = []   # [(pno, depth)]
        for _, row in base_df.iterrows():
            pn = re.sub(r'[^A-Z0-9]', '', str(row.get(pno_col, '')).strip().upper())
            if not pn:
                continue
            lvl_s = str(row.get(lvl_col, '')).strip() if lvl_col else ''
            depth = lvl_s.count('.') if lvl_s.startswith('.') else 0
            pno_to_info[pn] = {
                "부품명": str(row.get(desc_col, '')).strip() if desc_col else '',
                "품번":   pn,
                "레벨":   lvl_s,
                "수량":   str(row.get(qty_col, '')).strip() if qty_col else '1',
                "유형":   str(row.get(type_col, '')).strip() if type_col else '',
            }
            while stack and stack[-1][1] >= depth:
                stack.pop()
            if stack:
                child_to_parent[pn] = stack[-1][0]
            stack.append((pn, depth))

    # ── 각 그룹에 Assy 행 삽입 ──
    for g in groups:
        header   = g["header"]
        child_df = g["children"]
        if child_df is None or len(child_df) == 0:
            continue

        skip       = header.get("skip", False)
        header_pno = re.sub(r'[^A-Z0-9]', '', str(header.get("part_no", '')).strip().upper())

        # (1) children에서 변경 부품 품번 수집
        changed_pnos = set()
        for _, row in child_df.iterrows():
            pn = re.sub(r'[^A-Z0-9]', '', str(row.get("품번", '')).strip().upper())
            if pn:
                changed_pnos.add(pn)

        # (2) 각 변경부품 → 상위 체인 역추적
        assy_map = {}   # norm_pno → info  (중복 방지용 dict)
        for cpno in changed_pnos:
            cur = cpno
            visited = set()
            while cur in child_to_parent and cur not in visited:
                visited.add(cur)
                parent = child_to_parent[cur]
                # 이미 변경 부품에 있으면 스킵
                if parent not in changed_pnos and parent in pno_to_info:
                    assy_map[parent] = pno_to_info[parent]
                cur = parent

        # (3) header Assy 자체도 포함 (skip이 아니고, children에 없을 때)
        if not skip and header_pno and header_pno not in changed_pnos:
            if header_pno in pno_to_info:
                assy_map[header_pno] = pno_to_info[header_pno]
            else:
                assy_map[header_pno] = {
                    "부품명": header.get("desc", ''),
                    "품번":   header_pno,
                    "레벨":   header.get("lvl", ''),
                    "수량":   "1",
                    "유형":   "",
                }

        if not assy_map:
            continue

        # (4) Assy 행 생성 → depth 오름차순 정렬 (상위 먼저)
        assy_rows = []
        for apno, info in assy_map.items():
            assy_rows.append({
                "부품명":   info["부품명"],
                "품번":     info["품번"],
                "레벨":     info["레벨"],
                "수량":     info.get("수량", "1"),
                "유형":     info.get("유형", ""),
                "변경유형": "변경",
                "변경사유": "하위 부품 변경",
                "분류":     "",
                "비고":     "▲ 상위 Assy",
                "출처":     "Base BOM",
            })
        assy_rows.sort(key=lambda r: r.get("레벨", "").count('.'))

        # (5) children DataFrame 앞에 삽입
        assy_df = pd.DataFrame(assy_rows)
        for col in child_df.columns:
            if col not in assy_df.columns:
                assy_df[col] = ""
        assy_df = assy_df[child_df.columns]

        g["children"]  = pd.concat([assy_df, child_df], ignore_index=True)
        g["header"]["n_children"] = len(g["children"])

    return groups

def _build_retain_df(proposals):
    rows = []
    for p in proposals:
        for pt in (p.get("existing_parts") or []):
            rows.append({"부품명":pt.get("part_name") or "","품번":pt.get("part_no") or pt.get("display_pno") or "","레벨":pt.get("lvl") or "","수량":str(pt.get("qty") or "1"),"유형":pt.get("base_type") or ""})
    df = pd.DataFrame(rows) if rows else pd.DataFrame(columns=["부품명","품번","레벨","수량","유형"])
    if len(df) > 0:
        df["_s"] = df["레벨"].apply(lambda v: str(v).count(".") if str(v).startswith(".") else 998)
        df = df.sort_values("_s").reset_index(drop=True).drop(columns=["_s"])
    return df

import pandas as pd

def _norm_pno(x: str) -> str:
    s = "" if x is None else str(x).strip().upper()
    return re.sub(r"[^A-Z0-9]", "", s)

def _lvl_depth(v: str) -> int:
    """
    base BOM의 Lvl이 '.1', '..2' 같은 형태라고 가정하고 depth는 '.' 개수로 계산
    (값이 다르면 0 처리)
    """
    s = "" if v is None else str(v).strip()
    return s.count(".") if s.startswith(".") else 0

def extract_base_subtree(base_df: pd.DataFrame, assembly_pno: str = "", assembly_desc: str = "") -> pd.DataFrame:
    """
    Base BOM에서 assembly(상위 어셈블리) 1행을 찾고,
    그 아래 하위(children) 행들을 Lvl depth 기반으로 끝까지 잘라 반환.
    - 가장 안전한 방식: '상위 행 index' 이후로 depth가 다시 올라오다가,
      상위 depth 이하로 내려오면 subtree 종료
    """
    if not isinstance(base_df, pd.DataFrame) or len(base_df) == 0:
        return pd.DataFrame()

    # app.py에 이미 존재하는 pick_col()을 사용 (헤더 동의어 대응) [1](https://lgeteams-my.sharepoint.com/personal/heeyeong_oh_lge_com/Documents/Microsoft%20Copilot%20Chat%20%ED%8C%8C%EC%9D%BC/app.py)
    pc = pick_col(base_df, ["P/NO","P/NO.","P/no.","Part No","품번","부품번호","PNO","PARTNO"])
    dc = pick_col(base_df, ["Description","DESC","DESC.","Desc","부품명","품명"])
    lc = pick_col(base_df, ["Lvl","LVL","Level","LEVEL","레벨"])

    if not lc:
        return pd.DataFrame()

    target_pno = _norm_pno(assembly_pno)
    target_desc = (assembly_desc or "").strip().upper()

    # 1) 상위 어셈블리 행 찾기 (품번 우선, 없으면 부품명으로)
    start_idx = None
    if pc and target_pno:
        for i, row in base_df.iterrows():
            if _norm_pno(row.get(pc, "")) == target_pno:
                start_idx = i
                break

    if start_idx is None and dc and target_desc:
        for i, row in base_df.iterrows():
            desc = str(row.get(dc, "")).strip().upper()
            if desc and target_desc in desc:
                start_idx = i
                break

    if start_idx is None:
        return pd.DataFrame()

    base_depth = _lvl_depth(base_df.loc[start_idx, lc])

    # 2) 하위 subtree 수집 (상위 다음 행부터)
    rows = []
    for j in range(start_idx + 1, len(base_df)):
        d = _lvl_depth(base_df.iloc[j][lc])
        if d <= base_depth:
            break
        rows.append(base_df.iloc[j].to_dict())

    return pd.DataFrame(rows)

# -------------------------------
# 3) UI
# -------------------------------
# =========================================================
# ✅ UI 메인 (chatbot_flow 모듈 사용)
# =========================================================
st.title("개발부품Master & BOM 초안 생성 Agent")

init_review_state()
init_feedback_chat()

## ── FAST MODE (개발용) ──
if FAST_MODE:
    from dotenv import load_dotenv; load_dotenv(Path(__file__).resolve().parent / ".env")
    ss = st.session_state
    ss.setdefault("dev_grade", "B")
    ss.setdefault("target_model", "KWS9D7687M")
    ss.setdefault("change_items", ["도어에 카메라 추가"])
    ss.setdefault("ref_model", "")
    ss.setdefault("event", "DV")
    ss["_was_fast"]    = True

    # ✅ BOM 자동 로드 (EJHV 방식)
    if not ss.get("bom_uploaded") is True:
        import os
        bom_path = os.path.join(os.path.dirname(__file__), "data", "uploads", "base_bom.xlsx")
        if os.path.exists(bom_path):
            _df_raw = pd.read_excel(bom_path, header=0, dtype=str)
            from chatbot_flow import _filter_bom_base_rows, _detect_model_from_bom
            ss["base_df_raw"]  = _df_raw
            ss["base_df"]      = _filter_bom_base_rows(_df_raw)
            ss["bom_model"]    = _detect_model_from_bom(_df_raw)
            ss["bom_uploaded"] = True

## ─────────────────────────────────
## 📋 상단: 프로젝트 정보 입력 폼
## ─────────────────────────────────
with st.container():
    st.subheader("📋 프로젝트 정보")

    uploaded_bom = st.file_uploader(
        "Base BOM 업로드",
        type=["xlsx", "xls"],
        key="_input_bom_file",
    )

    # ── BOM 업로드 처리 + 모델코드 자동 추출 ──
    if uploaded_bom is not None and not st.session_state.get("bom_uploaded"):
        try:
            df_bom, hdr_row = read_excel_auto_header(uploaded_bom)
            df_bom = df_bom.dropna(how="all").fillna("")
            df_bom.columns = [str(c).strip() for c in df_bom.columns]

            st.session_state["base_df"]     = df_bom
            st.session_state["base_df_raw"] = df_bom.copy()
            st.session_state["bom_uploaded"] = True

            # D-012: bom_model 자동 추출. 첫 데이터 row(Lvl=0 또는 그 근처)에서
            # P/NO 컬럼 값을 보고 9자리 model core를 뽑음. 못 찾으면 추가 시도:
            #   1) Part No 컬럼명의 모든 variant (P/no., 품번, 부품번호 등)
            #   2) 첫 5행에서 .iloc[i]['P/NO']에 model 토큰 발견되는 첫 값
            #   3) chatbot_flow._detect_model_from_bom 백업
            pno_col = pick_col(df_bom, [
                "P/NO", "P/NO.", "P/no.", "Part No", "Part No.",
                "PartNo", "PNO", "품번", "부품번호",
            ])
            detected_model = ""
            if pno_col:
                # 첫 5행 스캔 — root row가 첫 row 아닐 수도 있음
                for i in range(min(5, len(df_bom))):
                    cand_pno = str(df_bom.iloc[i].get(pno_col, "")).strip()
                    if not cand_pno:
                        continue
                    info = extract_model_code(cand_pno)
                    if info.get("core9"):
                        detected_model = info["core9"]
                        break
                    # core9 없으면 . 앞 부분 fallback
                    if "." in cand_pno and not detected_model:
                        detected_model = cand_pno.split(".")[0]
            # 그래도 없으면 chatbot_flow의 detection 시도
            if not detected_model:
                try:
                    from chatbot_flow import _detect_model_from_bom
                    detected_model = _detect_model_from_bom(df_bom) or ""
                except Exception:
                    pass
            if detected_model:
                st.session_state["bom_model"] = detected_model
            else:
                # 디버그: 어떤 컬럼/값이 있어서 못 잡았는지 표시
                st.warning(
                    f"⚠️ Base 모델 자동 추출 실패. P/NO 컬럼 = {pno_col!r}. "
                    f"수동 입력하시거나 BOM 첫 행 점검 필요."
                )

            snap = make_base_snapshot(df_bom, uploaded_bom.name)
            st.session_state["base_snapshot"] = snap

            cur_file_key = _file_digest(uploaded_bom)
            if st.session_state.get("_structured_key") != cur_file_key and cur_file_key:
                docs = build_structured_docs_from_base(
                    df_bom,
                    st.session_state.get("bom_model", ""),
                    st.session_state.get("dev_grade", "B"),
                )
                if docs:
                    upsert_structured_docs(docs)
                    reset_structured_collection()
                    st.session_state["_structured_key"] = cur_file_key
                    st.session_state["_structured_upserted"] = True

            st.success(f"✅ BOM 업로드 완료 ({len(df_bom)}행) | 기준모델: {st.session_state.get('bom_model','')}")
        except Exception as e:
            import traceback
            tb_text = traceback.format_exc()
            st.error(f"BOM 업로드 실패: {type(e).__name__}: {e}")
            with st.expander("🔍 디버그 상세 (스택 추적)"):
                st.code(tb_text, language="python")
            # 콘솔 로그도 같이
            print(f"[BOM_UPLOAD_ERROR] {type(e).__name__}: {e}\n{tb_text}")

    st.markdown("**개발 유형 및 등급 확정 심의회 PPT 업로드**")
    ppt_up = st.file_uploader(
        "심의회 PPT 업로드 (.pptx)",
        type=["pptx"],
        key="_review_pptx_file",
    )

    # PPT 파일을 session_state에 저장해서 rerun/rerender 후에도 유지
    if ppt_up is not None:
        st.session_state["_ppt_file_uploaded"] = ppt_up
    ppt_up_final = st.session_state.get("_ppt_file_uploaded")

    if ppt_up_final is not None:
        st.caption(f"📎 업로드됨: {ppt_up_final.name}")

    uploads_ready = bool(st.session_state.get("bom_uploaded")) and bool(ppt_up_final is not None)
    review_btn = st.button(
        "🔎 모델 정보 확인",
        use_container_width=True,
        key="_review_project_inputs_btn",
        disabled=not uploads_ready,
    )
    if review_btn:
        with st.spinner("PPT 분석 중..."):
            try:
                ppt_bytes = ppt_up_final.getvalue()
                ctx = {
                    "source_model": st.session_state.get("bom_model", ""),
                    "target_model": st.session_state.get("target_model", ""),
                    "key_diff": "",
                    "dev_grade": st.session_state.get("dev_grade", "B"),
                    "product_type": st.session_state.get("product_type", ""),
                }
                ppt_result = extract_change_review_from_pptx_bytes(ppt_bytes, ctx=ctx)
                st.session_state["ppt_extraction_result"] = ppt_result
                st.session_state["ppt_extraction_name"] = ppt_up_final.name
                _pm = (ppt_result or {}).get("project_meta") or {}
                st.session_state["dev_grade"] = str(_pm.get("dev_grade") or st.session_state.get("dev_grade", "B") or "B").strip().upper()

                # D-012 fix: 원본은 PPT에서 base_model/target_model 추출해도 session_state에
                # 안 채워서 화면 input이 비어 보임. 추출값을 명시적으로 wire-up.
                _ppt_base = str(_pm.get("base_model") or _pm.get("source_model") or "").strip()
                _ppt_target = str(_pm.get("target_model") or "").strip()
                if _ppt_base and not st.session_state.get("bom_model"):
                    st.session_state["bom_model"] = _ppt_base
                if _ppt_target:
                    st.session_state["target_model"] = _ppt_target
                # 부가 메타도 같이 wire-up (있으면)
                for _k in ("product_type", "target_country", "rating"):
                    _v = str(_pm.get(_k) or "").strip()
                    if _v:
                        st.session_state[_k] = _v

                st.session_state["_show_project_input_review"] = True
                st.success(
                    f"✅ 모델 정보 확인 완료 — Base: {st.session_state.get('bom_model','')} / "
                    f"Target: {st.session_state.get('target_model','(미추출)')}"
                )
            except Exception as e:
                st.error(f"⚠️ 모델 정보 확인 실패(PPT 추출): {e}")

    fields_unlocked = bool(st.session_state.get("_show_project_input_review"))
    if not fields_unlocked:
        st.info("Base BOM + 심의회 PPT 업로드 후 `모델 정보 확인`을 누르면 PPT 추출과 입력 정보 확인이 함께 진행됩니다.")

    # ── 모델 정보 ──
    row2_c1, row2_c2, row2_c3 = st.columns([1, 1, 1])

    with row2_c1:
        base_model_display = st.session_state.get("bom_model", "(BOM 업로드 필요)")
        st.text_input("기준모델 (Base)", value=base_model_display, disabled=True, key="_disp_base_model")

    with row2_c2:
        target_val = st.session_state.get("target_model", "")
        # D-012: 신규모델은 PPT 추출이 디자인 의도상 안 되므로 사용자 직접 입력.
        # BOM 업로드만으로도 입력 가능하도록 PPT 확인 의존성 제거.
        new_model = st.text_input(
            "신규모델",
            value=target_val,
            key="_input_target_model",
            placeholder="예: WSED7613B (선택 — 비워도 분석 가능)",
        )
        st.session_state["target_model"] = new_model

    with row2_c3:
        event_options = ["CP", "DV", "PV", "PreMP"]
        cur_event = st.session_state.get("event", "CP")
        ev_idx = event_options.index(cur_event) if cur_event in event_options else 0
        selected_event = st.selectbox("Event", event_options, index=ev_idx, key="_input_event", disabled=not fields_unlocked)
        st.session_state["event"] = selected_event

    # v4.0: region은 프로젝트 메타(국가/정격)에서 자동 설정
    region = _infer_project_region()
    st.session_state["region"] = region

    if st.session_state.get("ppt_extraction_result"):
        _ppt_res = st.session_state.get("ppt_extraction_result") or {}
        _pm = _ppt_res.get("project_meta") or {}
        _cp = _ppt_res.get("change_points_for_bom") or []
        _md = _ppt_res.get("module_details") or []

        with st.expander("✅ PPT 추출값 확인/수정", expanded=True):
            c1, c2 = st.columns(2)
            with c1:
                _product_type = st.text_input(
                    "제품군",
                    value=str(_pm.get("product_type") or st.session_state.get("product_type", "")),
                    key="_ppt_edit_product_type",
                )
                _base_model = st.text_input(
                    "Base 모델",
                    value=str(_pm.get("base_model") or st.session_state.get("bom_model", "")),
                    key="_ppt_edit_base_model",
                )
                _target_model = st.text_input(
                    "대상 모델",
                    value=str(st.session_state.get("target_model", "")),
                    key="_ppt_edit_target_model",
                )
            with c2:
                _grade_raw = str(_pm.get("dev_grade") or st.session_state.get("dev_grade", "B")).strip().upper()
                _grade_norm_map = {"CA": "CA", "CB": "CB", "CSW": "CSW", "S": "S", "A": "A", "B": "B", "D": "D"}
                _grade_pick = _grade_norm_map.get(_grade_raw, "B")
                _grade = st.text_input(
                    "개발등급",
                    value=_grade_pick,
                    key="_ppt_edit_dev_grade",
                    disabled=True,
                )
                _target_country = st.text_input(
                    "출시 국가",
                    value=str(_pm.get("target_country") or st.session_state.get("target_country", "")),
                    key="_ppt_edit_target_country",
                )
                _rating = st.text_input(
                    "정격",
                    value=str(_pm.get("rating") or st.session_state.get("rating", "")),
                    key="_ppt_edit_rating",
                )

            def _is_noise_change_point(_desc: Any, _module: Any) -> bool:
                mod = str(_module or "")
                desc = re.sub(r"\s+", "", str(_desc or "").upper())
                if "과전압" in mod:
                    return True
                if desc in {"O", "X", "OX", "(O,X)", "(O, X)"}:
                    return True
                return False

            def _infer_ui_change_type(_txt: Any) -> str:
                s = str(_txt or "").upper()
                if any(k in s for k in ["삭제", "제거", "미적용", "DELETED", "REMOVE"]):
                    return "삭제"
                if any(k in s for k in ["추가", "신규", "ADD", "NEW"]):
                    return "NEW"
                return "Changing"

            _cp_rows = []
            for i, r in enumerate(_cp, 1):
                _row = {
                    # PPT의 번호는 ① 같은 형식일 수 있어 문자열로 보존한다.
                    "No": str(r.get("id") or i),
                    "변경점": str(r.get("description") or "").strip(),
                    "상위 Ass'y": str(r.get("module") or "").strip(),
                    "분야": str(r.get("discipline") or "").strip(),
                    "변경 유형": str(r.get("type") or "").strip(),
                }
                if _is_noise_change_point(_row.get("변경점"), _row.get("상위 Ass'y")):
                    continue
                _cp_rows.append(_row)

            # 보조 안전망: cp가 비었는데 상세표는 있는 경우, 사용자 편집용 임시 행 제공
            if not _cp_rows and _md:
                for i, md in enumerate(_md, 1):
                    part = str((md.get("new") or {}).get("description") or "").strip()
                    if not part:
                        part = str((md.get("base") or {}).get("description") or "").strip()
                    _row = {
                        "No": str(i),
                        "변경점": part,
                        "상위 Ass'y": str(md.get("module") or "").strip(),
                        "분야": "",
                        "변경 유형": _infer_ui_change_type(part),
                    }
                    if _is_noise_change_point(_row.get("변경점"), _row.get("상위 Ass'y")):
                        continue
                    _cp_rows.append(_row)
                if _cp_rows:
                    st.caption("참고: MAIN 변경표 추출이 비어 상세표 기반 임시 변경점으로 채웠습니다. 필요 시 수정해 주세요.")
            if not _cp_rows:
                _cp_rows = [{"No": "1", "변경점": "", "상위 Ass'y": "", "분야": "", "변경 유형": ""}]

            _cp_df = pd.DataFrame(_cp_rows)
            if len(_cp_df) > 0:
                _cp_df = _cp_df[
                    ~_cp_df.apply(
                        lambda rr: _is_noise_change_point(rr.get("변경점"), rr.get("상위 Ass'y")),
                        axis=1,
                    )
                ].reset_index(drop=True)
            _cp_edit = st.data_editor(
                _cp_df,
                use_container_width=True,
                hide_index=True,
                num_rows="dynamic",
                key="_ppt_change_points_editor",
            )

            apply_ppt_btn = st.button(
                "✅ 확인 후 입력값 반영",
                type="primary",
                use_container_width=True,
                key="_apply_ppt_inputs_btn",
            )

            if apply_ppt_btn:
                _applied_changes = [
                    str(v).strip()
                    for v in (_cp_edit.get("변경점", pd.Series(dtype=str)).tolist() if isinstance(_cp_edit, pd.DataFrame) else [])
                    if str(v or "").strip()
                ]
                if not _applied_changes:
                    st.warning("변경점이 비어 있습니다. 최소 1개 이상 확인/입력해 주세요.")
                else:
                    st.session_state["product_type"] = str(_product_type or "").strip()
                    st.session_state["bom_model"] = str(_base_model or "").strip() or st.session_state.get("bom_model", "")
                    st.session_state["target_model"] = str(_target_model or "").strip() or st.session_state.get("target_model", "")
                    st.session_state["dev_grade"] = str(_grade or "B").strip().upper()
                    st.session_state["target_country"] = str(_target_country or "").strip()
                    st.session_state["rating"] = str(_rating or "").strip()
                    st.session_state["change_items"] = _applied_changes
                    st.session_state["ppt_input_confirmed"] = True
                    st.success("✅ PPT 추출값 반영 완료. 필요하면 변경점을 추가/수정 후 분석 시작하세요.")
                    st.rerun()

            if st.session_state.get("ppt_input_confirmed"):
                st.success("PPT 추출값 반영 상태입니다. 수동 입력 없이 분석 시작할 수 있습니다.")

        with st.expander("추출 원본 JSON", expanded=False):
            st.json(_ppt_res)

        import datetime as _dt
        ts_ppt = _dt.datetime.now().strftime("%Y%m%d_%H%M")
        out_name = f"PPT_Change_Extraction_{ts_ppt}.json"
        st.download_button(
            label="⬇️ PPT 추출결과 JSON 다운로드",
            data=json.dumps(_ppt_res, ensure_ascii=False, indent=2).encode("utf-8"),
            file_name=out_name,
            mime="application/json",
            use_container_width=True,
            key="_download_ppt_extract_json",
        )

    if st.session_state.get("_show_project_input_review"):
        st.markdown("**입력 정보 요약**")
        _changes = st.session_state.get("change_items") or []
        _changes = [str(x).strip() for x in _changes if str(x or "").strip()]

        info_rows = [
            {"항목": "기준모델(Base)", "값": str(st.session_state.get("bom_model") or "(미확인)")},
            {"항목": "신규모델", "값": str(st.session_state.get("target_model") or "(미입력)")},
            {"항목": "개발등급", "값": str(st.session_state.get("dev_grade") or "(미입력)")},
            {"항목": "Event", "값": str(st.session_state.get("event") or "(미입력)")},
            {"항목": "제품군", "값": str(st.session_state.get("product_type") or "(미입력)")},
            {"항목": "출시 국가", "값": str(st.session_state.get("target_country") or "(미입력)")},
            {"항목": "정격", "값": str(st.session_state.get("rating") or "(미입력)")},
            {"항목": "Base BOM 업로드", "값": "완료" if st.session_state.get("bom_uploaded") else "미완료"},
            {"항목": "PPT 반영", "값": "완료" if st.session_state.get("ppt_input_confirmed") else "미완료"},
            {"항목": "변경점 개수", "값": str(len(_changes))},
        ]
        st.dataframe(pd.DataFrame(info_rows), use_container_width=True, hide_index=True)

        if _changes:
            st.write("**변경점 목록**")
            st.dataframe(
                pd.DataFrame({"No": list(range(1, len(_changes) + 1)), "변경점": _changes}),
                use_container_width=True,
                hide_index=True,
            )
        else:
            st.caption("변경점이 아직 없습니다. PPT 추출 결과를 확인해 주세요.")

    # ── 🚀 분석 시작 버튼 ──
    st.divider()

    # D-012: 분석 시작 조건 완화 — PPT 확인 + target_model 둘 다 optional로.
    # 최소 필수는 BOM 업로드 + 변경점 입력. PPT/target_model 있으면 정확도↑.
    can_start = (
        st.session_state.get("bom_uploaded")
        and len(st.session_state.get("change_items", [])) > 0
    )

    run_col1, run_col2 = st.columns([3, 1])
    with run_col1:
        run_btn = st.button(
            "🚀 분석 시작",
            type="primary",
            use_container_width=True,
            disabled=not can_start,
        )
    with run_col2:
        rerun_btn = st.button(
            "🔄 재분석",
            use_container_width=True,
            disabled=not st.session_state.get("proposals"),
        )

    if not can_start and not FAST_MODE:
        st.info("Base BOM/PPT 업로드 후 `모델 정보 확인`으로 입력값을 점검하고 분석을 시작하세요.")

        # ── 분석 실행 (★ EJHV 방식: compare_all_proposals 없이 바로 저장) ──
    if run_btn: # ⚡️ 병목 해결: 변경점 일괄 처리로 변경
        ss = st.session_state
        import time as _time
        analysis_t0 = _time.perf_counter()

        change_items = ss.get("change_items", [])
        target_model = ss.get("target_model", "")
        dev_grade = ss.get("dev_grade", "B")

        ss["related_parts_analysis_prompt"] = build_related_parts_analysis_prompt(
            change_items=change_items,
            source_model=ss.get("bom_model", ""),
            target_model=target_model,
            key_diff=ss.get("key_diff", ""),
            product_type=ss.get("product_type", ""),
        )

        _bar = st.progress(0, text="⏳ 준비 중...")
        def _set(target, text):
            _bar.progress(int(target), text=text)

        step_sec = {"snapshot": 0.0, "structured_db": 0.0, "search": 0.0, "proposal": 0.0}

        _set(15, "📦 Base BOM 스냅샷 구축 중...")
        step1_t0 = _time.perf_counter()
        if not ss.get("base_snapshot"):
            _bdf = ss.get("base_df")
            if isinstance(_bdf, pd.DataFrame) and len(_bdf) > 0:
                ss["base_snapshot"] = make_base_snapshot(_bdf, "base_bom.xlsx")
        step_sec["snapshot"] = round(_time.perf_counter() - step1_t0, 3)

        _set(30, "🗄️ 구조화 DB 구축 중...")
        step2_t0 = _time.perf_counter()
        if not (ss.get("_structured_upserted") or ss.get("_structured_key")):
            _bdf = ss.get("base_df")
            if isinstance(_bdf, pd.DataFrame) and len(_bdf) > 0:
                try:
                    docs = build_structured_docs_from_base(
                        _bdf, ss.get("bom_model",""), ss.get("dev_grade","B"))
                    if docs:
                        upsert_structured_docs(docs)
                        reset_structured_collection()
                        ss["_structured_upserted"] = True
                        ss["_structured_key"] = f"fallback:{len(_bdf)}:{ss.get('bom_model','')}"
                except Exception:
                    pass
        step_sec["structured_db"] = round(_time.perf_counter() - step2_t0, 3)

        _set(60, f"🔍 변경점 {len(change_items)}건 일괄 검색 중...")
        search_t0 = _time.perf_counter()
        search_result = run_search(change_items, target_model, dev_grade)
        step_sec["search"] = round(_time.perf_counter() - search_t0, 3)

        _set(85, "🧩 제안 생성 및 병합 중...")
        proposal_t0 = _time.perf_counter()
        proposals_raw = generate_proposals_from_docs(
            search_result.get("primary_docs", []),
            ss.get("base_snapshot") or {},
            change_items,
        )
        all_proposals = merge_proposals_order_independent(proposals_raw)
        step_sec["proposal"] = round(_time.perf_counter() - proposal_t0, 3)

        # D-012 fix: 디버그 사이드바가 볼 수 있게 session_state에 저장
        _raw_primary = search_result.get("primary_docs") or []
        _raw_secondary = search_result.get("secondary_docs") or []
        _raw_sdbg = search_result.get("search_debug") or {}

        # D-012 진단: stdout이 안 잡혀서 파일에 직접 append
        try:
            from pathlib import Path as _DbgPath
            _dbg_path = _DbgPath(__file__).resolve().parent / "_dbg.log"
            with open(_dbg_path, "a", encoding="utf-8") as _df:
                from datetime import datetime as _Dt
                _df.write(
                    f"[{_Dt.now().isoformat(timespec='seconds')}] "
                    f"primary_at_assign={len(_raw_primary)} "
                    f"secondary_at_assign={len(_raw_secondary)} "
                    f"proposals_at_assign={len(all_proposals or [])} "
                    f"sdbg={_raw_sdbg}\n"
                )
        except Exception:
            pass

        ss["primary_docs"] = _raw_primary
        ss["secondary_docs"] = _raw_secondary
        ss["proposals"] = all_proposals or []
        ss["search_debug"] = _raw_sdbg
        # 진단용 — 저장 시점에 실제 list 길이를 별도 키로도 박아둠
        ss["_dbg_primary_at_assign"] = len(_raw_primary)
        ss["_dbg_secondary_at_assign"] = len(_raw_secondary)
        ss["_dbg_proposals_at_assign"] = len(all_proposals or [])

        sdbg = search_result.get("search_debug") or {}
        timing_logs = [{
            "change_item": "일괄 처리",
            "search_sec": step_sec["search"],
            "proposal_sec": step_sec["proposal"],
            "total_sec": step_sec["search"] + step_sec["proposal"],
            "primary_docs": len(search_result.get("primary_docs") or []),
            "proposals_generated": len(all_proposals or []),
            "docs_pre_sim": int(sdbg.get("pre_sim_docs") or 0),
            "docs_post_sim": int(sdbg.get("post_sim_docs") or 0),
            "sim_threshold": float(sdbg.get("sim_threshold") or 0.0),
        }]

        ss["primary_docs"]   = []
        ss["secondary_docs"] = []
        ss["proposals"]      = all_proposals
        ss["analysis_timing_logs"] = timing_logs
        ss["analysis_step_sec"] = step_sec
        ss["analysis_total_sec"] = round(_time.perf_counter() - analysis_t0, 3)

        # ── 완료 ──
        _set(100, "✅ 완료!")
        _time.sleep(0.8)
        _bar.empty()

        if all_proposals:
            st.success(f"✅ {len(all_proposals)}건 제안 생성 완료! (변경점 {len(change_items)}개 일괄 처리)")
        else:
            st.warning("⚠️ 추천 결과가 없습니다.")

        if timing_logs:
            st.caption(f"⏱️ 분석 소요시간(전체): {ss.get('analysis_total_sec', 0.0)}초")
            st.caption(
                f"단계별: snapshot {step_sec['snapshot']}s | "
                f"structured_db {step_sec['structured_db']}s | "
                f"search {step_sec['search']}s | proposal {step_sec['proposal']}s"
            )
            with st.expander("🕒 분석 소요시간 상세 로그", expanded=False):
                st.dataframe(pd.DataFrame(timing_logs), use_container_width=True, hide_index=True)

    if rerun_btn:
        st.session_state["proposals"] = []
        st.session_state["decisions"] = {}
        st.session_state["analysis_timing_logs"] = []
        st.session_state["analysis_total_sec"] = 0.0
        st.session_state["analysis_step_sec"] = {}
        st.rerun()

    saved_timing = st.session_state.get("analysis_timing_logs") or []
    if saved_timing and not run_btn:
        _step = st.session_state.get("analysis_step_sec") or {}
        st.caption(f"⏱️ 최근 분석 소요시간(전체): {st.session_state.get('analysis_total_sec', 0.0)}초")
        if _step:
            st.caption(
                f"단계별: snapshot {_step.get('snapshot', 0.0)}s | "
                f"structured_db {_step.get('structured_db', 0.0)}s | search {_step.get('search', 0.0)}s | "
                f"proposal {_step.get('proposal', 0.0)}s"
            )
        with st.expander("🕒 분석 소요시간 상세 로그 (최근 실행)", expanded=False):
            st.dataframe(pd.DataFrame(saved_timing), use_container_width=True, hide_index=True)

## ─────────────────────────────────
## 🧩 하단: 결과 (제안 검토 + 추출)
## ─────────────────────────────────
st.divider()

st.subheader("🧩 제안 검토")
all_edited = {}
if st.session_state.get("pending_recheck"):
    process_pending_recheck_request()
prepare_active_review()
proposals = st.session_state.get("proposals") or []
if not proposals:
    st.info("아직 추천이 없습니다.")
else:
    base_df = st.session_state.get("base_df")
    groups = _build_proposal_df(proposals, base_df)
    total_parts = sum(g["header"]["n_children"] for g in groups)
    TYPE_OPTIONS = ["추가", "변경(시방)", "삭제", "제외", "⚠️확인필요"]
    st.caption(f"추천 부품 총 **{total_parts}건** | 변경유형을 선택해서 확정해 주세요")
    all_edited = {}
    for g_idx, g in enumerate(groups):
        header = g["header"]
        child_df = g["children"]
        desc = header["desc"]
        pno = header["part_no"]
        lvl = header["lvl"]
        n = header["n_children"]
        skip = header["skip"]
        n_unresolved = int((child_df["변경유형"] == "⚠️확인필요").sum()) if (isinstance(child_df, pd.DataFrame) and "변경유형" in child_df.columns) else 0
        pno_display = f" | {pno}" if pno else ""
        if n_unresolved > 0:
            label = f"🔴 ■ {desc} ({n}건 | ⚠️ {n_unresolved}건 미확정){pno_display} | Lvl {lvl}"
            expanded_default = True
        else:
            label = f"✅ ■ {desc} ({n}건){pno_display} | Lvl {lvl}"
            expanded_default = (g_idx == 0)

        with st.expander(label, expanded=expanded_default):
            if not skip and pno:
                st.markdown(f"**상위 어셈블리**: `{desc}` | `{pno}` | 출처: {header['source']}")
            if len(child_df) == 0:
                st.info("하위 변경 부품이 없습니다.")
            else:
                edited = st.data_editor(child_df, use_container_width=True, hide_index=True, column_config={"변경유형":st.column_config.SelectboxColumn("변경유형",options=TYPE_OPTIONS,width="small",required=True),"부품명":st.column_config.TextColumn("부품명",width="large"),"품번":st.column_config.TextColumn("품번",width="medium"),"레벨":st.column_config.TextColumn("레벨",width="small"),"수량":st.column_config.TextColumn("수량",width="small"),"유형":st.column_config.TextColumn("유형",width="small"),"변경사유":st.column_config.TextColumn("변경사유",width="large"),"분류":st.column_config.TextColumn("분류",width="small"),"비고":st.column_config.TextColumn("비고",width="large"),"출처":st.column_config.TextColumn("출처",width="medium")}, disabled=["부품명","품번","레벨","수량","유형","변경사유","분류","비고","출처"], key=f"proposal_editor_{g_idx}")
                all_edited[g_idx] = edited

                unresolved_rows = edited[edited["변경유형"] == "⚠️확인필요"] if "변경유형" in edited.columns else pd.DataFrame()
                if len(unresolved_rows) > 0:
                    st.warning(f"⚠️ 확인필요 {len(unresolved_rows)}건 - 변경유형을 확정해 주세요:")
                    for row_idx, (_, row) in enumerate(unresolved_rows.iterrows(), 1):
                        pno_v = str(row.get("품번", "")).strip()
                        pno_display2 = pno_v if pno_v and pno_v.lower() != "nan" else "(품번 미확정)"
                        name_v = str(row.get("부품명", "")).strip()
                        st.caption(f"  {row_idx}. {name_v} ({pno_display2})")

                show_keep = st.toggle(
                    "🔎 하위 전체(유지 부품) 보기",
                    value=False,
                    key=f"toggle_keep_{g_idx}",
                    help="가독성을 위해 기본 숨김. 켜면 Base BOM 기준으로 해당 어셈블리 하위 전체를 '유지'로 보여줍니다."
                )

                if show_keep:
                    if not isinstance(base_df, pd.DataFrame) or len(base_df) == 0:
                        st.info("Base BOM이 없어 하위 전체(유지)를 표시할 수 없습니다.")
                    else:
                        subtree = extract_base_subtree(base_df, assembly_pno=pno, assembly_desc=desc)
                        if len(subtree) == 0:
                            st.warning("Base BOM에서 해당 어셈블리 하위를 찾지 못했습니다.")
                        else:
                            st.dataframe(subtree, use_container_width=True, hide_index=True)

## ── 🆕 플로팅 피드백 챗봇 ──
if all_edited:
    # 현재 화면에서 편집된 값을 우선 세션에 저장한다.
    st.session_state["proposal_review_all"] = all_edited
else:
    # rerun 직후에는 세션의 마지막 편집본을 사용한다.
    all_edited = st.session_state.get("proposal_review_all") or {}

if all_edited:
        all_rows = pd.concat(all_edited.values(), ignore_index=True)
        n_add = len(all_rows[all_rows["변경유형"]=="추가"])
        n_mod = len(all_rows[all_rows["변경유형"]=="변경(시방)"])
        n_del = len(all_rows[all_rows["변경유형"]=="삭제"])
        n_exc = len(all_rows[all_rows["변경유형"]=="제외"])
        n_chk = len(all_rows[all_rows["변경유형"]=="⚠️확인필요"])
        st.write(f"📊 추가 **{n_add}** | 변경(시방) **{n_mod}** | 삭제 **{n_del}** | 제외 **{n_exc}** | ⚠️확인필요 **{n_chk}**")
        if n_chk > 0:
            st.warning(f"⚠️ 확인필요 **{n_chk}건**이 남아있습니다. 변경유형을 확정해 주세요!")

        # ───────────────────────────────────────
        # PATCH 2: 변경부품리스트 Excel 추출
        # ───────────────────────────────────────
        from io import BytesIO
        import datetime as _dt

        def _classify_part_type(type_val, tech_spec):
            """Base BOM의 Type + Technical Spec → 개발부품마스터 Part Type"""
            t = str(type_val or "").strip().upper()
            spec = str(tech_spec or "").strip().upper()
            if "ASSEMBLY" in t:
                return "Assy"
            if t in ("PCBASSEMBLYPART", "PCBPART"):
                return "회로"
            if t == "CIRCUITCOMPONENTPART":
                return "전장"
            if t == "MATERIALPART":
                return "원재료"
            if t in ("LABELPART", "MANUALPART", "BOXPART"):
                return "인쇄물"
            if t == "FASTENERPART":
                return "체결"
            if t == "SOFTWAREPART":
                return "S/W"
            if spec.startswith("MOLD"):
                return "사출"
            if spec.startswith("PRESS") or spec.startswith("BENDING") or spec.startswith("COMPLEX"):
                return "판금"
            if spec.startswith("CUTTING"):
                return "절삭"
            if spec.startswith("EXTRUSION"):
                return "압출"
            if spec.startswith("DIECASTING"):
                return "주조"
            return "기타"

        def _classify_change_type(chg_type, note, group_has_add=False):
            c = str(chg_type or "").strip()
            n = str(note or "").strip()
            if c == "추가":
                return "NEW"
            if c == "삭제" or c.upper() == "DELETE":
                return "삭제"
            if "▲ 상위 Assy" in n:
                return "NEW" if group_has_add else "Changing"
            if c == "변경(시방)":
                return "Changing"
            return "Common"
        
        def _make_change_parts_list(all_rows, ss):
            df = all_rows.copy()
            df = df[df["변경유형"].isin(["추가", "변경(시방)", "삭제"])].copy()
            if "skeleton_order" in df.columns:
                df["_sort_skel"] = pd.to_numeric(df["skeleton_order"], errors="coerce").fillna(999999).astype(int)
                df = df.sort_values(["_sort_skel"], kind="stable").reset_index(drop=True)
                df = df.drop(columns=["_sort_skel"], errors="ignore")
            base_df = ss.get("base_df")
            if isinstance(base_df, pd.DataFrame) and len(base_df) > 0:
                bom_cols = [c for c in base_df.columns if not str(c).startswith("Unnamed")]
            else:
                bom_cols = []
            MAPPING = {
                "부품명": ["Description","DESC","DESC.","Desc","부품명","품명","Part Name","Part Name(자)"],
                "품번":   ["P/NO","P/NO.","P/no.","Part No","품번","부품번호","PNO","PARTNO"],
                "레벨":   ["Lvl","LVL","Level","LEVEL","레벨"],
                "수량":   ["Qty","QTY","수량","개수","수량(EA)"],
                "유형":   ["CMDT","공정","분류","Type","TYPE","유형"],
            }
            def _find_our_col(bom_col_name):
                n = bom_col_name.strip().upper()
                for our_col, synonyms in MAPPING.items():
                    for s in synonyms:
                        if s.strip().upper() == n:
                            return our_col
                return None
            if bom_cols:
                result_rows = []
                for _, row in df.iterrows():
                    new_row = {}
                    for bc in bom_cols:
                        our = _find_our_col(bc)
                        if our and our in row.index:
                            new_row[bc] = row[our]
                        else:
                            new_row[bc] = ""
                    result_rows.append(new_row)
                result = pd.DataFrame(result_rows, columns=bom_cols)
            else:
                result = df[["변경유형","부품명","품번","레벨","수량","유형","변경사유","비고","출처"]].copy()
            result["변경유형"] = df["변경유형"].values
            result["변경사유"] = df["변경사유"].values if "변경사유" in df.columns else ""
            result["비고"] = df["비고"].values if "비고" in df.columns else ""
            result["출처"] = df["출처"].values if "출처" in df.columns else ""
            return result

        def _to_excel_bytes(df, sheet_name="ChangePartsList"):
            bio = BytesIO()
            with pd.ExcelWriter(bio, engine="openpyxl") as writer:
                df.to_excel(writer, index=False, sheet_name=sheet_name, startcol=1)
            return bio.getvalue()
        
        def _build_ref_bom_index():
            from pathlib import Path
            from enrich import _norm_pno, COL_SYNONYMS, _pick_col
            ref_dir = Path(__file__).resolve().parent / "data" / "ref_boms"
            index = {}
            if not ref_dir.exists():
                st.warning(f"참고 BOM 폴더 없음: {ref_dir}")
                return index
            def _process_file(f, engine=None):
                try:
                    kwargs = {"header": 0, "dtype": str}
                    if engine:
                        kwargs["engine"] = engine
                    df_ref = pd.read_excel(f, **kwargs)
                    df_ref = df_ref.dropna(how="all").fillna("")
                    df_ref.columns = [str(c).strip() for c in df_ref.columns]
                    pno_col = _pick_col(df_ref, "part_no")
                    if not pno_col:
                        return
                    col_map = {}
                    for std_key in COL_SYNONYMS:
                        found = _pick_col(df_ref, std_key)
                        if found:
                            col_map[std_key] = found
                    source_file = f.name
                    for _, row in df_ref.iterrows():
                        raw_pno = str(row.get(pno_col, "")).strip()
                        pno_key = _norm_pno(raw_pno)
                        if not pno_key or len(pno_key) < 5:
                            continue
                        if pno_key in index:
                            continue
                        entry = {"_source_file": source_file}
                        for std_key, real_col in col_map.items():
                            val = str(row.get(real_col, "")).strip()
                            if val and val.lower() not in ("nan", "none", ""):
                                entry[std_key] = val
                        index[pno_key] = entry
                except Exception as e:
                    st.warning(f"참고 BOM 로드 실패 [{f.name}]: {e}")
            for f in ref_dir.glob("*.xlsx"):
                _process_file(f)
            for f in ref_dir.glob("*.xls"):
                _process_file(f, engine="xlrd")
            return index
        

        def _make_dev_parts_master(change_df, ss, db_index=None):
            """변경부품리스트 → 개발부품 마스터 DataFrame"""
            if change_df is None or len(change_df) == 0:
                return pd.DataFrame()

            if "skeleton_order" in change_df.columns:
                _tmp_skel = pd.to_numeric(change_df["skeleton_order"], errors="coerce").fillna(999999).astype(int)
                change_df = change_df.assign(_sort_skel=_tmp_skel).sort_values(["_sort_skel"], kind="stable").drop(columns=["_sort_skel"], errors="ignore").reset_index(drop=True)

            base_raw = ss.get("base_df_raw")

            # ── Base BOM 인덱스 ──
            base_map = {}
            if isinstance(base_raw, pd.DataFrame) and len(base_raw) > 0:
                _bpno = pick_col(base_raw, ["P/NO","P/NO.","Part No","품번"])
                _btype = pick_col(base_raw, ["Type","TYPE","유형"])
                _bspec = pick_col(base_raw, ["Technical Spec","TECHNICAL SPEC","기술규격"])
                _bmaker = pick_col(base_raw, ["Maker","MAKER","제조사"])
                _bsupply = pick_col(base_raw, ["Supply Type"])
                _buit = pick_col(base_raw, ["UIT"])
                _bqty = pick_col(base_raw, ["Qty","QTY","수량"])
                if _bpno:
                    for _, row in base_raw.iterrows():
                        pn = re.sub(r'[^A-Z0-9]','',str(row.get(_bpno,'')).strip().upper())
                        if pn:
                            base_map[pn] = {
                                "type": str(row.get(_btype,'')).strip() if _btype else '',
                                "spec": str(row.get(_bspec,'')).strip() if _bspec else '',
                                "maker": str(row.get(_bmaker,'')).strip() if _bmaker else '',
                                "supply": str(row.get(_bsupply,'')).strip() if _bsupply else '',
                                "uit": str(row.get(_buit,'')).strip().upper() if _buit else '',
                                "qty": str(row.get(_bqty,'')).strip() if _bqty else '1',
                            }

            # ── 컬럼 탐색 ──
            c_pno  = pick_col(change_df, ["P/NO","P/NO.","Part No","품번"]) or "Part No"
            c_desc = pick_col(change_df, ["Description","DESC","부품명","Part Name(자)"]) or "Description"
            # Prefer change-list level columns first to preserve toggle-style level strings.
            c_lvl  = pick_col(change_df, ["레벨","Lvl","LVL","Level","LEVEL"]) or "Lvl"
            c_qty  = pick_col(change_df, ["Qty","QTY","수량"]) or "Qty"
            c_spec = pick_col(change_df, ["Technical Spec","기술규격"]) or "Technical Spec"
            c_type_col = pick_col(change_df, ["Type","TYPE","유형"])

            # ── 그룹별 "추가" 여부 판단 ──
            # _group_parent_pno가 있으면 그룹 기준, 없으면 전체 기준
            group_has_add_map = {}
            if "_group_parent_pno" in change_df.columns and "변경유형" in change_df.columns:
                for gpno, gdf in change_df.groupby("_group_parent_pno", sort=False):
                    has_add = (gdf["변경유형"] == "추가").any()
                    group_has_add_map[str(gpno)] = has_add
            else:
                # 전체에 추가가 있으면 True
                global_has_add = ("변경유형" in change_df.columns and
                                  (change_df["변경유형"] == "추가").any())
                group_has_add_map["__all__"] = global_has_add

            # ── 참고품번 추출 패턴 ──
            ref_pno_re = re.compile(r"참고[:：]?\s*([A-Z0-9]{5,})", re.I)

            rows = []
            for idx, (_, row) in enumerate(change_df.iterrows(), 1):
                def _clean(v: Any) -> str:
                    s = str(v or "").strip()
                    return "" if s.lower() in ("nan", "none") else s

                pno_raw = str(row.get(c_pno, '')).strip()
                pno_norm = re.sub(r'[^A-Z0-9]','', pno_raw.upper())
                desc = _clean(row.get(c_desc, ''))
                lvl = _clean(row.get(c_lvl, ''))
                qty_new = _clean(row.get(c_qty, '1'))
                spec_val = _clean(row.get(c_spec, ''))
                chg_type = _clean(row.get("변경유형", '')) if "변경유형" in change_df.columns else ""
                chg_reason = _clean(row.get("변경사유", '')) if "변경사유" in change_df.columns else ""
                note = _clean(row.get("비고", '')) if "비고" in change_df.columns else ""
                gpno = str(row.get("_group_parent_pno", '')).strip() if "_group_parent_pno" in change_df.columns else "__all__"

                # ── ✅ 핵심: Base P/No 결정 ──
                # 1순위: 이 부품 자체가 Base BOM에 있으면 → 자기 자신이 Base
                # 2순위: 비고의 참고품번(DB의 new_part_no) → 우리의 Base P/No
                # 3순위: 없으면 빈칸
                base_pno = ""
                if pno_norm and pno_norm in base_map:
                    # 1순위: Base BOM에 있는 품번 → 자기가 Base
                    base_pno = pno_raw
                else:
                    # 2순위: 비고의 참고품번
                    ref_match = ref_pno_re.search(note)
                    if ref_match:
                        base_pno = ref_match.group(1)
                    else:
                        # 3순위: ref_boms(출처 파일)에서 부품명으로 역조회
                        desc_norm = re.sub(r'[^A-Z0-9가-힣]','', desc.upper())
                        if desc_norm and db_index:
                            for rpno, rentry in db_index.items():
                                r_name = re.sub(r'[^A-Z0-9가-힣]','',
                                    str(rentry.get("part_name","")).upper())
                                if r_name == desc_norm:
                                    # 원본 품번 → 우리의 Base P/No
                                    base_pno = rentry.get("part_no", rpno)
                                    break

                # ── Base 정보 가져오기 (Base P/No 기준) ──
                base_pno_norm = re.sub(r'[^A-Z0-9]','', base_pno.upper()) if base_pno else ""
                bi = base_map.get(pno_norm, {}) or base_map.get(base_pno_norm, {})

                type_raw = str(row.get(c_type_col, '')).strip() if c_type_col else bi.get("type", "")
                if not spec_val:
                    spec_val = bi.get("spec", "")

                part_type = _classify_part_type(type_raw, spec_val)

                # ── 분류: 그룹 내 추가 여부 반영 ──
                group_has_add = group_has_add_map.get(gpno, group_has_add_map.get("__all__", False))
                classification = _classify_change_type(chg_type, note, group_has_add=group_has_add)

                maker = bi.get("maker", "")
                qty_base = bi.get("qty", "") if base_pno else ""

                # 금형 추론
                mold = "○" if part_type == "사출" else "X"
                # 사내제작 추론
                uit = bi.get("uit", "")
                # bi에 UIT 없으면 db_index에서 참고품번으로 조회
                if not uit and base_pno and db_index:
                    ref_key = re.sub(r'[^A-Z0-9]','', base_pno.upper())
                    ref_entry = db_index.get(ref_key, {})
                    uit = str(ref_entry.get("uit", "")).strip().upper()
                if uit in ("F", "P", "T"):
                    in_house = "○"
                elif uit in ("Q", "D", "G", "M", "R", "S", "X"):
                    in_house = "X"
                else:
                    in_house = ""   # UIT 없으면 빈칸 (신규 부품 등)

                # New P/No 결정
                new_pno = pno_raw
                # Base에 있는 품번이 그대로 쓰이면서 상위 Assy면 → 채번 필요
                if base_pno == pno_raw and "▲ 상위 Assy" in note:
                    new_pno = "(채번 필요)"

                rows.append({
                    "No.": idx,
                    "BOM\nLevel": lvl,
                    "Part Type": part_type,
                    "Base P/No": base_pno,
                    "New P/No": new_pno,
                    "부품명\nClass Desc.(Part Name)": desc,
                    "Q'ty\nBase": qty_base,
                    "Q'ty\nNew": qty_new,
                    "변경점\nChanging Point": chg_reason,
                    "변경사유\nChanging Reason": chg_reason,
                    "양산처\nSupplier": maker,
                    "신규/변경\n분류\nClassification": classification,
                    "금형\n개발/수정": mold,
                    "사내제작\nIn-house": in_house,
                    "부품인정시험\nPart Approval Test": "○",
                })

            return pd.DataFrame(rows)

        def _master_to_excel_bytes(master_df, target_model, base_model, dev_grade, event=""):
            """개발부품 마스터 → 엑셀 (2단 헤더 병합 + A열 빈칸)"""
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
            from openpyxl.utils import get_column_letter

            wb = Workbook()
            ws = wb.active
            ws.title = "Master"

            thin = Side(style='thin')
            border = Border(top=thin, left=thin, right=thin, bottom=thin)
            header_font = Font(bold=True, size=9)
            data_font = Font(size=9)
            center = Alignment(horizontal='center', vertical='center', wrap_text=True)
            left_align = Alignment(horizontal='left', vertical='center', wrap_text=True)

            # 색상
            fill_common  = PatternFill("solid", fgColor="DAEEF3")  # 연한 파랑
            fill_grade   = PatternFill("solid", fgColor="FDE9D9")  # 연한 주황
            fill_dev     = PatternFill("solid", fgColor="EBF1DE")  # 연한 초록
            fill_fmea    = PatternFill("solid", fgColor="E4DFEC")  # 연한 보라
            fill_hsms    = PatternFill("solid", fgColor="F2DCDB")  # 연한 빨강
            fill_info    = PatternFill("solid", fgColor="FFFF00")  # 노랑

            # ── Row 1: 모델 정보 ──
            ws.merge_cells("B1:D1")
            ws["B1"] = "Base Model/Grade"
            ws["B1"].font = header_font
            ws.merge_cells("E1:H1")
            ws["E1"] = f"{base_model} / {dev_grade}"
            ws["E1"].font = data_font

            ws.merge_cells("J1:L1")
            ws["J1"] = "New Model/Grade"
            ws["J1"].font = header_font
            ws.merge_cells("M1:P1")
            ws["M1"] = f"{target_model} / {dev_grade}"
            ws["M1"].font = data_font

            # ── Row 3: 섹션 헤더 (1단) ──
            # 공통: B~P (15열)
            ws.merge_cells("B3:P3")
            ws["B3"] = "공통 Common"
            ws["B3"].font = header_font
            ws["B3"].fill = fill_common
            ws["B3"].alignment = center

            # 부품등급심의: Q~X (8열)
            ws.merge_cells("Q3:X3")
            ws["Q3"] = "부품등급심의 Parts Grade Review"
            ws["Q3"].font = header_font
            ws["Q3"].fill = fill_grade
            ws["Q3"].alignment = center

            # 부품개발완료: Y~AK (13열)
            ws.merge_cells("Y3:AK3")
            ws["Y3"] = "부품개발완료 Parts dev. completion"
            ws["Y3"].font = header_font
            ws["Y3"].fill = fill_dev
            ws["Y3"].alignment = center

            # FMEA: AL~AO (4열)
            ws.merge_cells("AL3:AO3")
            ws["AL3"] = "FMEA"
            ws["AL3"].font = header_font
            ws["AL3"].fill = fill_fmea
            ws["AL3"].alignment = center

            # HSMS: AP~BB (13열)
            ws.merge_cells("AP3:BB3")
            ws["AP3"] = "HSMS"
            ws["AP3"].font = header_font
            ws["AP3"].fill = fill_hsms
            ws["AP3"].alignment = center

            # ── Row 4: 세부 컬럼 헤더 (2단) ──
            col_headers = [
                # 공통 (B~P, 15개)
                ("B", "No.", fill_common),
                ("C", "BOM\nLevel", fill_common),
                ("D", "Part Type", fill_common),
                ("E", "P/No\n(Base)", fill_common),
                ("F", "P/No\n(New)", fill_common),
                ("G", "부품명\nClass Desc.", fill_common),
                ("H", "Q'ty\n(Base)", fill_common),
                ("I", "Q'ty\n(New)", fill_common),
                ("J", "변경점\nChanging Point", fill_common),
                ("K", "변경사유\nChanging Reason", fill_common),
                ("L", "양산처\nSupplier", fill_common),
                ("M", "신규/변경\nClassification", fill_common),
                ("N", "금형 개발/수정\nMold", fill_common),
                ("O", "사내제작\nIn-house", fill_common),
                ("P", "부품인정시험\nApproval Test", fill_common),
                # 부품등급심의 (Q~X, 8개)
                ("Q", "설계FMEA\nDesign FMEA", fill_grade),
                ("R", "미실시 사유", fill_grade),
                ("S", "부품 등급\nPart Grade", fill_grade),
                ("T", "수율관리\nYield", fill_grade),
                ("U", "부품추적성\nTracking", fill_grade),
                ("V", "s-APQP", fill_grade),
                ("W", "공정/공장심사\nAudit", fill_grade),
                ("X", "Remark", fill_grade),
                # 부품개발완료 (Y~AK, 13개)
                ("Y", "부품승인", fill_dev),
                ("Z", "부품인정", fill_dev),
                ("AA", "시방인정", fill_dev),
                ("AB", "선정배경", fill_dev),
                ("AC", "전수검사\n구축", fill_dev),
                ("AD", "도면리뷰\n부품", fill_dev),
                ("AE", "칸칸대차\n적용", fill_dev),
                ("AF", "SQA\n초품검사", fill_dev),
                ("AG", "SQA\n수입검사", fill_dev),
                ("AH", "DQMS\n부품인정", fill_dev),
                ("AI", "부품승인\n(단가포함)", fill_dev),
                ("AJ", "s-APQP", fill_dev),
                ("AK", "공정/공장\n심사", fill_dev),
                # FMEA (AL~AO, 4개)
                ("AL", "설계FMEA\n전개", fill_fmea),
                ("AM", "미전개\n사유", fill_fmea),
                ("AN", "공정FMEA\n전개", fill_fmea),
                ("AO", "미전개\n사유", fill_fmea),
                # HSMS (AP~BB, 13개)
                ("AP", "HSMS\n등록유무", fill_hsms),
                ("AQ", "승인날짜\nApproved", fill_hsms),
                ("AR", "미등록\n사유", fill_hsms),
                ("AS", "Food\nContact", fill_hsms),
                ("AT", "FCM Part\nHSMS설정", fill_hsms),
                ("AU", "신규재질\n/색상", fill_hsms),
                ("AV", "신규재질\n/업체", fill_hsms),
                ("AW", "항균/살균\n기능부품", fill_hsms),
                ("AX", "Biocidal\n물질명", fill_hsms),
                ("AY", "살생물질\n공급업체", fill_hsms),
                ("AZ", "정밀분석\n성적서", fill_hsms),
                ("BA", "Plastic\nmarking", fill_hsms),
                ("BB", "재질표기\n사진", fill_hsms),
            ]

            for col_letter, header_text, fill in col_headers:
                cell = ws[f"{col_letter}4"]
                cell.value = header_text
                cell.font = header_font
                cell.fill = fill
                cell.alignment = center
                cell.border = border

            # ── Row 5+: 데이터 ──
            common_cols = [
                "No.", "BOM\nLevel", "Part Type",
                "Base P/No", "New P/No",
                "부품명\nClass Desc.(Part Name)",
                "Q'ty\nBase", "Q'ty\nNew",
                "변경점\nChanging Point", "변경사유\nChanging Reason",
                "양산처\nSupplier", "신규/변경\n분류\nClassification",
                "금형\n개발/수정", "사내제작\nIn-house", "부품인정시험\nPart Approval Test",
            ]

            for r_idx, (_, row) in enumerate(master_df.iterrows()):
                excel_row = r_idx + 5   # 데이터는 5행부터
                for c_idx, col_name in enumerate(common_cols):
                    col_letter = get_column_letter(c_idx + 2)  # B=2
                    cell = ws[f"{col_letter}{excel_row}"]
                    val = row.get(col_name, "")
                    if str(val).lower() in ('nan', 'none'):
                        val = ""
                    cell.value = val
                    cell.font = data_font
                    cell.alignment = center if c_idx < 5 else left_align
                    cell.border = border

                # 빈 셀에도 border 적용 (등급심의~HSMS)
                for c_idx in range(len(common_cols), len(col_headers)):
                    col_letter = get_column_letter(c_idx + 2)
                    cell = ws[f"{col_letter}{excel_row}"]
                    cell.value = ""
                    cell.border = border

            # ── 열 너비 ──
            widths = {"B":5, "C":8, "D":8, "E":16, "F":16, "G":28,
                      "H":6, "I":6, "J":22, "K":22, "L":16, "M":10,
                      "N":8, "O":8, "P":8}
            for col, w in widths.items():
                ws.column_dimensions[col].width = w
            # 나머지 열
            for i in range(17, 55):
                ws.column_dimensions[get_column_letter(i)].width = 10

            bio = BytesIO()
            wb.save(bio)
            return bio.getvalue()

        st.divider()
        st.subheader("⬇️ 변경부품리스트 Excel 추출")

        unresolved = all_rows[all_rows["변경유형"] == "⚠️확인필요"]
        change_df = None
        if len(unresolved) > 0:
            st.error(f"⚠️ 확인필요 {len(unresolved)}건이 남아있어 Excel을 추출할 수 없습니다:")
            for g_idx in sorted(all_edited.keys()):
                g = groups[g_idx]
                edf = all_edited[g_idx]
                if "변경유형" not in edf.columns:
                    continue
                g_unresolved = edf[edf["변경유형"] == "⚠️확인필요"]
                if len(g_unresolved) == 0:
                    continue
                header_desc = g["header"].get("desc", "(미분류)")
                st.markdown(f"**🔴 [{header_desc}]**")
                for _, row in g_unresolved.iterrows():
                    pno_v = str(row.get("품번", "")).strip()
                    pno_display2 = pno_v if pno_v and pno_v.lower() != "nan" else "(품번 미확정)"
                    name_v = str(row.get("부품명", "")).strip()
                    st.caption(f"  • {name_v} ({pno_display2})")
            st.info("💡 위 부품의 변경유형을 확정하거나, 챗봇에서 번호로 처리할 수 있습니다.")
        else:
            change_df = _make_change_parts_list(all_rows, st.session_state)

            import enrich
            import importlib
            importlib.reload(enrich)
            from enrich import enrich_change_parts, _build_base_index, expand_with_assy_tree

            base_raw = st.session_state.get("base_df_raw")

            db_index = _build_ref_bom_index()
            if isinstance(base_raw, pd.DataFrame) and len(base_raw) > 0:
                base_as_index = _build_base_index(base_raw)
                for k, v in base_as_index.items():
                    if k not in db_index:
                        db_index[k] = v

            # 1단계: 빈칸 채우기
            change_df = enrich_change_parts(
                change_df,
                base_raw if isinstance(base_raw, pd.DataFrame) else pd.DataFrame(),
                db_index=db_index,
            )

            # ✅ 디버그: change_df 상태 즉시 확인
            st.write("🔎 change_df.shape:", change_df.shape)
            st.write("🔎 change_df.columns:", list(change_df.columns))
            st.write("🔎 '변경유형' 있음:", "변경유형" in change_df.columns)
            if "변경유형" in change_df.columns:
                st.write("🔎 변경유형 값:", change_df["변경유형"].tolist())
            st.write("🔎 '비고' 샘플:", change_df["비고"].tolist()[:5] if "비고" in change_df.columns else "없음")

            # ✅ 2단계: 그룹별 상위 Assy 태깅 → 트리 확장
            _tags = []
            for g_idx in sorted(all_edited.keys()):
                g = groups[g_idx]
                h = g["header"]
                edf = all_edited[g_idx]
                n = len(edf[edf["변경유형"].isin(["추가", "변경(시방)", "삭제"])])
                _tags.extend([h.get("part_no", "")] * n)

            if len(_tags) == len(change_df):
                change_df["_group_parent_pno"] = _tags
            else:
                st.warning(f"⚠️ 태그 수({len(_tags)}) ≠ change_df 행 수({len(change_df)}), Assy 삽입 스킵")

            if isinstance(base_raw, pd.DataFrame) and len(base_raw) > 0:
                try:
                    change_df = expand_with_assy_tree(change_df, base_raw)
                except Exception as e:
                    st.error(f"expand_with_assy_tree 오류: {e}")
                    import traceback
                    st.code(traceback.format_exc())

                n_assy = len(change_df[change_df["비고"] == "▲ 상위 Assy"]) if "비고" in change_df.columns else 0
                n_change = len(change_df) - n_assy
                st.info(f"📋 변경부품: {n_change}건 | 상위 Assy: {n_assy}건")

            # ✅ 웹 미리보기 테이블
            st.markdown("### 👀 변경부품리스트 미리보기")

            # ✅ change_df 실제 컬럼 기준으로 표시 순서만 정의
            PREVIEW_PRIORITY = [
                "변경유형", "변경사유", "비고",
                "Part No", "Lvl", "Part Name(자)", "Description",
                "Technical Spec", "Qty", "UOM", "Type", "Supply Type",
                "I.S", "MR", "Document", "Check Out", "Change", "R",
                "Parent Part No(모)", "Maker", "Standard",
                "출처", "filled_from",
            ]
            PREVIEW_EXCLUDE = {"신규모델", "기준모델(Base)", "Unnamed: 0", "Unnamed: 2"}

            # 실제 컬럼 중 우선순위 목록에 있는 것 먼저
            preview_cols = [c for c in PREVIEW_PRIORITY if c in change_df.columns]
            # 나머지 컬럼 뒤에 추가
            preview_cols += [c for c in change_df.columns
                             if c not in preview_cols and c not in PREVIEW_EXCLUDE]

            st.caption(
                f"총 **{len(change_df)}건** 중 상위 **{min(20, len(change_df))}건** 미리보기 | "
                f"실제 컬럼 수: {len(preview_cols)} | "
                "🟢 추가 🟠 변경 🔴 삭제 ⬜ 참고"
            )

            # ✅ 컬럼이 비었으면 경고
            if not preview_cols:
                st.error(f"❌ 미리보기 컬럼 없음. change_df 컬럼: {list(change_df.columns)}")
            else:
                preview_df = change_df[preview_cols].head(20).copy().reset_index(drop=True)

                # ✅ pandas Series는 .get() 없음 → row["컬럼명"] 방식으로 수정
                def _row_style(row):
                    vt = row.get("변경유형", "")
                    if vt == "추가":
                        return "background-color: #e8f5e9"
                    elif vt == "변경(시방)":
                        return "background-color: #fff3e0"
                    elif vt == "삭제":
                        return "background-color: #fce4ec"
                    elif vt == "[참고]":
                        return "background-color: #f5f5f5; color: #888"
                    return ""

                try:
                    # ⚡️ 성능 개선: apply를 사용하여 각 행에 함수를 적용하는 대신,
                    # 스타일을 한 번에 지정하는 Styler.applymap 또는 Styler.map을 사용하는 것이 더 효율적입니다.
                    # 여기서는 행 전체에 스타일을 적용하므로 Styler.apply를 유지하되, 불필요한 로직을 정리합니다.
                    styled = preview_df.style.apply(lambda row: [_row_style(row)]*len(row), axis=1)
                    st.dataframe(
                        styled,
                        use_container_width=True,
                        hide_index=True,
                        height=min(38 * (len(preview_df) + 1), 650),
                    )
                except Exception as e:
                    # ✅ 스타일 실패 시 일반 dataframe으로 폴백
                    st.warning(f"스타일 적용 실패({e}), 기본 테이블로 표시")
                    st.dataframe(
                        preview_df,
                        use_container_width=True,
                        hide_index=True,
                        height=min(38 * (len(preview_df) + 1), 650),
                    )

                if len(change_df) > 20:
                    st.caption(f"... 이하 {len(change_df) - 20}건은 Excel 다운로드에서 확인")

            # ✅ 디버그 expander
            with st.expander("🔍 Enrich 디버그", expanded=False):
                import enrich as _enrich_mod
                from enrich import _pick_col as _ep_col, _norm_pno as _ep_norm

                st.write("**change_df 컬럼:**", list(change_df.columns))
                base_raw2 = st.session_state.get("base_df_raw")
                st.write("**base_df_raw 컬럼:**",
                         list(base_raw2.columns) if isinstance(base_raw2, pd.DataFrame) else "없음")

                pno_col_dbg = _enrich_mod._pick_change_col(change_df, "part_no") if hasattr(_enrich_mod, "_pick_change_col") else None
                st.write("**품번 컬럼 매칭:**", pno_col_dbg)
                st.write(f"**참고 BOM 인덱스:** {len(db_index)}건")
                if db_index:
                    sample_keys = list(db_index.keys())[:3]
                    for k in sample_keys:
                        st.write(f"  `{k}` →", db_index[k])

                st.write("**filled_from 분포:**",
                         change_df["filled_from"].value_counts().to_dict()
                         if "filled_from" in change_df.columns else "없음")

                n_change2 = len(change_df[change_df["변경유형"] != "[참고]"])
                n_ref2    = len(change_df[change_df["변경유형"] == "[참고]"])
                st.write(f"변경: {n_change2}건 | 참고: {n_ref2}건")

            # ✅ 다운로드 버튼
            ts = _dt.datetime.now().strftime("%Y%m%d_%H%M")
            model_tag = st.session_state.get("target_model", "MODEL")
            # ✅ 엑셀 출력 시 Base BOM 양식 외 컬럼 제거
            EXPORT_EXCLUDE = {"신규모델", "기준모델(Base)", "변경유형", "변경사유",
                              "비고", "출처", "filled_from", "_group_parent_pno"}
            export_df = change_df.drop(
                columns=[c for c in change_df.columns
                         if c in EXPORT_EXCLUDE or str(c).startswith("Unnamed")],
                errors="ignore",
            )
            excel_bytes = _to_excel_bytes(export_df)

            st.download_button(
                label="⬇️ 변경부품리스트 Excel 다운로드",
                data=excel_bytes,
                file_name=f"ChangePartsList_{model_tag}_{ts}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                type="primary",
                use_container_width=True,
            )
        
        # ── 개발부품 마스터 ──
        st.divider()
        st.subheader("📋 개발부품 마스터 Excel 추출")
        if change_df is None:
            st.info("'확인필요' 항목을 확정하면 개발부품 마스터를 생성할 수 있습니다.")
        else:
            master_df = _make_dev_parts_master(change_df, st.session_state, db_index=db_index)

            if len(master_df) > 0:
                st.caption(f"총 **{len(master_df)}건** | 공통 섹션 자동 채움, 나머지 빈칸")
                st.dataframe(master_df.head(15), use_container_width=True, hide_index=True)

                master_bytes = _master_to_excel_bytes(
                    master_df,
                    target_model=st.session_state.get("target_model", ""),
                    base_model=st.session_state.get("bom_model", ""),
                    dev_grade=st.session_state.get("dev_grade", ""),
                )

                ts2 = _dt.datetime.now().strftime("%Y%m%d_%H%M")
                st.download_button(
                    label="📋 개발부품 마스터 Excel 다운로드",
                    data=master_bytes,
                    file_name=f"DevPartsMaster_{model_tag}_{ts2}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                )
            else:
                st.warning("마스터 생성할 데이터가 없습니다.")

# ── 🆕 플로팅 피드백 챗봇 ──

render_floating_chat()

    # v2.5: 제품군 추출 강화 + 모듈 매칭 정확도 개선 (2026-05-12)